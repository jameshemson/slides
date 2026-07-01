"""Shared python-pptx helpers for the build-deck scripts.

Decision D-002 is absolute for this module: it never sets fonts, colours, or
coordinates. The user's template carries all visual design. Code here only
opens templates, lists and resolves layouts, fills the placeholders the
template already defines, and — for `apply_theme` — writes brand values that
the *caller* supplies (read from brand.json) into the template's theme XML.
There are no font, colour, or coordinate literals anywhere in this file.

Chart drawing is a deliberately separate concern: it lives in `charts.py`
(matplotlib) and is placed by render.py, which derives the picture's geometry
from the template. Composed-slide drawing is likewise separate: token-bound
primitives live in `primitives.py` (the only module that emits literals, all of
them from brand.json design tokens) and are gated by `lint.py`. pptxlib itself
still adds no shape and holds no literals.

Public surface — six helpers:

    load_template(path)              -> pptx.Presentation
    list_layouts(prs)                -> list[dict]   (index, name, placeholders)
    resolve_role(prs, layout_map, role) -> SlideLayout
    fill_placeholders(slide, fields) -> None
    apply_theme(prs, brand)          -> None
    read_theme(prs)                  -> dict   (fonts, colours — inverse of apply_theme)

`read_theme` only reads the template's own theme values back out; it introduces
no literal of its own, so the no-literals invariant holds.
"""
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

# Placeholder types that hold slide *content* (a role's fields fill these).
CONTENT_PLACEHOLDER_TYPES = frozenset(
    {
        PP_PLACEHOLDER.TITLE,
        PP_PLACEHOLDER.CENTER_TITLE,
        PP_PLACEHOLDER.SUBTITLE,
        PP_PLACEHOLDER.BODY,
        PP_PLACEHOLDER.OBJECT,
    }
)

# Furniture placeholders — date, footer, slide number. Never filled, never
# counted, never removed. Inherited from the master and left to the template.
FURNITURE_PLACEHOLDER_TYPES = frozenset(
    {
        PP_PLACEHOLDER.DATE,
        PP_PLACEHOLDER.FOOTER,
        PP_PLACEHOLDER.SLIDE_NUMBER,
    }
)

# DrawingML namespace, used only to locate elements when editing theme XML.
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
# PresentationML namespace — used to read the master's txStyles (font sizes).
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def read_type_scale(prs):
    """Read the master's title/body default font sizes (in points).

    The theme's font scheme carries typefaces, not sizes — the real slide type
    sizes live in the slide master's `p:txStyles` (titleStyle/bodyStyle level-1
    `defRPr sz`, in hundredths of a point). Returns {"title": float?, "body":
    float?} for whichever are present; an empty dict when unreadable, so the
    caller (tokens.type_scale_from_master) falls back to the generic scale.
    """
    try:
        masters = list(prs.slide_masters)
    except AttributeError:
        return {}
    for master in masters:
        tx = master.element.find(f"{{{_P_NS}}}txStyles")
        if tx is None:
            continue
        out = {}
        for tag, key in (("titleStyle", "title"), ("bodyStyle", "body")):
            style = tx.find(f"{{{_P_NS}}}{tag}")
            if style is None:
                continue
            lvl1 = style.find(f"{{{_A_NS}}}lvl1pPr")
            def_rpr = lvl1.find(f"{{{_A_NS}}}defRPr") if lvl1 is not None else None
            sz = def_rpr.get("sz") if def_rpr is not None else None
            if sz:
                try:
                    out[key] = int(sz) / 100.0
                except ValueError:
                    pass
        return out
    return {}


def load_template(path):
    """Open a .pptx/.potx template and return a python-pptx Presentation.

    A template and a deck are the same thing to python-pptx: both carry a
    master and named layouts. Raises FileNotFoundError / PackageNotFoundError
    style errors up to the caller, which names the offending input.
    """
    return Presentation(path)


def list_layouts(prs):
    """Describe every slide layout in `prs`.

    Returns a list of dicts, one per layout, in index order:

        {"index": int, "name": str,
         "placeholders": [{"idx": int, "type": str}, ...]}

    Placeholder type is the PP_PLACEHOLDER member name (e.g. "BODY"). All
    placeholders are listed, including furniture, so inspect_template.py shows
    the template exactly as it is.
    """
    layouts = []
    for index, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            fmt = ph.placeholder_format
            placeholders.append(
                {
                    "idx": fmt.idx,
                    "type": _type_name(fmt.type),
                }
            )
        placeholders.sort(key=lambda p: p["idx"])
        layouts.append(
            {
                "index": index,
                "name": layout.name,
                "placeholders": placeholders,
            }
        )
    return layouts


def resolve_role(prs, layout_map, role):
    """Resolve a semantic role to one of the template's slide layouts.

    `layout_map` maps a role name to a layout index (brand.json's layout_map).
    Raises KeyError if the role is not mapped, IndexError if the mapped index
    is out of range — the caller turns these into a named, non-zero exit.
    """
    layout_idx = layout_map[role]
    return prs.slide_layouts[layout_idx]


def fill_placeholders(slide, fields):
    """Fill a slide's placeholders from an ordered list of role fields.

    `fields` is a list of (name, value) pairs in the role's field order. The
    first field fills the slide's title placeholder; the rest fill the
    remaining content placeholders in idx order. Furniture placeholders (date,
    footer, slide number) are never touched or counted.

    A value may be a string (one paragraph) or a list of strings (one
    paragraph per item — a bullet list or short paragraphs).

    Returns the list of content placeholders that were NOT filled, so the
    caller may optionally drop them to avoid empty "click to add" prompts.
    Raises ValueError if there are more fields than content placeholders; the
    caller names the slide and role in its exit message.
    """
    # Content placeholders in idx order: the lowest-idx one is the slide's
    # title placeholder (idx-0 TITLE/CENTER_TITLE), the rest follow it. The
    # role's first field fills the title; the remaining fields fill the rest.
    # python-pptx returns a fresh wrapper per access, so placeholders are
    # keyed and compared by idx, never by object identity.
    content_phs = _content_placeholders(slide)

    if not fields:
        return content_phs

    if len(fields) > len(content_phs):
        raise ValueError(
            f"role needs {len(fields)} placeholders but the layout has "
            f"{len(content_phs)}"
        )

    used_idx = set()
    for (_, value), placeholder in zip(fields, content_phs):
        _write_text(placeholder, value)
        used_idx.add(placeholder.placeholder_format.idx)

    return [
        ph for ph in content_phs if ph.placeholder_format.idx not in used_idx
    ]


def apply_theme(prs, brand):
    """Write brand fonts and colours into the template's theme XML.

    Pragmatic, not a design showcase: the theme's major/minor fonts are set
    from `brand["fonts"]` and the theme colour scheme's accent slots from
    `brand["colours"]`. Every value comes from the caller-supplied `brand`
    dict (read from brand.json) — this function holds no font or colour
    literal of its own. Used by make_template.py to theme a starter deck.

    No-ops gracefully if the package has no theme part.
    """
    for theme_part in _theme_parts(prs):
        root = etree.fromstring(theme_part.blob)
        theme_elements = root.find(f"{{{_A_NS}}}themeElements")
        if theme_elements is None:
            continue

        fonts = brand.get("fonts") or {}
        heading = fonts.get("heading")
        body = fonts.get("body")
        if heading or body:
            _apply_font_scheme(theme_elements, heading, body)

        colours = brand.get("colours") or {}
        if colours:
            _apply_colour_scheme(theme_elements, colours)

        theme_part._blob = etree.tostring(
            root, xml_declaration=True, encoding="UTF-8", standalone=True
        )


# The clrScheme slots read into a brand palette, and the brand.json name each
# maps to. accent1 becomes the primary `accent` (the key the chart colour
# resolver and teach-slides look for first); dk1/lt1 become ink/paper. The
# other scheme slots (dk2/lt2/hlink/folHlink) are not brand-palette colours and
# are skipped. Mirrors apply_theme, which writes accent1..6 and leaves dk/lt.
_BRAND_COLOUR_SLOTS = (
    ("accent1", "accent"),
    ("accent2", "accent2"),
    ("accent3", "accent3"),
    ("accent4", "accent4"),
    ("accent5", "accent5"),
    ("accent6", "accent6"),
    ("dk1", "ink"),
    ("lt1", "paper"),
)


def read_theme(prs):
    """Read brand fonts and colours back out of the template's theme XML.

    The inverse of `apply_theme`: it navigates the same theme parts and elements
    but reads instead of writes. Returns:

        {"fonts": {"heading": str?, "body": str?},
         "colours": {name: "#RRGGBB", ...}}

    Fonts come from the theme's major/minor font scheme (heading/body); colours
    from the colour scheme slots in `_BRAND_COLOUR_SLOTS`, reading either an
    `srgbClr` or a `sysClr` child. Used by extract_brand.py so teach-slides can
    pre-fill brand.json from a deck the user already has. Reads the first theme
    part that carries themeElements (a package has one theme); no-ops to empty
    maps if there is no theme.
    """
    for theme_part in _theme_parts(prs):
        root = etree.fromstring(theme_part.blob)
        theme_elements = root.find(f"{{{_A_NS}}}themeElements")
        if theme_elements is None:
            continue
        return {
            "fonts": _read_font_scheme(theme_elements),
            "colours": _read_colour_scheme(theme_elements),
        }
    return {"fonts": {}, "colours": {}}


# --- internal helpers --------------------------------------------------------


def _type_name(ph_type):
    """PP_PLACEHOLDER member name as a string, or None if unset."""
    if ph_type is None:
        return None
    name = getattr(ph_type, "name", None)
    return name if name is not None else str(ph_type)


def _content_placeholders(slide):
    """Slide placeholders that hold content, in idx order. No furniture."""
    phs = [
        ph
        for ph in slide.placeholders
        if ph.placeholder_format.type in CONTENT_PLACEHOLDER_TYPES
    ]
    phs.sort(key=lambda ph: ph.placeholder_format.idx)
    return phs


def _write_text(placeholder, value):
    """Write a string or list-of-strings into a placeholder's text frame.

    One paragraph per list item; a bare string is a single paragraph. Only
    text is set — no run-level font, size, or colour, so the template's
    inherited formatting carries through unchanged.
    """
    tf = placeholder.text_frame
    if isinstance(value, (list, tuple)):
        items = [str(v) for v in value]
    else:
        items = [str(value)]
    if not items:
        items = [""]

    tf.paragraphs[0].text = items[0]
    # Drop any extra inherited paragraphs beyond the first.
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)
    for line in items[1:]:
        tf.add_paragraph().text = line


def _theme_parts(prs):
    """Yield the package's theme parts (reached via each slide master).

    A theme is a generic OPC part whose XML lives in its `.blob`; python-pptx
    does not model it as an element graph, so callers parse and rewrite the
    blob directly.
    """
    seen = set()
    try:
        masters = list(prs.slide_masters)
    except AttributeError:
        return
    for master in masters:
        for rel in master.part.rels.values():
            if "theme" in rel.reltype:
                part = rel.target_part
                key = str(part.partname)
                if key not in seen:
                    seen.add(key)
                    yield part


def _apply_font_scheme(theme_elements, heading, body):
    """Set the major (heading) and minor (body) Latin typefaces in the theme."""
    font_scheme = theme_elements.find(f"{{{_A_NS}}}fontScheme")
    if font_scheme is None:
        return
    mapping = (("majorFont", heading), ("minorFont", body))
    for tag, typeface in mapping:
        if not typeface:
            continue
        font = font_scheme.find(f"{{{_A_NS}}}{tag}")
        if font is None:
            continue
        latin = font.find(f"{{{_A_NS}}}latin")
        if latin is not None:
            latin.set("typeface", typeface)


def _apply_colour_scheme(theme_elements, colours):
    """Set theme colour-scheme accent slots from brand colours.

    Brand colours arrive as a name->hex dict; values are mapped onto the
    theme's accent1..accent6 slots in dict order. Hex strings are normalised
    to the 6-digit, no-hash form OOXML expects. dk1/lt1/dk2/lt2 are left as
    the template defines them.
    """
    clr_scheme = theme_elements.find(f"{{{_A_NS}}}clrScheme")
    if clr_scheme is None:
        return

    accent_slots = [f"accent{n}" for n in range(1, 7)]
    hex_values = [_normalise_hex(v) for v in colours.values() if _normalise_hex(v)]

    for slot, hex_value in zip(accent_slots, hex_values):
        slot_el = clr_scheme.find(f"{{{_A_NS}}}{slot}")
        if slot_el is None:
            continue
        # Replace whatever colour child the slot holds with an srgbClr.
        for child in list(slot_el):
            slot_el.remove(child)
        srgb = slot_el.makeelement(f"{{{_A_NS}}}srgbClr", {"val": hex_value})
        slot_el.append(srgb)


def _read_font_scheme(theme_elements):
    """Read {heading, body} from the theme's major/minor Latin typefaces."""
    fonts = {}
    font_scheme = theme_elements.find(f"{{{_A_NS}}}fontScheme")
    if font_scheme is None:
        return fonts
    for tag, key in (("majorFont", "heading"), ("minorFont", "body")):
        font = font_scheme.find(f"{{{_A_NS}}}{tag}")
        if font is None:
            continue
        latin = font.find(f"{{{_A_NS}}}latin")
        if latin is not None and latin.get("typeface"):
            fonts[key] = latin.get("typeface")
    return fonts


def _read_colour_scheme(theme_elements):
    """Read named brand colours from the theme's colour-scheme slots.

    Each slot in `_BRAND_COLOUR_SLOTS` is read to a `#RRGGBB` value; slots that
    are absent or carry no parseable colour are skipped.
    """
    colours = {}
    clr_scheme = theme_elements.find(f"{{{_A_NS}}}clrScheme")
    if clr_scheme is None:
        return colours
    for slot, name in _BRAND_COLOUR_SLOTS:
        slot_el = clr_scheme.find(f"{{{_A_NS}}}{slot}")
        if slot_el is None:
            continue
        hex_value = _read_slot_hex(slot_el)
        if hex_value:
            colours[name] = "#" + hex_value
    return colours


def _read_slot_hex(slot_el):
    """Hex (6 upper digits, no '#') from a clrScheme slot, or None.

    Reads an `srgbClr val=` first, then falls back to a `sysClr lastClr=`
    (used by dk1/lt1, which carry a system colour with a resolved hex).
    """
    srgb = slot_el.find(f"{{{_A_NS}}}srgbClr")
    if srgb is not None:
        normalised = _normalise_hex(srgb.get("val"))
        if normalised:
            return normalised
    sys_clr = slot_el.find(f"{{{_A_NS}}}sysClr")
    if sys_clr is not None:
        return _normalise_hex(sys_clr.get("lastClr"))
    return None


def _normalise_hex(value):
    """Return a 6-hex-digit uppercase colour string, or None if unparseable."""
    if not isinstance(value, str):
        return None
    text = value.strip().lstrip("#")
    if len(text) == 3:
        text = "".join(ch * 2 for ch in text)
    if len(text) != 6:
        return None
    try:
        int(text, 16)
    except ValueError:
        return None
    return text.upper()
