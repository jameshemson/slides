"""primitives.py — D-002 carve-out: the ONLY module that emits colour and
coordinate literals to pptx slide objects.

Every colour value and every EMU coordinate written here is taken directly
from the caller-supplied tokens dict.  There are no brand literals in this
module.  The two numeric constants below (EMU_PER_PT, LINE_HEIGHT) are generic
typography conversion factors, not brand values.

Relationship to charts.py: this module plays the same isolation role for
shape/text drawing that charts.py plays for chart images — render.py delegates
all literal-emitting work here so the rest of the renderer stays literal-free.
"""

from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

# Generic typography/layout constants — the ONLY non-token numeric constants
# allowed (generic, not brand values).
EMU_PER_PT = 12700
LINE_HEIGHT = 1.2
# Optical centre: place the row's vertical centre slightly above the geometric
# centre (~45% from the top). Dead-centre reads as marginally low; the optical
# centre is the conventional resting point for a single focal element.
OPTICAL_CENTRE = 0.45

# Draw order (z): filled panels sit behind connectors, which sit behind text.
# A card is a container box with its text layered on top; the lint permits the
# overlap only because the box is a container (see lint.check_no_overlap), and
# draw() must paint the box first or the text would be hidden. Elements with an
# equal z keep their planned order (stable sort).
_Z_BY_KIND = {"box": 0, "edge": 1, "connector": 1, "icon": 2, "text": 2}


class ShapeError(Exception):
    """A primitive that cannot be drawn. render.py turns this into a SpecError."""


def _normalise_hex(value) -> "str | None":
    """Accept '#abc', 'abc', '#aabbcc', 'aabbcc'; return '#RRGGBB' uppercase or None."""
    if not isinstance(value, str):
        return None
    s = value.lstrip("#").strip()
    if len(s) == 3 and all(c in "0123456789abcdefABCDEF" for c in s):
        s = s[0] * 2 + s[1] * 2 + s[2] * 2
    if len(s) == 6 and all(c in "0123456789abcdefABCDEF" for c in s):
        return "#" + s.upper()
    return None


def plan_stat_row(stats, tokens, slide_w, slide_h, region=None) -> list:
    """PURE function — no pptx objects.  Returns a list of element dicts.

    Each element dict has keys:
      role, text, left, top, width, height, font_pt, colour
    where left/top/width/height are in EMU and colour is '#RRGGBB'.
    """
    grid = tokens["grid"]
    ts = tokens["type_scale"]
    roles = tokens["colour_roles"]

    # Validate required keys.
    missing_roles = [k for k in ("accent", "ink") if k not in roles]
    if missing_roles:
        raise ShapeError(f"colour_roles missing: {', '.join(missing_roles)}")
    missing_ts = [k for k in ("display", "caption") if k not in ts]
    if missing_ts:
        raise ShapeError(f"type_scale missing: {', '.join(missing_ts)}")

    N = len(stats)
    if N == 0:
        raise ShapeError("stat_row needs at least one stat")

    gutter = grid["gutter"]
    baseline = grid["baseline"]

    content_left, content_w = _content_span(tokens, slide_w, region)
    cell_w = (content_w - (N - 1) * gutter) // N

    number_pt = ts["display"]
    label_pt = ts["caption"]

    number_h = round(number_pt * EMU_PER_PT * LINE_HEIGHT)
    label_h = round(label_pt * EMU_PER_PT * LINE_HEIGHT)
    row_block_h = number_h + baseline + label_h

    if region is None:
        band_top = grid["margin_top"]
        band_bottom = slide_h - grid["margin_bottom"]
    else:
        _l, t, _w, h = region
        band_top = t
        band_bottom = t + h

    row_top = band_top + round((band_bottom - band_top - row_block_h) * OPTICAL_CENTRE)

    elements = []
    for i in range(N):
        cell_left = content_left + i * (cell_w + gutter)
        if i < N - 1:
            width_i = cell_w
        else:
            # Last cell absorbs rounding remainder.
            width_i = content_w - (N - 1) * (cell_w + gutter)

        elements.append({
            "role": "stat-number",
            "text": str(stats[i]["value"]),
            "left": cell_left,
            "top": row_top,
            "width": width_i,
            "height": number_h,
            "font_pt": number_pt,
            "colour": roles["accent"],
        })
        elements.append({
            "role": "stat-label",
            "text": str(stats[i]["label"]),
            "left": cell_left,
            "top": row_top + number_h + baseline,
            "width": width_i,
            "height": label_h,
            "font_pt": label_pt,
            "colour": roles["ink"],
        })

    return elements


# --- shared geometry helpers for the box-based primitives --------------------
#
# Every primitive below is a PURE planner: it returns a list of element dicts
# (box/text/connector) and touches no pptx object. Geometry comes from the grid
# tokens; every colour is a colour_roles value. The look is grounded in the
# design research (James's own decks + the Visme guide + Gestalt Common Region):
# equal panels grouped by a shared fill/outline, 3-5 siblings, grey field with
# one accent, hierarchy by size. draw() renders whatever these return.

_STROKE_EMU = EMU_PER_PT  # 1pt hairline panel outline (generic, not a brand value)


def _line_h(font_pt) -> int:
    """Line box height in EMU for a point size (matches plan_stat_row)."""
    return round(font_pt * EMU_PER_PT * LINE_HEIGHT)


def _require(tokens, ts_keys):
    """Validate the token sub-dicts a primitive needs; return (grid, ts, roles).

    Raises ShapeError naming the missing role or type-scale step — the same
    failure mode plan_stat_row uses, so a thin brand profile fails loudly.
    """
    roles = tokens.get("colour_roles", {}) or {}
    missing = [k for k in ("accent", "ink", "paper", "muted") if k not in roles]
    if missing:
        raise ShapeError(f"colour_roles missing: {', '.join(missing)}")
    ts = tokens.get("type_scale", {}) or {}
    missing_ts = [k for k in ts_keys if k not in ts]
    if missing_ts:
        raise ShapeError(f"type_scale missing: {', '.join(missing_ts)}")
    return tokens["grid"], ts, roles


def _content_span(tokens, slide_w, region=None):
    """(left, width) of the horizontal span a primitive fills.

    A placement `region` narrows the span (a block placed in the left columns
    draws in the left columns); with no region it spans margin to margin."""
    if region is not None:
        return region[0], region[2]
    margin_x = tokens["grid"]["margin_x"]
    return margin_x, slide_w - 2 * margin_x


def _band(tokens, slide_h, region):
    """(top, bottom) of the vertical band a primitive fills — the region under a
    title if one was reserved, else the full margin-to-margin band."""
    grid = tokens["grid"]
    if region is None:
        return grid["margin_top"], slide_h - grid["margin_bottom"]
    _l, t, _w, h = region
    return t, t + h


def _even_cells(content_left, content_w, n, gutter):
    """n equal cells across the content width, separated by `gutter`.

    The last cell absorbs the rounding remainder so the row's right edge lands
    exactly on the content margin (same rule as plan_stat_row)."""
    cell_w = (content_w - (n - 1) * gutter) // n
    cells = []
    for i in range(n):
        left = content_left + i * (cell_w + gutter)
        width = cell_w if i < n - 1 else content_w - (n - 1) * (cell_w + gutter)
        cells.append((left, width))
    return cells


def _split_body(text):
    """A card/panel body: ' / ' marks a line break, so a few terse points can
    share one text element without exploding the element count."""
    return "\n".join(part.strip() for part in str(text).split(" / ") if part.strip())


# --- card-grid ---------------------------------------------------------------


def plan_card_grid(cards, tokens, slide_w, slide_h, region=None) -> list:
    """A row of 3-5 equal panels — 'cluster by message: three or five topics'.

    cards: list of {"label": str, "body": str?, "emphasis": bool?}. One card may
    be marked emphasis to lead (accent fill, reversed text); the rest are paper
    panels with a hairline outline (Gestalt Common Region — the box binds its
    contents). Grouping is the point; hierarchy is by the one emphasised card.
    """
    grid, ts, roles = _require(tokens, ("h1", "body"))
    n = len(cards)
    if n == 0:
        raise ShapeError("card-grid needs at least one card")
    if n > 6:
        raise ShapeError("card-grid takes at most 6 cards; split the slide")

    content_left, content_w = _content_span(tokens, slide_w, region)
    gutter = grid["gutter"]
    pad = gutter
    baseline = grid["baseline"]
    cells = _even_cells(content_left, content_w, n, gutter)
    band_top, band_bottom = _band(tokens, slide_h, region)
    band_h = band_bottom - band_top

    label_pt, body_pt = ts["h1"], ts["body"]
    label_h = _line_h(label_pt)
    has_body = any((c.get("body") or "").strip() for c in cards)
    body_alloc = 2 * _line_h(body_pt) if has_body else 0

    inner_h = label_h + (baseline + body_alloc if has_body else 0)
    card_h = min(inner_h + 2 * pad, band_h)
    card_top = band_top + round((band_h - card_h) * OPTICAL_CENTRE)

    elements = []
    for i, c in enumerate(cards):
        cl, cw = cells[i]
        emph = bool(c.get("emphasis"))
        panel = {
            "role": "card-panel", "kind": "box", "container": True,
            "left": cl, "top": card_top, "width": cw, "height": card_h,
            "fill": roles["accent"] if emph else roles["paper"],
        }
        if not emph:
            panel["stroke"] = roles["ink"]
            panel["stroke_w"] = _STROKE_EMU
        elements.append(panel)
        text_colour = roles["paper"] if emph else roles["ink"]
        elements.append({
            "role": "card-label", "text": str(c.get("label", "")),
            "left": cl + pad, "top": card_top + pad,
            "width": cw - 2 * pad, "height": label_h,
            "font_pt": label_pt, "colour": text_colour, "bold": True,
        })
        body = (c.get("body") or "").strip()
        if has_body and body:
            elements.append({
                "role": "card-body", "text": _split_body(body),
                "left": cl + pad, "top": card_top + pad + label_h + baseline,
                "width": cw - 2 * pad, "height": body_alloc,
                "font_pt": body_pt, "colour": text_colour,
            })
    return elements


# --- comparison / two-panel --------------------------------------------------


def plan_comparison(sides, tokens, slide_w, slide_h, region=None) -> list:
    """Two panels set side by side so a difference is unmissable.

    sides: exactly two {"header": str, "body": str?, "emphasis": bool?}. A
    comparison must RESOLVE, not just balance — mark the winning side emphasis
    ('order for impact') and it fills with the accent; the design tilts to the
    turn instead of sitting symmetric.
    """
    grid, ts, roles = _require(tokens, ("h1", "body"))
    if len(sides) != 2:
        raise ShapeError("comparison needs exactly two panels")

    content_left, content_w = _content_span(tokens, slide_w, region)
    gutter = grid["gutter"]
    pad = gutter
    baseline = grid["baseline"]
    cells = _even_cells(content_left, content_w, 2, gutter)
    band_top, band_bottom = _band(tokens, slide_h, region)
    band_h = band_bottom - band_top

    header_pt, body_pt = ts["h1"], ts["body"]
    header_h = _line_h(header_pt)
    max_lines = 1
    for s in sides:
        lines = _split_body(s.get("body") or "").split("\n")
        max_lines = max(max_lines, len([x for x in lines if x]))
    max_lines = min(max_lines, 6)
    body_alloc = max_lines * _line_h(body_pt)
    inner_h = header_h + baseline + body_alloc
    panel_h = min(inner_h + 2 * pad, band_h)
    panel_top = band_top + round((band_h - panel_h) * OPTICAL_CENTRE)

    elements = []
    for i, s in enumerate(sides):
        cl, cw = cells[i]
        emph = bool(s.get("emphasis"))
        panel = {
            "role": "comparison-panel", "kind": "box", "container": True,
            "left": cl, "top": panel_top, "width": cw, "height": panel_h,
            "fill": roles["accent"] if emph else roles["paper"],
        }
        if not emph:
            panel["stroke"] = roles["ink"]
            panel["stroke_w"] = _STROKE_EMU
        elements.append(panel)
        text_colour = roles["paper"] if emph else roles["ink"]
        elements.append({
            "role": "comparison-header", "text": str(s.get("header", "")),
            "left": cl + pad, "top": panel_top + pad,
            "width": cw - 2 * pad, "height": header_h,
            "font_pt": header_pt, "colour": text_colour, "bold": True,
        })
        body = (s.get("body") or "").strip()
        if body:
            elements.append({
                "role": "comparison-body", "text": _split_body(body),
                "left": cl + pad, "top": panel_top + pad + header_h + baseline,
                "width": cw - 2 * pad, "height": body_alloc,
                "font_pt": body_pt, "colour": text_colour,
            })
    return elements


# --- process / flow ----------------------------------------------------------


def plan_process(steps, tokens, slide_w, slide_h, region=None) -> list:
    """3 (up to 5) numbered steps left to right, joined by arrows.

    steps: list of {"label": str, "detail": str?}, numbered by order. Each step
    is a paper box with a big accent number, a bold label, and an optional light
    detail line; a muted arrow bridges the gap to the next. This is James's real
    'Plan -> Create -> Deliver' pattern, deliberately NOT a SmartArt chevron
    ribbon.
    """
    grid, ts, roles = _require(tokens, ("h1", "body", "caption"))
    n = len(steps)
    if n == 0:
        raise ShapeError("process needs at least one step")
    if n > 5:
        raise ShapeError("process takes at most 5 steps; keep it to 3-5")

    content_left, content_w = _content_span(tokens, slide_w, region)
    gutter = grid["gutter"]
    pad = gutter
    baseline = grid["baseline"]
    cells = _even_cells(content_left, content_w, n, gutter)
    band_top, band_bottom = _band(tokens, slide_h, region)
    band_h = band_bottom - band_top

    num_pt, label_pt, detail_pt = ts["h1"], ts["body"], ts["caption"]
    num_h, label_h, detail_h = _line_h(num_pt), _line_h(label_pt), _line_h(detail_pt)
    has_detail = any((s.get("detail") or "").strip() for s in steps)
    inner_h = num_h + baseline + label_h + (baseline + detail_h if has_detail else 0)
    box_h = min(inner_h + 2 * pad, band_h)
    box_top = band_top + round((band_h - box_h) * OPTICAL_CENTRE)

    elements = []
    for i, s in enumerate(steps):
        cl, cw = cells[i]
        elements.append({
            "role": "process-step", "kind": "box", "container": True,
            "left": cl, "top": box_top, "width": cw, "height": box_h,
            "fill": roles["paper"], "stroke": roles["ink"], "stroke_w": _STROKE_EMU,
        })
        elements.append({
            "role": "process-number", "text": str(i + 1),
            "left": cl + pad, "top": box_top + pad,
            "width": cw - 2 * pad, "height": num_h,
            "font_pt": num_pt, "colour": roles["accent"], "bold": True,
        })
        elements.append({
            "role": "process-label", "text": str(s.get("label", "")),
            "left": cl + pad, "top": box_top + pad + num_h + baseline,
            "width": cw - 2 * pad, "height": label_h,
            "font_pt": label_pt, "colour": roles["ink"], "bold": True,
        })
        detail = (s.get("detail") or "").strip()
        if has_detail and detail:
            elements.append({
                "role": "process-detail", "text": _split_body(detail),
                "left": cl + pad,
                "top": box_top + pad + num_h + baseline + label_h + baseline,
                "width": cw - 2 * pad, "height": detail_h,
                "font_pt": detail_pt, "colour": roles["ink"],
            })
        if i < n - 1:
            gap_left = cl + cw
            gap_right = cells[i + 1][0]
            inset = (gap_right - gap_left) // 6
            arrow_top = box_top + (box_h - label_h) // 2
            elements.append({
                "role": "process-connector", "kind": "connector",
                "shape": "right_arrow",
                "left": gap_left + inset, "top": arrow_top,
                "width": (gap_right - gap_left) - 2 * inset, "height": label_h,
                "fill": roles["ink"],
            })
    return elements


# --- timeline / milestones ---------------------------------------------------


def plan_timeline(nodes, tokens, slide_w, slide_h, region=None) -> list:
    """Dated milestones as dots on a rail — Start ...o...o...o... End.

    nodes: list of {"date": str, "event": str, "emphasis": bool?}. One milestone
    is the turn: mark it emphasis and it gets a larger accent dot and a bold
    label while the rest stay muted — a timeline still obeys grey-push + one
    accent, so it reads as a sequence with a hero beat, not an even dotted rule.
    """
    grid, ts, roles = _require(tokens, ("body",))
    n = len(nodes)
    if n == 0:
        raise ShapeError("timeline needs at least one milestone")
    if n > 6:
        raise ShapeError("timeline takes at most 6 milestones")

    content_left, content_w = _content_span(tokens, slide_w, region)
    gutter = grid["gutter"]
    baseline = grid["baseline"]
    cells = _even_cells(content_left, content_w, n, gutter)
    band_top, band_bottom = _band(tokens, slide_h, region)
    band_h = band_bottom - band_top

    label_pt = ts["body"]
    line_h = _line_h(label_pt)
    label_h = 2 * line_h
    dot_d = line_h
    emph_dot_d = dot_d + baseline
    rail_h = max(EMU_PER_PT // 2, dot_d // 8)

    block_h = min(emph_dot_d + baseline + label_h, band_h)
    block_top = band_top + round((band_h - block_h) * OPTICAL_CENTRE)
    rail_y = block_top + emph_dot_d // 2
    label_top = rail_y + emph_dot_d // 2 + baseline

    elements = []
    dot_edges = []  # (left, right) of each drawn dot, for exact rail seams
    for i, nd in enumerate(nodes):
        cl, cw = cells[i]
        centre = cl + cw // 2
        emph = bool(nd.get("emphasis"))
        d = emph_dot_d if emph else dot_d
        dleft = centre - d // 2
        dot_edges.append((dleft, dleft + d))
        elements.append({
            "role": "timeline-dot", "kind": "box", "shape": "oval",
            "left": dleft, "top": rail_y - d // 2,
            "width": d, "height": d,
            "fill": roles["accent"] if emph else roles["ink"],
        })
        date = str(nd.get("date", "")).strip()
        event = str(nd.get("event", "")).strip()
        text = f"{date}\n{event}" if date and event else (date or event)
        elements.append({
            "role": "timeline-label", "text": text,
            "left": cl, "top": label_top, "width": cw, "height": label_h,
            "font_pt": label_pt, "align": "center",
            "colour": roles["accent"] if emph else roles["ink"], "bold": emph,
        })
    for i in range(n - 1):
        left = dot_edges[i][1]        # right edge of dot i
        right = dot_edges[i + 1][0]   # left edge of dot i+1
        if right > left:
            elements.append({
                "role": "timeline-rail", "kind": "connector", "shape": "rect",
                "left": left, "top": rail_y - rail_h // 2,
                "width": right - left, "height": rail_h,
                "fill": roles["ink"],
            })
    return elements


# --- freeform ----------------------------------------------------------------


def _subrect(base, placement, grid):
    """A sub-rectangle of `base` for a {cols, rows} placement on its 12x12 grid.

    Mirrors render._place_region but relative to whatever band the freeform
    block was given, so element placement composes under block placement. A
    small inset keeps adjacent elements from touching."""
    bl, bt, bw, bh = base
    cols_n = grid.get("columns", 12) or 12
    rows_n = 12
    cols = placement.get("cols")
    rows = placement.get("rows")
    left = bl + (cols[0] - 1) * bw // cols_n if cols else bl
    right = bl + cols[1] * bw // cols_n if cols else bl + bw
    top = bt + (rows[0] - 1) * bh // rows_n if rows else bt
    bottom = bt + rows[1] * bh // rows_n if rows else bt + bh
    pad_x = grid.get("gutter", 0) // 2
    pad_y = grid.get("baseline", 0)
    return (left + pad_x, top + pad_y,
            max(1, (right - pad_x) - (left + pad_x)),
            max(1, (bottom - pad_y) - (top + pad_y)))


def plan_freeform(elements, tokens, slide_w, slide_h, region=None) -> list:
    """The escape hatch: place token-bound elements the named primitives don't
    cover — a matrix, a quadrant, a node graph, an annotated layout.

    Freedom in the arrangement, the SAME guardrails as every other composed
    block: each element's colour is a role name (ink/paper/accent/muted) and
    each text size a scale name (display/h1/body/caption), both resolved to the
    brand's tokens here, and every element still passes the mechanical lint
    (on-token, on-grid, no overlap outside a container, under the cap). What the
    lint does NOT prove is that the arrangement is *well* composed — that is the
    author's judgement, nudged by the one freeform advisory (grey-push).

    Each element dict (parsed by render._parse_freeform_element) carries:
      kind: box|panel|text|arrow|dot|line, a placement {cols,rows}, and either
      a `fill`(+`stroke`) / `colour` role name, plus `scale`+`text` for text.
    """
    grid, ts, roles = _require(tokens, ())
    if not elements:
        raise ShapeError("freeform needs at least one element")
    if region is None:
        mx = grid["margin_x"]
        band = (mx, grid["margin_top"], slide_w - 2 * mx,
                slide_h - grid["margin_top"] - grid["margin_bottom"])
    else:
        band = region

    out = []
    for el in elements:
        l, t, w, h = _subrect(band, el["placement"], grid)
        kind = el["kind"]
        if kind in ("box", "panel"):
            box = {
                "role": "freeform-panel", "kind": "box", "container": True,
                "left": l, "top": t, "width": w, "height": h,
                "fill": roles[el["fill"]],
            }
            if el.get("stroke"):
                box["stroke"] = roles[el["stroke"]]
                box["stroke_w"] = _STROKE_EMU
            out.append(box)
        elif kind == "text":
            out.append({
                "role": "freeform-text", "text": el["text"],
                "left": l, "top": t, "width": w, "height": h,
                "font_pt": ts[el["scale"]], "colour": roles[el["colour"]],
                "anchor": "middle",
            })
        elif kind == "arrow":
            out.append({
                "role": "freeform-arrow", "kind": "connector",
                "shape": "right_arrow",
                "left": l, "top": t, "width": w, "height": h,
                "fill": roles[el["colour"]],
            })
        elif kind == "dot":
            out.append({
                "role": "freeform-dot", "kind": "box", "shape": "oval",
                "left": l, "top": t, "width": w, "height": h,
                "fill": roles[el["colour"]],
            })
        elif kind == "line":
            lh = min(h, EMU_PER_PT * 2)  # a hairline divider, centred in its cell
            out.append({
                "role": "freeform-line", "kind": "connector", "shape": "rect",
                "left": l, "top": t + (h - lh) // 2, "width": w, "height": lh,
                "fill": roles[el["colour"]],
            })
        elif kind == "icon":
            # A square icon centred in its placement cell. render.py resolves the
            # element to a recoloured PNG after the lint clears its token colour.
            side = min(w, h)
            out.append({
                "role": "freeform-icon", "kind": "icon", "name": el["name"],
                "left": l + (w - side) // 2, "top": t + (h - side) // 2,
                "width": side, "height": side,
                "colour": roles[el["colour"]],
            })
    return out


# Autoshape names a box element may request via its "shape" key. Generic
# geometry, not brand values. Default is a rounded rectangle (the card/panel).
_SHAPE_BY_NAME = {
    "rect": MSO_SHAPE.RECTANGLE,
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded": MSO_SHAPE.ROUNDED_RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "chevron": MSO_SHAPE.CHEVRON,
    "right_arrow": MSO_SHAPE.RIGHT_ARROW,
    "line": MSO_SHAPE.RECTANGLE,
}

_ALIGN_BY_NAME = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "centre": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}

_ANCHOR_BY_NAME = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}


def draw(slide, elements) -> list:
    """Add one shape per element to `slide`; return the added shapes.

    An element's `kind` selects the shape: "box"/"connector" draw a filled
    autoshape (card, panel, step, connector), anything else (the default, and
    every stat-row element) draws a text box. Elements are painted in z-order
    (boxes behind connectors behind text) so a card's text lands on top of its
    panel; within a z-band the planned order is kept (stable sort). Every
    colour written here comes from the element dict, which the caller filled
    from brand tokens — this module still emits no brand literal of its own.
    """
    ordered = sorted(
        elements, key=lambda el: _Z_BY_KIND.get(el.get("kind", "text"), 2)
    )
    added = []
    for el in ordered:
        kind = el.get("kind")
        if kind in ("box", "connector"):
            added.append(_add_box(slide, el))
        elif kind == "edge":
            added.append(_add_edge(slide, el))
        elif kind == "icon":
            shp = _add_icon(slide, el)
            if shp is not None:
                added.append(shp)
        else:
            added.append(_add_text(slide, el))
    return added


def _add_edge(slide, el):
    """Draw one tree edge as a single elbow connector, parent -> child.

    An elbow connector is ONE shape per edge (so a tree stays under the element
    cap), routed by PowerPoint between the two points. The line colour is a
    token; the lint exempts connectors/edges from the overlap rule because a
    1-D line crossing a box is not a composition fault."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.ELBOW,
        Emu(el["x1"]), Emu(el["y1"]), Emu(el["x2"]), Emu(el["y2"]),
    )
    conn.line.color.rgb = RGBColor.from_string(el["colour"].lstrip("#"))
    if el.get("stroke_w"):
        conn.line.width = Emu(int(el["stroke_w"]))
    return conn


def _add_icon(slide, el):
    """Place a pre-rendered icon PNG on the grid. render.py recolours + rasterises
    the icon to `el['png']` after the lint clears it; an icon whose rasteriser was
    absent has no `png` and is dropped by render.py before draw, so this only ever
    places a real file. Returns None defensively if `png` is missing."""
    png = el.get("png")
    if not png:
        return None
    return slide.shapes.add_picture(
        png, Emu(el["left"]), Emu(el["top"]),
        width=Emu(el["width"]), height=Emu(el["height"]),
    )


def _add_text(slide, el):
    """Draw one text element as a textbox. Colour and size come from the token
    element dict; an optional `font`, `align`, and `anchor` refine it."""
    tb = slide.shapes.add_textbox(
        Emu(el["left"]), Emu(el["top"]), Emu(el["width"]), Emu(el["height"]),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    anchor = _ANCHOR_BY_NAME.get(el.get("anchor"))
    if anchor is not None:
        tf.vertical_anchor = anchor
    # Trim the textbox's default internal padding so text seats tight inside a
    # card; harmless for a free-standing stat where the box hugs the text.
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    align = _ALIGN_BY_NAME.get(el.get("align"))
    # A `\n` in the text splits into paragraphs (a card body's few lines);
    # every paragraph inherits the element's size, colour, weight, and font.
    lines = str(el["text"]).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(el["font_pt"])
        run.font.color.rgb = RGBColor.from_string(el["colour"].lstrip("#"))
        if el.get("bold"):
            run.font.bold = True
        if el.get("font"):
            run.font.name = el["font"]
    return tb


def _add_box(slide, el):
    """Draw one filled-shape element (card, panel, step, connector).

    Fill (and optional stroke) come from the token element dict. Shadows are
    turned off — the on-brand look is a flat token fill, not a drop shadow
    (guards the gradient/shadow SaaS cliché the composition rules warn against).
    """
    shape_enum = _SHAPE_BY_NAME.get(el.get("shape"), MSO_SHAPE.ROUNDED_RECTANGLE)
    shp = slide.shapes.add_shape(
        shape_enum,
        Emu(el["left"]), Emu(el["top"]), Emu(el["width"]), Emu(el["height"]),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(el["fill"].lstrip("#"))
    stroke = el.get("stroke")
    if stroke:
        shp.line.color.rgb = RGBColor.from_string(stroke.lstrip("#"))
        if el.get("stroke_w"):
            shp.line.width = Emu(int(el["stroke_w"]))
    else:
        shp.line.fill.background()
    try:
        shp.shadow.inherit = False
    except Exception:  # noqa: BLE001 — some shapes lack a shadow element
        pass
    return shp
