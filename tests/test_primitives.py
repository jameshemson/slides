"""Tests for primitives.py — D-002 carve-out."""
import sys
import os
import unittest

# Add the scripts directory to sys.path so `import primitives` resolves.
_SCRIPTS_DIR = os.path.join(
    os.path.dirname(__file__),
    "..",
    "source",
    "skills",
    "build-deck",
    "scripts",
)
sys.path.insert(0, os.path.abspath(_SCRIPTS_DIR))

import primitives
import lint
from primitives import (
    plan_stat_row, plan_card_grid, plan_comparison, plan_process,
    plan_timeline, plan_freeform, plan_tree, plan_icon_list,
    plan_cycle, plan_matrix, ShapeError, _normalise_hex,
)


def _node(label, *children, emphasis=False, icon=None):
    return {"label": label, "emphasis": emphasis, "icon": icon,
            "children": list(children)}
from pptx import Presentation
from pptx.dml.color import RGBColor

TEMPLATE = os.path.join(os.path.dirname(__file__), "fixtures", "sample-template.pptx")

TOKENS = {
    "grid": {
        "margin_x": 457200,
        "margin_top": 1600200,
        "margin_bottom": 731837,
        "columns": 12,
        "gutter": 152400,
        "baseline": 91211,
    },
    "type_scale": {
        "display": 40.0,
        "h1": 28.0,
        "body": 18.0,
        "caption": 12.0,
    },
    "colour_roles": {
        "ink": "#000000",
        "paper": "#FFFFFF",
        "accent": "#4F81BD",
        "muted": "#C0504D",
    },
}

SLIDE_W, SLIDE_H = 9144000, 6858000
STATS = [
    {"value": "56", "label": "Days"},
    {"value": "4%", "label": "Win rate"},
    {"value": "120", "label": "Deals"},
]


class TestPlanStatRow(unittest.TestCase):

    def setUp(self):
        self.els = plan_stat_row(STATS, TOKENS, SLIDE_W, SLIDE_H)

    def test_element_count(self):
        els = self.els
        self.assertEqual(len(els), 6)
        numbers = [e for e in els if e["role"] == "stat-number"]
        labels = [e for e in els if e["role"] == "stat-label"]
        self.assertEqual(len(numbers), 3)
        self.assertEqual(len(labels), 3)

    def test_margin_alignment(self):
        els = self.els
        left_min = min(e["left"] for e in els)
        right_max = max(e["left"] + e["width"] for e in els)
        self.assertEqual(left_min, 457200)
        self.assertEqual(right_max, SLIDE_W - 457200)
        self.assertEqual(right_max, 8686800)

    def test_number_and_label_style(self):
        els = self.els
        number = next(e for e in els if e["role"] == "stat-number")
        label = next(e for e in els if e["role"] == "stat-label")
        self.assertEqual(number["font_pt"], 40.0)
        self.assertEqual(number["colour"], "#4F81BD")
        self.assertEqual(label["font_pt"], 12.0)
        self.assertEqual(label["colour"], "#000000")

    def test_cells_separated_by_gutter(self):
        gutter = TOKENS["grid"]["gutter"]
        numbers = sorted(
            [e for e in self.els if e["role"] == "stat-number"],
            key=lambda e: e["left"],
        )
        for idx in range(len(numbers) - 1):
            prev = numbers[idx]
            nxt = numbers[idx + 1]
            # No overlap.
            self.assertLessEqual(prev["left"] + prev["width"], nxt["left"])
            # Gap is exactly gutter for all but the last boundary (last cell
            # absorbs rounding remainder, so gap >= gutter).
            if idx < len(numbers) - 2:
                self.assertEqual(nxt["left"], prev["left"] + prev["width"] + gutter)
            else:
                self.assertGreaterEqual(nxt["left"], prev["left"] + prev["width"] + gutter)

    def test_number_above_label_no_overlap(self):
        baseline = TOKENS["grid"]["baseline"]
        numbers = {i: e for i, e in enumerate(self.els) if e["role"] == "stat-number"}
        labels = {i: e for i, e in enumerate(self.els) if e["role"] == "stat-label"}
        # Elements alternate: number, label, number, label...
        num_list = [e for e in self.els if e["role"] == "stat-number"]
        lab_list = [e for e in self.els if e["role"] == "stat-label"]
        for num, lab in zip(num_list, lab_list):
            # Number top + height + baseline == label top.
            self.assertEqual(num["top"] + num["height"] + baseline, lab["top"])
            # No vertical overlap.
            self.assertLessEqual(num["top"] + num["height"], lab["top"])

    def test_empty_stats_raises(self):
        with self.assertRaises(ShapeError):
            plan_stat_row([], TOKENS, SLIDE_W, SLIDE_H)

    def test_draw_adds_nonplaceholder_shapes(self):
        prs = Presentation(TEMPLATE)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        els = self.els
        shapes = primitives.draw(slide, els)
        non_placeholder = [s for s in slide.shapes if not s.is_placeholder]
        self.assertEqual(len(non_placeholder), 6)

    def test_normalise_hex(self):
        self.assertEqual(_normalise_hex("#abc"), "#AABBCC")
        self.assertEqual(_normalise_hex("4f81bd"), "#4F81BD")
        self.assertIsNone(_normalise_hex("nope"))

    def test_draw_renders_box_with_token_fill(self):
        prs = Presentation(TEMPLATE)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        box = {
            "role": "card-panel", "kind": "box", "container": True,
            "fill": "#4F81BD",
            "left": 457200, "top": 1600200, "width": 3000000, "height": 2000000,
        }
        text = {
            "role": "card-title", "text": "Hi",
            "left": 600000, "top": 1700000, "width": 2000000, "height": 400000,
            "font_pt": 18.0, "colour": "#FFFFFF",
        }
        shapes = primitives.draw(slide, [text, box])  # unordered on purpose
        non_placeholder = [s for s in slide.shapes if not s.is_placeholder]
        self.assertEqual(len(non_placeholder), 2)
        # The box must paint behind the text: it is added first despite being
        # second in the element list (z-order sort).
        self.assertEqual(shapes[0].fill.fore_color.rgb, RGBColor.from_string("4F81BD"))

    def test_optical_centre_top_biased(self):
        # Default placement sits at the optical centre (~45% from top): the row's
        # vertical centre is ABOVE the band's geometric centre, but still in the
        # middle band (not jammed to the top).
        top = min(e["top"] for e in self.els)
        bottom = max(e["top"] + e["height"] for e in self.els)
        row_centre = (top + bottom) / 2
        band_top = TOKENS["grid"]["margin_top"]
        band_bottom = SLIDE_H - TOKENS["grid"]["margin_bottom"]
        band_centre = (band_top + band_bottom) / 2
        self.assertLess(row_centre, band_centre)
        self.assertGreater(row_centre, band_top + 0.25 * (band_bottom - band_top))


def _roles(els):
    return [e["role"] for e in els]


def _fills(els):
    return {e.get("fill") for e in els if e.get("kind") == "box"}


class TestNewPrimitives(unittest.TestCase):
    """card-grid, comparison, process, timeline: planners are pure, produce
    lint-clean geometry, keep colours on-token, and honour emphasis."""

    def _lint_clean(self, els):
        # Every primitive's default output must pass the hard gate.
        self.assertIsNone(lint.check(els, TOKENS, SLIDE_W, SLIDE_H))

    def test_card_grid_shape_and_lint(self):
        els = plan_card_grid(
            [{"label": "Size"}, {"label": "Know"}, {"label": "Aim"}],
            TOKENS, SLIDE_W, SLIDE_H)
        self._lint_clean(els)
        # 3 panels + 3 labels, no body.
        self.assertEqual(_roles(els).count("card-panel"), 3)
        self.assertEqual(_roles(els).count("card-label"), 3)
        # Panels are paper with an ink outline (grouped, grey field).
        self.assertEqual(_fills(els), {"#FFFFFF"})

    def test_card_grid_emphasis_fills_accent(self):
        els = plan_card_grid(
            [{"label": "a"}, {"label": "b", "emphasis": True}],
            TOKENS, SLIDE_W, SLIDE_H)
        self._lint_clean(els)
        self.assertIn("#4F81BD", _fills(els))  # the hero card fills with accent

    def test_card_grid_rejects_too_many(self):
        with self.assertRaises(ShapeError):
            plan_card_grid([{"label": str(i)} for i in range(7)],
                           TOKENS, SLIDE_W, SLIDE_H)

    def test_comparison_two_panels(self):
        els = plan_comparison(
            [{"header": "Before", "body": "a / b"},
             {"header": "After", "body": "b / a", "emphasis": True}],
            TOKENS, SLIDE_W, SLIDE_H)
        self._lint_clean(els)
        self.assertEqual(_roles(els).count("comparison-panel"), 2)
        self.assertIn("#4F81BD", _fills(els))  # the winning side tilts to accent

    def test_comparison_needs_exactly_two(self):
        with self.assertRaises(ShapeError):
            plan_comparison([{"header": "solo"}], TOKENS, SLIDE_W, SLIDE_H)

    def test_process_numbered_with_connectors(self):
        els = plan_process(
            [{"label": "Plan"}, {"label": "Create"}, {"label": "Deliver"}],
            TOKENS, SLIDE_W, SLIDE_H)
        self._lint_clean(els)
        self.assertEqual(_roles(els).count("process-step"), 3)
        self.assertEqual(_roles(els).count("process-connector"), 2)  # n-1
        nums = [e["text"] for e in els if e["role"] == "process-number"]
        self.assertEqual(nums, ["1", "2", "3"])
        # The number is the accent; the label is ink (lead by size + colour).
        num = next(e for e in els if e["role"] == "process-number")
        self.assertEqual(num["colour"], "#4F81BD")

    def test_process_rejects_over_five(self):
        with self.assertRaises(ShapeError):
            plan_process([{"label": str(i)} for i in range(6)],
                         TOKENS, SLIDE_W, SLIDE_H)

    def test_timeline_dots_rail_labels(self):
        els = plan_timeline(
            [{"date": "26", "event": "Kick"},
             {"date": "27", "event": "Ship", "emphasis": True},
             {"date": "28", "event": "Scale"}],
            TOKENS, SLIDE_W, SLIDE_H)
        self._lint_clean(els)
        self.assertEqual(_roles(els).count("timeline-dot"), 3)
        self.assertEqual(_roles(els).count("timeline-rail"), 2)   # n-1 segments
        self.assertEqual(_roles(els).count("timeline-label"), 3)
        # The emphasised milestone dot is the accent; the rest are muted to ink.
        dots = [e for e in els if e["role"] == "timeline-dot"]
        self.assertIn("#4F81BD", {d["fill"] for d in dots})

    def test_freeform_plans_lint_clean_kit(self):
        # The whole freeform kit (panel/text/arrow/dot/line) plans lint-clean
        # when placed without overlap; text sits inside its container panel.
        els = plan_freeform([
            {"kind": "box", "fill": "paper", "stroke": "ink",
             "placement": {"cols": (1, 6), "rows": (1, 10)}},
            {"kind": "text", "scale": "h1", "colour": "ink", "text": "In panel",
             "placement": {"cols": (1, 6), "rows": (1, 4)}},
            {"kind": "arrow", "colour": "ink",
             "placement": {"cols": (7, 8), "rows": (5, 6)}},
            {"kind": "dot", "colour": "accent",
             "placement": {"cols": (9, 10), "rows": (1, 3)}},
            {"kind": "line", "colour": "ink",
             "placement": {"cols": (1, 12), "rows": (12, 12)}},
        ], TOKENS, SLIDE_W, SLIDE_H)
        self._lint_clean(els)
        roles = set(_roles(els))
        self.assertEqual(roles, {"freeform-panel", "freeform-text",
                                 "freeform-arrow", "freeform-dot", "freeform-line"})

    def test_freeform_colours_resolve_to_tokens(self):
        els = plan_freeform([
            {"kind": "box", "fill": "accent",
             "placement": {"cols": (1, 12), "rows": (1, 8)}},
        ], TOKENS, SLIDE_W, SLIDE_H)
        self.assertEqual(els[0]["fill"], "#4F81BD")  # role name -> brand hex

    def test_cycle_ring(self):
        for n in (2, 3, 4, 5, 6):
            els = plan_cycle([{"label": f"S{i}"} for i in range(n)],
                             TOKENS, SLIDE_W, SLIDE_H)
            self._lint_clean(els)
            self.assertEqual(_roles(els).count("cycle-node"), n)
            self.assertEqual(_roles(els).count("cycle-edge"), n)  # closes the loop

    def test_cycle_rejects_over_six(self):
        with self.assertRaises(ShapeError):
            plan_cycle([{"label": str(i)} for i in range(7)], TOKENS, SLIDE_W, SLIDE_H)

    def test_matrix_quadrants(self):
        spec = {"x": "Effort", "y": "Impact", "quadrants": [
            {"label": "QW"}, {"label": "BB", "emphasis": True},
            {"label": "DP"}, {"label": "FI"}]}
        els = plan_matrix(spec, TOKENS, SLIDE_W, SLIDE_H)
        self._lint_clean(els)
        self.assertEqual(_roles(els).count("matrix-cell"), 4)
        self.assertEqual(_roles(els).count("matrix-axis"), 2)  # x + y captions
        self.assertIn("#4F81BD", _fills(els))  # emphasised quadrant

    def test_matrix_needs_four(self):
        with self.assertRaises(ShapeError):
            plan_matrix({"quadrants": [{"label": "a"}, {"label": "b"}]},
                        TOKENS, SLIDE_W, SLIDE_H)

    def test_process_icon_replaces_number(self):
        els = plan_process([{"label": "Plan", "icon": "idea"},
                            {"label": "Ship", "icon": "fast"}],
                           TOKENS, SLIDE_W, SLIDE_H)
        self._lint_clean(els)
        self.assertEqual(_roles(els).count("process-icon"), 2)
        self.assertEqual(_roles(els).count("process-number"), 0)  # icon takes the slot

    def test_comparison_icon(self):
        els = plan_comparison([{"header": "Before", "body": "slow", "icon": "decline"},
                               {"header": "After", "body": "fast", "icon": "growth",
                                "emphasis": True}], TOKENS, SLIDE_W, SLIDE_H)
        self._lint_clean(els)
        self.assertEqual(_roles(els).count("comparison-icon"), 2)

    def test_icon_list_rows(self):
        rows = [{"icon": "growth", "text": "Up"}, {"icon": "team", "text": "Bigger"}]
        els = plan_icon_list(rows, TOKENS, SLIDE_W, SLIDE_H)
        self._lint_clean(els)
        self.assertEqual(_roles(els).count("iconlist-icon"), 2)
        self.assertEqual(_roles(els).count("iconlist-text"), 2)
        icon = next(e for e in els if e["role"] == "iconlist-icon")
        self.assertEqual(icon["kind"], "icon")
        self.assertEqual(icon["colour"], "#4F81BD")  # accent marker

    def test_card_with_icon_lint_clean(self):
        cards = [{"label": "Grow", "icon": "growth"},
                 {"label": "Serve", "icon": "team"}]
        els = plan_card_grid(cards, TOKENS, SLIDE_W, SLIDE_H)
        self._lint_clean(els)
        self.assertEqual(_roles(els).count("card-icon"), 2)

    def test_panel_corner_honours_shape_token(self):
        # Default (no shape token) -> rounded; a sharp brand -> plain rectangle.
        cards = [{"label": "A"}, {"label": "B"}, {"label": "C"}]
        default = plan_card_grid(cards, TOKENS, SLIDE_W, SLIDE_H)
        panel = next(e for e in default if e["role"] == "card-panel")
        self.assertEqual(panel["shape"], "rounded_rectangle")

        sharp = dict(TOKENS, shape={"corner": "sharp"})
        panel2 = next(e for e in plan_card_grid(cards, sharp, SLIDE_W, SLIDE_H)
                      if e["role"] == "card-panel")
        self.assertEqual(panel2["shape"], "rect")

    def test_tree_org_chart_lint_clean(self):
        root = _node("CEO", _node("Eng"), _node("Sales"), _node("Ops"))
        els = plan_tree(root, TOKENS, SLIDE_W, SLIDE_H)
        self._lint_clean(els)
        self.assertEqual(_roles(els).count("tree-node"), 4)
        self.assertEqual(_roles(els).count("tree-edge"), 3)  # N-1 edges

    def test_tree_depth_three_lint_clean(self):
        root = _node("R", _node("A", _node("A1"), _node("A2")), _node("B"))
        self._lint_clean(plan_tree(root, TOKENS, SLIDE_W, SLIDE_H))

    def test_tree_emphasis_fills_accent(self):
        root = _node("Vision", _node("Now", emphasis=True), _node("Next"))
        els = plan_tree(root, TOKENS, SLIDE_W, SLIDE_H)
        self.assertIn("#4F81BD", _fills(els))

    def test_tree_with_icons_lower_cap(self):
        # 7 icon'd nodes exceeds the with-icons cap of 6.
        kids = [_node(f"c{i}", icon="growth") for i in range(6)]
        root = _node("root", *kids)  # 7 nodes, all-but-root icon'd
        with self.assertRaises(ShapeError):
            plan_tree(root, TOKENS, SLIDE_W, SLIDE_H)

    def test_tree_over_cap_raises(self):
        root = _node("R", *[_node(f"c{i}") for i in range(9)])  # 10 nodes
        with self.assertRaises(ShapeError):
            plan_tree(root, TOKENS, SLIDE_W, SLIDE_H)

    def test_tree_too_deep_raises(self):
        root = _node("a", _node("b", _node("c", _node("d", _node("e")))))
        with self.assertRaises(ShapeError):
            plan_tree(root, TOKENS, SLIDE_W, SLIDE_H)

    def test_all_defaults_stay_under_cap(self):
        # 5-step process with details is the worst case; must fit the cap.
        els = plan_process(
            [{"label": f"S{i}", "detail": "why & who"} for i in range(5)],
            TOKENS, SLIDE_W, SLIDE_H)
        self.assertLessEqual(len(els), lint.ELEMENT_CAP)
        self._lint_clean(els)


# A small header + two data rows: one text column, two money/percent columns.
TABLE = {
    "header": ["Metric", "Q1", "Q2"],
    "rows": [
        {"cells": ["Revenue", "$1.2M", "$1.5M"], "emphasis": False},
        {"cells": ["Growth", "4%", "9%"], "emphasis": False},
    ],
}


class TestPlanTable(unittest.TestCase):
    """plan_table (D-001/D-002/D-004/D-008/D-009): ONE native-table element that
    carries its colours/sizes as vectors and per-column alignment; hard caps and
    a band-fit guard raise ShapeError.

    Contract pinned here (the implementer follows these exact keys):
      element keys : kind "table", role "table-grid", text, left/top/width/height,
                     fills [hex], text_colours [hex], font_pts [pt], stroke (hairline),
                     col_aligns [ "left" | "right" ]  (per column, header order)
      alignment key: col_aligns
    """

    # Grid/band geometry the existing token fixture implies.
    MARGIN_X = 457200
    CONTENT_LEFT = 457200
    CONTENT_RIGHT = SLIDE_W - 457200        # 8686800
    BAND_TOP = 1600200
    BAND_BOTTOM = SLIDE_H - 731837          # 6126163

    def _plan(self, table, region=None):
        return primitives.plan_table(table, TOKENS, SLIDE_W, SLIDE_H, region)

    def test_returns_single_table_element(self):
        els = self._plan(TABLE)
        self.assertEqual(len(els), 1)          # D-001: one GraphicFrame, not per-cell boxes
        el = els[0]
        self.assertEqual(el["kind"], "table")
        self.assertEqual(el["role"], "table-grid")
        # A text key so a lint margins violation can name the element.
        self.assertIn("text", el)
        self.assertIsInstance(el["text"], str)
        # D-008 vector keys the lint validates.
        self.assertIsInstance(el["fills"], list)
        self.assertIsInstance(el["text_colours"], list)
        self.assertIsInstance(el["font_pts"], list)

    def test_geometry_within_content_span_and_band(self):
        el = self._plan(TABLE)[0]
        # Spans margin to margin (the _content_span rule the other planners use).
        self.assertEqual(el["left"], self.CONTENT_LEFT)
        self.assertEqual(el["left"] + el["width"], self.CONTENT_RIGHT)
        # Sits inside the band.
        self.assertGreaterEqual(el["top"], self.BAND_TOP)
        self.assertLessEqual(el["top"] + el["height"], self.BAND_BOTTOM)
        self.assertGreater(el["width"], 0)
        self.assertGreater(el["height"], 0)

    def test_optical_centre_top(self):
        # Same optical-centre placement as plan_card_grid: the block's vertical
        # centre sits above the band centre, but not jammed to the top.
        el = self._plan(TABLE)[0]
        band_centre = (self.BAND_TOP + self.BAND_BOTTOM) / 2
        row_centre = el["top"] + el["height"] / 2
        self.assertLess(row_centre, band_centre)
        self.assertGreater(
            row_centre, self.BAND_TOP + 0.25 * (self.BAND_BOTTOM - self.BAND_TOP))

    def test_default_fills_text_colours_and_font(self):
        # D-002: header = ink fill + paper text; data = paper fill + ink text;
        # a muted bottom hairline rides the scalar `stroke` key (D-008).
        el = self._plan(TABLE)[0]
        ink = TOKENS["colour_roles"]["ink"]
        paper = TOKENS["colour_roles"]["paper"]
        muted = TOKENS["colour_roles"]["muted"]
        self.assertIn(ink, el["fills"])       # header band
        self.assertIn(paper, el["fills"])     # data band
        self.assertIn(ink, el["text_colours"])
        self.assertIn(paper, el["text_colours"])
        self.assertEqual(el["stroke"], muted)  # row hairline
        # Body type-scale size throughout (no shrink-to-fit).
        self.assertIn(TOKENS["type_scale"]["body"], el["font_pts"])
        self.assertTrue(all(p == TOKENS["type_scale"]["body"] for p in el["font_pts"]))

    def test_rejects_nine_data_rows(self):
        table = {"header": ["A", "B", "C"],
                 "rows": [{"cells": [str(i), "x", "y"], "emphasis": False}
                          for i in range(9)]}          # cap is 8 data rows
        with self.assertRaises(ShapeError):
            self._plan(table)

    def test_rejects_six_columns(self):
        table = {"header": ["A", "B", "C", "D", "E", "F"],
                 "rows": [{"cells": ["1", "2", "3", "4", "5", "6"],
                           "emphasis": False}]}          # cap is 5 columns
        with self.assertRaises(ShapeError):
            self._plan(table)

    def test_rejects_one_column(self):
        table = {"header": ["Only"],
                 "rows": [{"cells": ["a"], "emphasis": False},
                          {"cells": ["b"], "emphasis": False}]}  # a list, not a table
        with self.assertRaises(ShapeError):
            self._plan(table)

    def test_band_overflow_raises_named(self):
        # A short placed region: (1 + n_rows) * row_h cannot fit, so the band-fit
        # guard fires before drawing and the message names how many rows fit /
        # tells the author to cut rows or split the slide (D-004).
        short = (self.MARGIN_X, self.BAND_TOP,
                 SLIDE_W - 2 * self.MARGIN_X, 400000)
        table = {"header": ["A", "B", "C"],
                 "rows": [{"cells": [str(i), "x", "y"], "emphasis": False}
                          for i in range(3)]}
        with self.assertRaises(ShapeError) as cm:
            self._plan(table, region=short)
        self.assertIn("row", str(cm.exception).lower())

    def test_emphasis_row_accent_fill_paper_text(self):
        # REQ-002 / D-002: an emphasised row fills with accent and reverses to
        # paper text.
        table = {"header": ["Metric", "Q1", "Q2"],
                 "rows": [
                     {"cells": ["Revenue", "$1.2M", "$1.5M"], "emphasis": False},
                     {"cells": ["Total", "$2.7M", "$3.0M"], "emphasis": True},
                 ]}
        el = self._plan(table)[0]
        self.assertIn(TOKENS["colour_roles"]["accent"], el["fills"])
        self.assertIn(TOKENS["colour_roles"]["paper"], el["text_colours"])

    def test_numeric_columns_right_text_left(self):
        # D-009: a column whose data cells are all numeric (money/percent) is
        # right-aligned; a text column stays left. col_aligns is header order.
        el = self._plan(TABLE)[0]
        aligns = el["col_aligns"]
        self.assertEqual(len(aligns), 3)
        self.assertEqual(aligns[0], "left")    # "Revenue" / "Growth" — text
        self.assertEqual(aligns[1], "right")   # "$1.2M" / "4%"      — numeric
        self.assertEqual(aligns[2], "right")   # "$1.5M" / "9%"      — numeric


if __name__ == "__main__":
    unittest.main()
