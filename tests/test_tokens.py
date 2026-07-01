"""Tests for source/skills/build-deck/scripts/tokens.py.

Stdlib unittest only — no pytest, no third-party test deps beyond python-pptx
(already a runtime dependency of the build-deck scripts).

Run from the repo root:

    python3 -m unittest tests.test_tokens -v

NOTE on fixture values
----------------------
derive_grid measures only the layouts the brand actually maps (its layout_map
values) — the right "brand surface". On sample-template.pptx the mapped role
layouts are indices 0,1,2,3,5 (title, title-content, section, two-column,
statement); measuring those yields margin_x=457200, margin_top=274638 (the
topmost content edge — a title placeholder), gutter=152400 (the true two-column
gutter), margin_bottom=731837 (slide 9144000 x 6858000 EMU). Measuring ALL
layouts instead pulls in the bundled default template's unused extras (e.g.
'Content with Caption'), which narrow the gutter to 109537 — exactly the
pollution the layout restriction avoids.
"""

import os
import sys
import unittest

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SCRIPTS = os.path.join(REPO_ROOT, "source", "skills", "build-deck", "scripts")
sys.path.insert(0, SCRIPTS)

import tokens  # noqa: E402

from pptx import Presentation  # noqa: E402

TEMPLATE = os.path.join(REPO_ROOT, "tests", "fixtures", "sample-template.pptx")


class TestGridFromRects(unittest.TestCase):
    """grid_from_rects: EMU geometry → grid sub-dict."""

    # The fixture's mapped role layouts (title, title-content, section,
    # two-column, statement). quote shares section in this fixture.
    ROLE_LAYOUTS = [0, 1, 2, 3, 5]

    def test_derive_grid_role_layouts(self):
        """derive_grid restricted to the brand's mapped layouts → clean values."""
        prs = Presentation(TEMPLATE)
        g = tokens.derive_grid(prs, self.ROLE_LAYOUTS)

        # Measured from layouts 0,1,2,3,5 of sample-template.pptx
        # (slide 9144000 x 6858000 EMU). margin_x = the symmetric content left
        # edge (457200); margin_top = the topmost content edge (a title
        # placeholder at 274638); gutter = the two-column gap (152400).
        self.assertEqual(g["margin_x"], 457200)
        self.assertEqual(g["margin_top"], 274638)
        self.assertEqual(g["gutter"], 152400)
        self.assertEqual(g["margin_bottom"], 731837)
        self.assertEqual(g["columns"], 12)
        self.assertEqual(g["baseline"], 91211)

    def test_derive_grid_all_layouts_is_superset(self):
        """No layout_indices → measure all layouts (the documented fallback).

        Adding the bundled default's extra layouts can only lower the minimum
        top and never raises margin_x above the role-layout value. Asserted as
        inequalities so the test does not pin python-pptx's internal template.
        """
        prs = Presentation(TEMPLATE)
        g = tokens.derive_grid(prs)
        self.assertEqual(g["margin_x"], 457200)
        self.assertEqual(g["columns"], 12)
        self.assertLessEqual(g["margin_top"], 274638)
        self.assertGreater(g["gutter"], 0)

    def test_grid_from_rects_fallback(self):
        """Empty rects → proportional fallback values."""
        g = tokens.grid_from_rects([], 9144000, 6858000)
        self.assertEqual(g["margin_x"], round(9144000 * 0.05))
        self.assertEqual(g["margin_x"], 457200)
        self.assertEqual(g["margin_top"], round(6858000 * 0.08))
        self.assertEqual(g["margin_bottom"], round(6858000 * 0.08))
        self.assertEqual(g["gutter"], round(9144000 * 0.0167))
        self.assertEqual(g["columns"], 12)

    def test_grid_from_rects_simple(self):
        """Single rect: margins come from the rect edges."""
        # One 1000×1000 rect at (200, 300) in a 2000×2000 slide.
        rects = [(200, 300, 1000, 1000)]
        g = tokens.grid_from_rects(rects, 2000, 2000)
        self.assertEqual(g["margin_x"], 200)
        self.assertEqual(g["margin_top"], 300)
        self.assertEqual(g["margin_bottom"], 2000 - (300 + 1000))
        self.assertEqual(g["columns"], 12)
        # No same-top pair → fallback gutter.
        self.assertEqual(g["gutter"], round(2000 * 0.0167))

    def test_grid_from_rects_gutter_detected(self):
        """Two rects at the same top → gutter is the gap between them."""
        # Rect A: left=100, top=200, width=500, height=300  → right=600
        # Rect B: left=650, top=200, width=500, height=300
        # gap = 650 - 600 = 50
        rects = [(100, 200, 500, 300), (650, 200, 500, 300)]
        g = tokens.grid_from_rects(rects, 2000, 1000)
        self.assertEqual(g["gutter"], 50)

    def test_grid_from_rects_gutter_tolerance(self):
        """Rects within 12700 EMU vertical offset are treated as same row."""
        # Rect A: left=100, top=200, width=400, height=300  → right=500
        # Rect B: left=560, top=12700, width=400, height=300
        # |200 - 12700| = 12500 ≤ 12700  → gap = 560 - 500 = 60
        rects = [(100, 200, 400, 300), (560, 12700, 400, 300)]
        g = tokens.grid_from_rects(rects, 2000, 600000)
        self.assertEqual(g["gutter"], 60)

    def test_grid_from_rects_gutter_outside_tolerance(self):
        """Rects more than 12700 EMU apart vertically produce fallback gutter."""
        # |200 - 213000| = 12800 > 12700  → no valid pair
        rects = [(100, 200, 400, 300), (560, 213000, 400, 300)]
        g = tokens.grid_from_rects(rects, 2000, 600000)
        self.assertEqual(g["gutter"], round(2000 * 0.0167))


class TestDefaultTypeScale(unittest.TestCase):
    """default_type_scale returns the correct generic defaults."""

    def test_default_type_scale(self):
        ts = tokens.default_type_scale()
        self.assertEqual(set(ts), {"display", "h1", "body", "caption"})
        self.assertEqual(ts["display"], 40.0)
        self.assertEqual(ts["h1"], 28.0)
        self.assertEqual(ts["body"], 18.0)
        self.assertEqual(ts["caption"], 12.0)


class TestDeriveScale(unittest.TestCase):
    """_derive_scale: a monotonic brand scale, or None to fall back."""

    def test_normal_case_monotonic(self):
        # display is a hero size above the title (44*1.4 = 62); h1 = title.
        scale = tokens._derive_scale(44, 32)
        self.assertEqual(scale, {"display": 62.0, "h1": 44.0,
                                 "body": 32.0, "caption": 24.0})
        self.assertTrue(scale["display"] > scale["h1"] > scale["body"]
                        > scale["caption"])
        # A hero number still dominates a caption label >= 2.5x.
        self.assertGreaterEqual(scale["display"] / scale["caption"], 2.5)

    def test_inverted_falls_back(self):
        # title <= body would invert the scale (the fixture's own risk) -> None.
        self.assertIsNone(tokens._derive_scale(32, 44))
        self.assertIsNone(tokens._derive_scale(32, 32))

    def test_missing_falls_back(self):
        self.assertIsNone(tokens._derive_scale(None, 32))
        self.assertIsNone(tokens._derive_scale(44, None))

    def test_narrow_range_still_monotonic(self):
        # title just above body still derives (h1=title > body); display leads.
        scale = tokens._derive_scale(33, 32)
        self.assertIsNotNone(scale)
        self.assertTrue(scale["display"] > scale["h1"] > scale["body"]
                        > scale["caption"])


class TestResolveColourRoles(unittest.TestCase):
    """resolve_colour_roles: name->hex dict → canonical colour roles."""

    def test_resolve_colour_roles_fixture(self):
        """Explicit ink / paper / accent / accent2 → deterministic mapping."""
        roles = tokens.resolve_colour_roles(
            {
                "accent": "#4F81BD",
                "accent2": "#C0504D",
                "ink": "#000000",
                "paper": "#FFFFFF",
            }
        )
        self.assertEqual(
            roles,
            {
                "ink": "#000000",
                "paper": "#FFFFFF",
                "accent": "#4F81BD",
                "muted": "#C0504D",
            },
        )

    def test_resolve_colour_roles_empty(self):
        """No valid hex values → empty dict."""
        self.assertEqual(tokens.resolve_colour_roles({}), {})

    def test_resolve_colour_roles_no_valid_hex(self):
        """Unparseable values only → empty dict."""
        self.assertEqual(
            tokens.resolve_colour_roles({"primary": "notacolour", "bg": "###"}),
            {},
        )

    def test_resolve_colour_roles_fallback_by_luminance(self):
        """Without explicit ink/paper keys, darkest → ink, lightest → paper.

        Luminance formula: 0.2126*R + 0.7152*G + 0.0722*B (0-255 channels).
          #4F81BD: R=79, G=129, B=189 → lum ≈ 122.70  (darker)
          #888888: R=G=B=136          → lum  = 136.00  (lighter)
        """
        roles = tokens.resolve_colour_roles(
            {"accent2": "#888888", "accent": "#4F81BD"}
        )
        # accent: explicit key present.
        self.assertEqual(roles["accent"], "#4F81BD")
        # ink: darkest hex by luminance.
        self.assertEqual(roles["ink"], "#4F81BD")
        # paper: lightest hex by luminance.
        self.assertEqual(roles["paper"], "#888888")
        # muted: accent2 fallback.
        self.assertEqual(roles["muted"], "#888888")

    def test_resolve_colour_roles_accent_first_when_no_key(self):
        """Without an 'accent' key, the first valid hex by insertion order is used."""
        roles = tokens.resolve_colour_roles(
            {"primary": "#FF0000", "secondary": "#0000FF"}
        )
        self.assertEqual(roles["accent"], "#FF0000")

    def test_resolve_colour_roles_normalises_rgb_shorthand(self):
        """'#RGB' shorthand is expanded to '#RRGGBB' uppercase."""
        roles = tokens.resolve_colour_roles({"ink": "#000", "paper": "#fff"})
        self.assertEqual(roles["ink"], "#000000")
        self.assertEqual(roles["paper"], "#FFFFFF")

    def test_resolve_colour_roles_normalises_no_hash(self):
        """Values without a leading '#' are accepted."""
        roles = tokens.resolve_colour_roles({"ink": "000000", "paper": "FFFFFF"})
        self.assertEqual(roles["ink"], "#000000")
        self.assertEqual(roles["paper"], "#FFFFFF")

    def test_resolve_colour_roles_muted_falls_back_to_accent(self):
        """With no 'muted' or 'accent2' key, muted equals accent."""
        roles = tokens.resolve_colour_roles(
            {"accent": "#4F81BD", "ink": "#000000", "paper": "#FFFFFF"}
        )
        self.assertEqual(roles["muted"], "#4F81BD")


class TestResolveTokens(unittest.TestCase):
    """resolve_tokens: merge derived defaults with explicit brand overrides."""

    def test_resolve_tokens_default_and_override(self):
        """Explicit tokens override derived values; unoverridden keys survive."""
        prs = Presentation(TEMPLATE)
        brand = {
            "colours": {"accent": "#4F81BD", "ink": "#000000", "paper": "#FFFFFF"},
            "layout_map": {"title": 0, "title-content": 1, "section": 2,
                           "two-column": 3, "statement": 5, "quote": 2},
            "tokens": {"type_scale": {"body": 21.0}},
        }
        t = tokens.resolve_tokens(brand, prs)

        # Override wins.
        self.assertEqual(t["type_scale"]["body"], 21.0)
        # The unoverridden keys are now the brand scale read from the master
        # (title 44 / body 32 -> hero display 62), not the generic default.
        self.assertEqual(t["type_scale"]["display"], 62.0)
        # Grid is measured from the MAPPED layouts (not all layouts): the
        # two-column layout's margins flow through resolve_tokens.
        self.assertEqual(t["grid"]["margin_x"], 457200)
        self.assertEqual(t["grid"]["margin_top"], 274638)
        self.assertEqual(t["grid"]["gutter"], 152400)
        # Colour role derived from brand colours.
        self.assertEqual(t["colour_roles"]["accent"], "#4F81BD")

    def test_resolve_tokens_no_overrides(self):
        """With no explicit tokens, the scale is derived from the master's real
        sizes (title 44 / body 32 -> 44 / 38 / 32 / 24), monotonic."""
        prs = Presentation(TEMPLATE)
        brand = {"colours": {"ink": "#111111", "paper": "#EEEEEE"}}
        t = tokens.resolve_tokens(brand, prs)

        ts = t["type_scale"]
        self.assertEqual(ts["display"], 62.0)
        self.assertEqual(ts["h1"], 44.0)
        self.assertEqual(ts["body"], 32.0)
        self.assertEqual(ts["caption"], 24.0)
        self.assertGreater(ts["display"], ts["h1"])
        self.assertGreater(ts["h1"], ts["body"])
        self.assertGreater(ts["body"], ts["caption"])
        self.assertIn("margin_x", t["grid"])

    def test_resolve_tokens_empty_brand(self):
        """Empty brand dict produces derived tokens with no colour_roles."""
        prs = Presentation(TEMPLATE)
        t = tokens.resolve_tokens({}, prs)
        self.assertIn("grid", t)
        self.assertIn("type_scale", t)
        # No valid colours → empty colour_roles.
        self.assertEqual(t["colour_roles"], {})

    def test_resolve_tokens_shape_default_and_override(self):
        prs = Presentation(TEMPLATE)
        # Default shape language is present and rounded (back-compat).
        t = tokens.resolve_tokens({}, prs)
        self.assertEqual(t["shape"]["corner"], "rounded")
        # A brand can set a sharp corner language.
        t2 = tokens.resolve_tokens({"tokens": {"shape": {"corner": "sharp"}}}, prs)
        self.assertEqual(t2["shape"]["corner"], "sharp")

    def test_resolve_tokens_explicit_grid_override(self):
        """Explicit grid values in brand['tokens'] override derived grid."""
        prs = Presentation(TEMPLATE)
        brand = {"tokens": {"grid": {"columns": 16, "gutter": 99999}}}
        t = tokens.resolve_tokens(brand, prs)
        self.assertEqual(t["grid"]["columns"], 16)
        self.assertEqual(t["grid"]["gutter"], 99999)
        # Non-overridden key still derived.
        self.assertIn("margin_x", t["grid"])


if __name__ == "__main__":
    unittest.main()
