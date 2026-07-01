"""Tests for init_brand.py — one-step brand.json (heuristic layout_map) from a deck.

Stdlib unittest. Runs init_brand.py as a subprocess (like test_render.py /
test_extract_brand.py) and also unit-tests the suggest_layout_map heuristic
directly. The end-to-end test feeds the generated brand.json to render.py and
renders the committed sample deck. Until init_brand.py exists these FAIL.

    python3 -m unittest tests.test_init_brand
    python3 -m unittest discover tests
"""
import json
import os
import shutil
import subprocess
import sys
import tempfile
import unittest

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SCRIPTS = os.path.join(REPO_ROOT, "source", "skills", "build-deck", "scripts")
FIXTURES = os.path.join(REPO_ROOT, "tests", "fixtures")
TEMPLATE = os.path.join(FIXTURES, "sample-template.pptx")
SAMPLE_DECK = os.path.join(FIXTURES, "sample-deck.md")
INIT_BRAND_PY = os.path.join(SCRIPTS, "init_brand.py")
RENDER_PY = os.path.join(SCRIPTS, "render.py")

CANONICAL_MAP = {
    "title": 0, "title-content": 1, "section": 2,
    "two-column": 3, "statement": 5, "quote": 2,
    "composed": 5,
}


class InitBrandTest(unittest.TestCase):
    """Happy path against the fixture, plus the end-to-end render."""

    @classmethod
    def setUpClass(cls):
        cls._tmp = tempfile.mkdtemp(prefix="slides-initbrand-")
        cls.result = subprocess.run(
            [sys.executable, INIT_BRAND_PY, TEMPLATE],
            capture_output=True, text=True,
        )

    @classmethod
    def tearDownClass(cls):
        shutil.rmtree(cls._tmp, ignore_errors=True)

    def _data(self):
        self.assertEqual(
            self.result.returncode, 0,
            f"init_brand exited {self.result.returncode}.\n"
            f"stdout: {self.result.stdout}\nstderr: {self.result.stderr}",
        )
        return json.loads(self.result.stdout)

    def test_complete_brand_json(self):
        data = self._data()
        for key in ("template", "fonts", "colours", "layout_map", "tokens"):
            self.assertIn(key, data)

    def test_layout_map_is_canonical(self):
        self.assertEqual(self._data()["layout_map"], CANONICAL_MAP)

    def test_each_role_layout_has_enough_placeholders(self):
        # Self-documenting sufficiency: every mapped layout has at least the
        # role's required_min content placeholders, so render never overflows.
        # "composed" is a canvas alias, not a render role; skip it here.
        sys.path.insert(0, SCRIPTS)
        import pptxlib
        import render
        prs = pptxlib.load_template(TEMPLATE)
        content = {t.name for t in pptxlib.CONTENT_PLACEHOLDER_TYPES}
        counts = {
            layout["index"]: sum(
                1 for p in layout["placeholders"] if p["type"] in content)
            for layout in pptxlib.list_layouts(prs)
        }
        for role, idx in self._data()["layout_map"].items():
            if role not in render.ROLE_FIELDS:
                continue  # skip canvas aliases like "composed"
            required_min = (len(render.ROLE_FIELDS[role])
                            - len(render.OPTIONAL_FIELDS.get(role, set())))
            self.assertGreaterEqual(
                counts[idx], required_min,
                f"role {role!r} -> layout {idx} has {counts[idx]} content "
                f"placeholders, needs >= {required_min}")

    def test_end_to_end_renders_sample_deck(self):
        # The generated brand.json must render the committed sample deck.
        data = self._data()
        data["template"] = TEMPLATE  # absolute, so render resolves it anywhere
        brand_path = os.path.join(self._tmp, "brand.json")
        out_path = os.path.join(self._tmp, "out.pptx")
        with open(brand_path, "w", encoding="utf-8") as fh:
            json.dump(data, fh)
        result = subprocess.run(
            [sys.executable, RENDER_PY, "--spec", SAMPLE_DECK,
             "--brand", brand_path, "--out", out_path],
            capture_output=True, text=True,
        )
        self.assertEqual(
            result.returncode, 0,
            f"render failed.\nstdout: {result.stdout}\nstderr: {result.stderr}")
        from pptx import Presentation
        self.assertEqual(len(Presentation(out_path).slides), 6)

    def test_composed_layout_map_entry(self):
        data = self._data()
        layout_map = data["layout_map"]
        self.assertIn("composed", layout_map,
                      "layout_map must contain a 'composed' key")
        # The fixture has a "statement" layout; composed should alias it.
        self.assertEqual(
            layout_map["composed"], layout_map["statement"],
            "composed must equal statement index when statement is mapped",
        )

    def test_tokens_structure(self):
        data = self._data()
        self.assertIn("tokens", data, "brand JSON must contain a 'tokens' object")
        tok = data["tokens"]

        # grid: must have margin_x and columns
        self.assertIn("grid", tok)
        grid = tok["grid"]
        self.assertIn("margin_x", grid)
        self.assertIn("columns", grid)

        # type_scale: the fixture master's real sizes (title 44 / body 32) derive
        # a hero display of 62 (44*1.4); a bare template would fall back to 40.
        self.assertIn("type_scale", tok)
        self.assertAlmostEqual(
            tok["type_scale"]["display"], 62.0,
            msg="type_scale.display should be the brand-derived hero size",
        )
        self.assertGreater(tok["type_scale"]["display"], tok["type_scale"]["h1"])

        # colour_roles: non-empty, must contain accent, ink, paper
        self.assertIn("colour_roles", tok)
        colour_roles = tok["colour_roles"]
        self.assertTrue(colour_roles, "colour_roles must be non-empty")
        for role in ("accent", "ink", "paper"):
            self.assertIn(role, colour_roles,
                          f"colour_roles must contain '{role}'")

    def test_template_ref_override(self):
        result = subprocess.run(
            [sys.executable, INIT_BRAND_PY, TEMPLATE,
             "--template-ref", "template.pptx"],
            capture_output=True, text=True,
        )
        self.assertEqual(result.returncode, 0, result.stderr)
        self.assertEqual(json.loads(result.stdout)["template"], "template.pptx")


class SuggestLayoutMapUnitTest(unittest.TestCase):
    """Direct unit test of the heuristic, incl. the unhostable-role error."""

    def setUp(self):
        sys.path.insert(0, SCRIPTS)
        import init_brand
        import pptxlib
        import render
        self.init_brand = init_brand
        self.role_fields = render.ROLE_FIELDS
        self.optional_fields = render.OPTIONAL_FIELDS
        self.content_names = {t.name for t in pptxlib.CONTENT_PLACEHOLDER_TYPES}

    def test_unhostable_role_raises(self):
        layouts = [{"index": 0, "name": "Blank", "placeholders": []}]
        with self.assertRaises(ValueError):
            self.init_brand.suggest_layout_map(
                layouts, self.role_fields, self.optional_fields,
                self.content_names)


class InitBrandErrorTest(unittest.TestCase):
    """Malformed input exits non-zero with a named error and no JSON."""

    @classmethod
    def setUpClass(cls):
        cls._tmp = tempfile.mkdtemp(prefix="slides-initbrand-error-")

    @classmethod
    def tearDownClass(cls):
        shutil.rmtree(cls._tmp, ignore_errors=True)

    def _assert_named_error(self, result, *needles):
        self.assertEqual(result.returncode, 1, f"stderr: {result.stderr}")
        msg = (result.stderr + result.stdout).lower()
        self.assertIn("error:", msg)
        self.assertEqual(result.stdout.strip(), "", "no JSON on error")
        for needle in needles:
            self.assertIn(needle.lower(), msg)

    def test_missing_file(self):
        missing = os.path.join(self._tmp, "nope.pptx")
        result = subprocess.run(
            [sys.executable, INIT_BRAND_PY, missing],
            capture_output=True, text=True)
        self._assert_named_error(result, "not found", "nope.pptx")

    def test_garbage_file(self):
        bad = os.path.join(self._tmp, "bad.pptx")
        with open(bad, "w", encoding="utf-8") as fh:
            fh.write("not a presentation")
        result = subprocess.run(
            [sys.executable, INIT_BRAND_PY, bad],
            capture_output=True, text=True)
        self._assert_named_error(result, "could not open")


if __name__ == "__main__":
    unittest.main()
