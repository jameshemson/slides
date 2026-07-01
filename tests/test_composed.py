"""End-to-end tests for the `composed` deck-spec role (render.py + primitives + lint).

Stdlib unittest only. The render path is exercised via subprocess (like
test_render.py); the off-token lint gate, which is not reachable from spec text,
is exercised by calling the exact lint.check the render branch calls.

Run from the repo root:

    python3 -m unittest tests.test_composed -v
"""
import json
import os
import shutil
import subprocess
import sys
import tempfile
import unittest

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIXTURES = os.path.join(REPO_ROOT, "tests", "fixtures")
TEMPLATE = os.path.join(FIXTURES, "sample-template.pptx")
COMPOSED_SPEC = os.path.join(FIXTURES, "composed-deck.md")
SAMPLE_SPEC = os.path.join(FIXTURES, "sample-deck.md")
SCRIPTS = os.path.join(REPO_ROOT, "source", "skills", "build-deck", "scripts")
RENDER_PY = os.path.join(SCRIPTS, "render.py")

# Measured content margin of the fixture's mapped layouts (the two-column left
# edge). The composed stat_row must align to it exactly (REQ-011).
MARGIN_X = 457200

# A brand.json with tokens left implicit (derived) and composed mapped to the
# statement layout (single title, index 5).
BRAND = {
    "template": TEMPLATE,
    "fonts": {"heading": "Calibri", "body": "Calibri"},
    "colours": {"accent": "#4F81BD", "accent2": "#C0504D",
                "ink": "#000000", "paper": "#FFFFFF"},
    "layout_map": {"title": 0, "title-content": 1, "section": 2,
                   "two-column": 3, "statement": 5, "quote": 2, "composed": 5},
}


def _run(spec_text, out_name="out.pptx", brand=None, env=None):
    """Render spec_text via render.py in a temp dir. Returns (proc, out_path).

    `env`, when given, replaces the subprocess environment (mirrors
    test_render.py's `_render` — used to hide matplotlib via a PYTHONPATH
    shim).
    """
    tmp = tempfile.mkdtemp()
    spec_path = os.path.join(tmp, "deck.md")
    with open(spec_path, "w", encoding="utf-8") as fh:
        fh.write(spec_text)
    brand_path = os.path.join(tmp, "brand.json")
    with open(brand_path, "w", encoding="utf-8") as fh:
        json.dump(brand or BRAND, fh)
    out_path = os.path.join(tmp, out_name)
    proc = subprocess.run(
        [sys.executable, RENDER_PY, "--spec", spec_path,
         "--brand", brand_path, "--out", out_path],
        capture_output=True, text=True, env=env,
    )
    return proc, out_path


def _run_file(spec_path, brand=None):
    tmp = tempfile.mkdtemp()
    brand_path = os.path.join(tmp, "brand.json")
    with open(brand_path, "w", encoding="utf-8") as fh:
        json.dump(brand or BRAND, fh)
    out_path = os.path.join(tmp, "out.pptx")
    proc = subprocess.run(
        [sys.executable, RENDER_PY, "--spec", spec_path,
         "--brand", brand_path, "--out", out_path],
        capture_output=True, text=True,
    )
    return proc, out_path


class ComposedRenderTest(unittest.TestCase):
    """The fixture composed deck renders a stat row aligned to template margins."""

    @classmethod
    def setUpClass(cls):
        cls.proc, cls.out = _run_file(COMPOSED_SPEC)

    def test_render_succeeds(self):
        self.assertEqual(self.proc.returncode, 0,
                         f"render failed: {self.proc.stderr}\n{self.proc.stdout}")
        self.assertTrue(os.path.isfile(self.out), "no .pptx written")

    def test_six_stat_shapes(self):
        prs = Presentation(self.out)
        slide = prs.slides[0]
        drawn = [s for s in slide.shapes if not s.is_placeholder]
        self.assertEqual(len(drawn), 6,
                         "expected 3 numbers + 3 labels as non-placeholder shapes")

    def test_stat_row_aligns_to_margins(self):
        """REQ-011: over the drawn (non-placeholder) shapes, the row's left edge
        is the derived margin and its right edge is the symmetric margin."""
        prs = Presentation(self.out)
        slide = prs.slides[0]
        drawn = [s for s in slide.shapes if not s.is_placeholder]
        self.assertTrue(drawn, "no drawn shapes")
        left = min(s.left for s in drawn)
        right = max(s.left + s.width for s in drawn)
        self.assertEqual(left, MARGIN_X)
        self.assertEqual(right, prs.slide_width - MARGIN_X)

    def test_title_placeholder_filled(self):
        prs = Presentation(self.out)
        slide = prs.slides[0]
        self.assertIsNotNone(slide.shapes.title)
        self.assertEqual(slide.shapes.title.text, "What moved this quarter")


class ComposedLintGateTest(unittest.TestCase):
    """The mechanical lint is wired into the render gate."""

    def test_over_cap_fails_render(self):
        """13 stats -> 26 elements -> over the element cap -> render fails, no file."""
        spec = (
            "---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
            "Block: stat-row\n"
            + "".join(f"{i} | label{i}\n" for i in range(13))
        )
        proc, out = _run(spec)
        self.assertNotEqual(proc.returncode, 0)
        combined = proc.stderr + proc.stdout
        self.assertIn("error:", combined)
        self.assertIn("[count]", combined)
        self.assertFalse(os.path.isfile(out), "a half-built .pptx was written")

    def test_off_token_colour_blocked_by_gate(self):
        """Off-token colour isn't reachable from spec text, so exercise the exact
        gate the render branch calls: plan a row, mutate a colour off-token,
        and confirm lint.check raises."""
        sys.path.insert(0, SCRIPTS)
        import tokens
        import primitives
        import lint

        prs = Presentation(TEMPLATE)
        toks = tokens.resolve_tokens(BRAND, prs)
        els = primitives.plan_stat_row(
            [{"value": "1", "label": "x"}, {"value": "2", "label": "y"}],
            toks, prs.slide_width, prs.slide_height,
        )
        els[0]["colour"] = "#ABCDEF"  # not in colour_roles
        with self.assertRaises(lint.LintError) as ctx:
            lint.check(els, toks, prs.slide_width, prs.slide_height)
        msg = str(ctx.exception)
        self.assertIn("[colour]", msg)


class ComposedSpecErrorTest(unittest.TestCase):
    """Malformed composed specs fail loudly with no output."""

    def _assert_fails(self, spec, needle):
        proc, out = _run(spec)
        self.assertNotEqual(proc.returncode, 0)
        combined = proc.stderr + proc.stdout
        self.assertIn("error:", combined)
        self.assertIn(needle, combined)
        self.assertFalse(os.path.isfile(out))

    def test_stat_line_without_pipe(self):
        spec = ("---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
                "Block: stat-row\n56 Days to close\n")
        self._assert_fails(spec, "value | label")

    def test_empty_block(self):
        spec = ("---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
                "Block: stat-row\n")
        self._assert_fails(spec, "stat-row")

    def test_no_block(self):
        spec = ("---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
                "Title: Just a title\n")
        self._assert_fails(spec, "Block")

    def test_unknown_block_type(self):
        spec = ("---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
                "Block: pie-tower\n1 | a\n")
        self._assert_fails(spec, "pie-tower")

    def test_mixed_placement_rejected(self):
        # Several blocks must all be placed or all auto-placed, never a mix.
        spec = ("---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
                "Block: card-grid at cols 1-6\nA | x\nBlock: process\nPlan\n")
        self._assert_fails(spec, "mix of placed")

    def test_too_many_blocks_rejected(self):
        spec = ("---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
                + "".join("Block: stat-row\n1 | a\n" for _ in range(5)))
        self._assert_fails(spec, "at most 4 blocks")

    def test_bad_placement_rejected(self):
        spec = ("---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
                "Block: card-grid at cols 1-99\nA | x\n")
        self._assert_fails(spec, "within 1-12")


class BackCompatTest(unittest.TestCase):
    """REQ-009: the existing six-role sample deck renders unchanged."""

    def test_sample_deck_still_renders(self):
        proc, out = _run_file(SAMPLE_SPEC)
        self.assertEqual(proc.returncode, 0,
                         f"sample deck failed: {proc.stderr}\n{proc.stdout}")
        prs = Presentation(out)
        self.assertEqual(len(prs.slides), 6)
        # No composed slides here, so no non-placeholder text shapes (the
        # existing invariant test_render.py also asserts).
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    self.assertTrue(
                        shape.is_placeholder,
                        "sample deck gained a non-placeholder text shape",
                    )


if __name__ == "__main__":
    unittest.main()


class CompositionAdvisoryTest(unittest.TestCase):
    """The advisory composition layer surfaces non-blocking notes in the summary."""

    def test_clean_composed_deck_has_no_advisory(self):
        # The fixture composed deck is good by construction -> no advisory lines.
        proc, out = _run_file(COMPOSED_SPEC)
        self.assertEqual(proc.returncode, 0, proc.stderr)
        self.assertNotIn("advisory", (proc.stdout + proc.stderr).lower())

    def test_weak_composed_deck_advises_without_blocking(self):
        # 6 stats (system cap is 12 elements, so the gate passes) + 5-word
        # labels -> stat-count + label-terseness advisories, render still succeeds.
        spec = (
            "---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
            "Block: stat-row\n"
            + "".join(f"{i} | this is a long label\n" for i in range(6))
        )
        proc, out = _run(spec)
        self.assertEqual(proc.returncode, 0,
                         f"advisory layer must not block: {proc.stderr}")
        self.assertTrue(os.path.isfile(out), "a .pptx should still be written")
        combined = proc.stdout + proc.stderr
        self.assertIn("advisory", combined.lower())
        self.assertIn("stat-count", combined)
        self.assertIn("label-terseness", combined)


class NewPrimitiveRenderTest(unittest.TestCase):
    """card-grid, comparison, process, timeline render as real drawn shapes."""

    def _render_block(self, block_text, out_name):
        spec = (
            "---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
            "Title: T\n" + block_text
        )
        return _run(spec, out_name=out_name)

    def _drawn(self, out):
        prs = Presentation(out)
        slide = prs.slides[0]
        self.assertEqual(slide.shapes.title.text, "T")
        return [s for s in slide.shapes if not s.is_placeholder]

    def test_card_grid_renders_panels(self):
        proc, out = self._render_block(
            "Block: card-grid\nSize | a\nKnowledge | b\n!Aim | c\n", "cards.pptx")
        self.assertEqual(proc.returncode, 0, proc.stderr)
        # 3 panels + 3 labels + 3 bodies = 9 drawn shapes.
        self.assertEqual(len(self._drawn(out)), 9)

    def test_comparison_renders_two_panels(self):
        proc, out = self._render_block(
            "Block: comparison\nBefore | slow\n!After | fast\n", "cmp.pptx")
        self.assertEqual(proc.returncode, 0, proc.stderr)
        self.assertEqual(len(self._drawn(out)), 6)  # 2 panels + 2 headers + 2 bodies

    def test_process_renders_steps_and_connectors(self):
        proc, out = self._render_block(
            "Block: process\nPlan\nCreate\nDeliver\n", "proc.pptx")
        self.assertEqual(proc.returncode, 0, proc.stderr)
        # 3 boxes + 3 numbers + 3 labels + 2 connectors = 11.
        self.assertEqual(len(self._drawn(out)), 11)

    def test_timeline_renders_dots_and_rail(self):
        proc, out = self._render_block(
            "Block: timeline\n2026 | Kickoff\n!2027 | Launch\n2028 | Scale\n",
            "tl.pptx")
        self.assertEqual(proc.returncode, 0, proc.stderr)
        # 3 dots + 3 labels + 2 rail segments = 8.
        self.assertEqual(len(self._drawn(out)), 8)

    def test_comparison_wrong_count_fails(self):
        proc, out = self._render_block("Block: comparison\nOnly one | x\n", "bad.pptx")
        self.assertNotEqual(proc.returncode, 0)
        self.assertIn("exactly two", proc.stderr + proc.stdout)

    def test_process_over_five_fails(self):
        block = "Block: process\n" + "".join(f"S{i}\n" for i in range(6))
        proc, out = self._render_block(block, "bad2.pptx")
        self.assertNotEqual(proc.returncode, 0)
        self.assertIn("at most 5", proc.stderr + proc.stdout)

    def test_two_blocks_stack(self):
        # Two unplaced blocks stack top to bottom on one slide, no overlap.
        proc, out = self._render_block(
            "Block: stat-row\n56 | Days\n4% | Rate\nBlock: process\nPlan\nShip\n",
            "stack.pptx")
        self.assertEqual(proc.returncode, 0, proc.stderr)
        drawn = self._drawn(out)
        # 2 numbers + 2 labels (stat-row) + 2 boxes + 2 numbers + 2 labels + 1
        # connector (process) = 11 shapes, all on one slide.
        self.assertEqual(len(drawn), 11)

    def test_placed_blocks_tile_side_by_side(self):
        # Left half card-grid, right half process — placed on the grid.
        proc, out = self._render_block(
            "Block: card-grid at cols 1-6\nWhy | reason\nWho | people\n"
            "Block: process at cols 7-12\nPlan\nShip\n",
            "tile.pptx")
        self.assertEqual(proc.returncode, 0, proc.stderr)
        prs = Presentation(out)
        slide = prs.slides[0]
        drawn = [s for s in slide.shapes if not s.is_placeholder]
        mid = prs.slide_width // 2
        # Card panels sit left of centre; process boxes sit right of it.
        left_boxes = [s for s in drawn if s.left + s.width <= mid + 10000]
        right_boxes = [s for s in drawn if s.left >= mid - 10000]
        self.assertTrue(left_boxes, "no shapes in the left column")
        self.assertTrue(right_boxes, "no shapes in the right column")

    def test_freeform_renders_node_graph(self):
        # The case the fixed primitives can't express: nodes that connect.
        block = (
            "Block: freeform\n"
            "panel paper outline ink at cols 1-4 rows 1-8\n"
            "text h1 ink at cols 1-4 rows 1-3 | UK\n"
            "arrow ink at cols 5-5 rows 4-5\n"
            "panel accent at cols 6-9 rows 1-8\n"
            "text h1 paper at cols 6-9 rows 1-3 | Platform\n"
        )
        proc, out = self._render_block(block, "ff.pptx")
        self.assertEqual(proc.returncode, 0, proc.stderr)
        self.assertEqual(len(self._drawn(out)), 5)

    def test_freeform_bad_colour_fails(self):
        proc, out = self._render_block(
            "Block: freeform\ntext h1 purple at cols 1-6 rows 1-3 | Hi\n",
            "ffb.pptx")
        self.assertNotEqual(proc.returncode, 0)
        self.assertIn("purple", proc.stderr + proc.stdout)

    def test_freeform_missing_at_fails(self):
        proc, out = self._render_block(
            "Block: freeform\ntext h1 ink hello world\n", "ffb2.pptx")
        self.assertNotEqual(proc.returncode, 0)
        self.assertIn("at <placement>", proc.stderr + proc.stdout)

    def test_freeform_overlap_still_fails_hard(self):
        # The guardrail holds: freedom of arrangement, but no overlap.
        block = (
            "Block: freeform\n"
            "panel paper outline ink at cols 1-8 rows 1-8\n"
            "panel accent at cols 5-12 rows 4-10\n"
        )
        proc, out = self._render_block(block, "ffo.pptx")
        self.assertNotEqual(proc.returncode, 0)
        self.assertIn("[overlap]", proc.stderr + proc.stdout)

    def test_freeform_icon_degrades_without_rasteriser(self):
        # No cairosvg here: the icon is dropped, the deck still builds, and the
        # summary names the fallback (mirrors the matplotlib chart fallback).
        import icons
        block = ("Block: freeform\n"
                 "icon growth accent at cols 1-3 rows 1-4\n"
                 "text h1 ink at cols 4-10 rows 1-4 | Up\n")
        proc, out = self._render_block(block, "ffi.pptx")
        self.assertEqual(proc.returncode, 0, proc.stderr)
        self.assertTrue(os.path.isfile(out))
        if not icons.cairosvg_available():
            self.assertIn("cairosvg not installed", proc.stdout + proc.stderr)
            # the text still drew; the icon was dropped
            drawn = self._drawn(out)
            self.assertEqual(len(drawn), 1)

    def test_tree_renders_nodes_and_edges(self):
        block = ("Block: tree\n"
                 "Product\n"
                 "  Discovery\n"
                 "  Delivery\n"
                 "  !Growth\n")
        proc, out = self._render_block(block, "tree.pptx")
        self.assertEqual(proc.returncode, 0, proc.stderr)
        # 4 nodes + 4 labels + 3 edges = 11 drawn shapes.
        self.assertEqual(len(self._drawn(out)), 11)

    def test_tree_indent_jump_fails(self):
        block = "Block: tree\nRoot\n    Grandchild\n"  # jumps root -> level 2
        proc, out = self._render_block(block, "treebad.pptx")
        self.assertNotEqual(proc.returncode, 0)
        self.assertIn("jumps a level", proc.stderr + proc.stdout)

    def test_tree_two_roots_fails(self):
        block = "Block: tree\nRootA\nRootB\n"
        proc, out = self._render_block(block, "tree2.pptx")
        self.assertNotEqual(proc.returncode, 0)
        self.assertIn("exactly one root", proc.stderr + proc.stdout)

    def test_icon_list_renders(self):
        import icons
        block = ("Block: icon-list\n"
                 "growth | Revenue up\n"
                 "team | Team scaled\n"
                 "fast | Shipping weekly\n")
        proc, out = self._render_block(block, "il.pptx")
        self.assertEqual(proc.returncode, 0, proc.stderr)
        drawn = self._drawn(out)
        # 3 text rows always draw; the 3 icons draw only with a rasteriser.
        expected = 6 if icons.cairosvg_available() else 3
        self.assertEqual(len(drawn), expected)

    def test_icon_list_bad_icon_fails(self):
        proc, out = self._render_block(
            "Block: icon-list\nno-such | x\n", "ilbad.pptx")
        self.assertNotEqual(proc.returncode, 0)
        self.assertIn("no-such", proc.stderr + proc.stdout)

    def test_card_icon_prefix_renders(self):
        block = ("Block: card-grid\n"
                 "[growth] Grow | markets\n"
                 "[team] Serve | users\n")
        proc, out = self._render_block(block, "cardicon.pptx")
        self.assertEqual(proc.returncode, 0, proc.stderr)

    def test_cycle_renders(self):
        block = "Block: cycle\nPlan\nBuild\nMeasure\nLearn\n"
        proc, out = self._render_block(block, "cycle.pptx")
        self.assertEqual(proc.returncode, 0, proc.stderr)
        # 4 nodes + 4 labels + 4 ring edges = 12 shapes.
        self.assertEqual(len(self._drawn(out)), 12)

    def test_matrix_renders(self):
        block = ("Block: matrix\nx: Effort\ny: Impact\n"
                 "Quick wins | do now\n!Big bets | plan\n"
                 "Deprioritise | skip\nFill-ins | maybe\n")
        proc, out = self._render_block(block, "matrix.pptx")
        self.assertEqual(proc.returncode, 0, proc.stderr)
        # 4 cells + 4 labels + 2 axis captions (+ up to 4 bodies if the region
        # is tall enough — bodies are clamped to the cell under a title).
        self.assertGreaterEqual(len(self._drawn(out)), 10)

    def test_matrix_wrong_count_fails(self):
        block = "Block: matrix\nOnly | one\nTwo | two\n"
        proc, out = self._render_block(block, "matbad.pptx")
        self.assertNotEqual(proc.returncode, 0)
        self.assertIn("exactly four", proc.stderr + proc.stdout)


class TableRenderTest(unittest.TestCase):
    """T-011 (REQ-001/002/004): a `Block: table` composed slide renders a native
    pptx GraphicFrame table, styled entirely from brand tokens; the hard caps fail
    loudly. The render path is exercised via the render.py subprocess and read back
    with python-pptx, matching the end-to-end style of the classes above.
    """

    # Header `Metric | Q1 | Q2`, three data rows, one `!` emphasis row (Total).
    TABLE_SPEC = (
        "---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
        "Title: Results\n"
        "Block: table\n"
        "Metric | Q1 | Q2\n"
        "Revenue | $1.2M | $1.5M\n"
        "Users | 4% | 9%\n"
        "!Total | $2.7M | $3.0M\n"
    )

    @classmethod
    def setUpClass(cls):
        cls.proc, cls.out = _run(cls.TABLE_SPEC, out_name="table.pptx")

    def _table(self):
        """Reopen the rendered deck; assert exactly one GraphicFrame table and
        return (presentation, table)."""
        prs = Presentation(self.out)
        slide = prs.slides[0]
        tables = [s for s in slide.shapes
                  if s.shape_type == MSO_SHAPE_TYPE.TABLE]
        self.assertEqual(len(tables), 1,
                         "expected exactly one GraphicFrame table on the slide")
        return prs, tables[0].table

    def test_table_slide_renders(self):
        self.assertEqual(self.proc.returncode, 0,
                         f"render failed: {self.proc.stderr}\n{self.proc.stdout}")
        self.assertTrue(os.path.isfile(self.out), "no .pptx written")

    def test_single_graphicframe_with_expected_cells(self):
        # One native table with the header and cells exactly as specified; the
        # `!` emphasis grammar does not alter cell text, only styling.
        _prs, tbl = self._table()
        self.assertEqual(len(tbl.rows), 4)      # header + 3 data rows
        self.assertEqual(len(tbl.columns), 3)
        header = [tbl.cell(0, c).text for c in range(3)]
        self.assertEqual(header, ["Metric", "Q1", "Q2"])
        data = [[tbl.cell(r, c).text for c in range(3)] for r in range(1, 4)]
        self.assertEqual(data, [
            ["Revenue", "$1.2M", "$1.5M"],
            ["Users", "4%", "9%"],
            ["Total", "$2.7M", "$3.0M"],
        ])

    def test_every_cell_fill_is_a_token_colour(self):
        # D-003: theme banding is stripped, so every cell is an explicit SOLID
        # token fill from the brand palette (ink header / paper data rows / accent
        # emphasis row) — never the pptx default blue banding.
        sys.path.insert(0, SCRIPTS)
        import tokens
        from pptx.enum.dml import MSO_FILL
        palette = {v.lstrip("#").upper()
                   for v in tokens.resolve_colour_roles(BRAND["colours"]).values()}
        _prs, tbl = self._table()
        for r in range(len(tbl.rows)):
            for c in range(len(tbl.columns)):
                cell = tbl.cell(r, c)
                self.assertEqual(
                    cell.fill.type, MSO_FILL.SOLID,
                    f"cell ({r},{c}) is not a solid fill (theme banding leaked)")
                rgb = str(cell.fill.fore_color.rgb).upper()
                self.assertIn(rgb, palette,
                              f"cell ({r},{c}) fill {rgb} is off-palette {palette}")

    def test_nine_data_rows_fails_with_row_cap(self):
        # D-004 (pinned in tests/test_primitives.py TestPlanTable): the cap allows
        # 8 data rows and rejects 9. A full band (no Title) isolates the row-count
        # cap as the failure mode, distinct from the band-fit guard below.
        spec = (
            "---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
            "Block: table\nMetric | Q1 | Q2\n"
            + "".join(f"r{i} | {i} | {i}\n" for i in range(9))
        )
        proc, out = _run(spec, out_name="ninerow.pptx")
        self.assertNotEqual(proc.returncode, 0)
        combined = proc.stderr + proc.stdout
        self.assertIn("error:", combined)
        self.assertIn("row", combined.lower())            # names the row cap
        # The row-count cap message, not the band-fit guard ("band fits only").
        self.assertIn("at most 8 data rows", combined)
        self.assertFalse(os.path.isfile(out),
                         "a half-built .pptx was written")

    def test_eight_data_rows_under_title_hits_band_fit_guard(self):
        # A separate, legitimate failure mode (D-004 band-fit guard): under a
        # Title the fixture's band is short, so even 8 data rows overflow it and
        # the guard fires naming how many rows fit — distinct from the row cap.
        spec = (
            "---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
            "Title: Big\nBlock: table\nMetric | Q1 | Q2\n"
            + "".join(f"r{i} | {i} | {i}\n" for i in range(8))
        )
        proc, out = _run(spec, out_name="eightrow.pptx")
        self.assertNotEqual(proc.returncode, 0)
        combined = proc.stderr + proc.stdout
        self.assertIn("error:", combined)
        self.assertIn("row", combined.lower())
        self.assertIn("band fits only", combined)
        self.assertFalse(os.path.isfile(out))

    def test_table_beside_stat_row_places_and_lints_clean(self):
        # Region placement for tables: a table `at left` beside a stat-row `at
        # right` renders exit 0 (lint clean — the render gate blocks on any lint
        # violation) with the table in the left half and the stats in the right.
        spec = (
            "---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
            "Title: Split\n"
            "Block: table at left\nMetric | Q1 | Q2\n"
            "Revenue | $1.2M | $1.5M\n!Total | $2.7M | $3.0M\n"
            "Block: stat-row at right\n56 | Days\n4% | Rate\n"
        )
        proc, out = _run(spec, out_name="tablesplit.pptx")
        self.assertEqual(proc.returncode, 0,
                         f"placed table failed: {proc.stderr}\n{proc.stdout}")
        self.assertTrue(os.path.isfile(out), "no .pptx written")
        self.assertNotIn("error:", proc.stderr + proc.stdout)
        prs = Presentation(out)
        slide = prs.slides[0]
        mid = prs.slide_width // 2
        tables = [s for s in slide.shapes
                  if s.shape_type == MSO_SHAPE_TYPE.TABLE]
        self.assertEqual(len(tables), 1)
        tbl_shape = tables[0]
        self.assertLessEqual(tbl_shape.left + tbl_shape.width, mid + 10000,
                             "table should sit in the left half")
        stats = [s for s in slide.shapes
                 if not s.is_placeholder
                 and s.shape_type != MSO_SHAPE_TYPE.TABLE]
        self.assertTrue(stats, "no stat-row shapes drawn")
        self.assertTrue(all(s.left >= mid - 10000 for s in stats),
                        "stat-row should sit in the right half")


class ComposedChartTest(unittest.TestCase):
    """REQ-006 / D-006: `Block: chart` on a composed slide reuses
    `_parse_chart_block`'s `key: value` grammar (D-006) and renders into the
    block's region — a picture via `_place_picture` by default, or a native
    GraphicFrame chart with `native: true`. The chart element is planned
    inline, linted like any other element, then removed before `primitives.draw`
    and fulfilled afterwards — it must never reach `_add_text`.

    Red-first: today 'chart' is not a member of COMPOSED_BLOCK_TYPES, so every
    case below fails at parse with "unknown composed block type 'chart'"
    (exit 1, no .pptx written). Each assertion below pins the END-STATE
    (post-implementation) behaviour, so these are red now and green once
    render.py/native_charts.py land (T-006/T-007).
    """

    # The same `key: value` grammar `Chart:` already parses via
    # _parse_chart_block (a category chart: type/categories/series).
    CHART_ITEMS = (
        "type: column\n"
        "categories: Q1, Q2, Q3\n"
        "series Rev: 10, 20, 30\n"
    )

    @classmethod
    def setUpClass(cls):
        spec = (
            "---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
            "Title: Chart slide\n"
            "Block: chart\n" + cls.CHART_ITEMS
        )
        cls.proc, cls.out = _run(spec, out_name="composed-chart.pptx")

    def _grid(self):
        """The fixture's derived grid tokens (margin_x/margin_top/margin_bottom),
        the same content-band margins ComposedRenderTest measures via MARGIN_X."""
        sys.path.insert(0, SCRIPTS)
        import tokens
        prs = Presentation(TEMPLATE)
        return tokens.resolve_tokens(BRAND, prs)["grid"]

    def test_chart_block_renders_as_picture(self):
        # Image path (native: absent): exit 0, exactly one picture, sitting
        # inside the content band (matplotlib is installed dev-side, so the
        # PNG draws for real rather than degrading to a note).
        self.assertEqual(self.proc.returncode, 0,
                         f"render failed: {self.proc.stderr}\n{self.proc.stdout}")
        prs = Presentation(self.out)
        slide = prs.slides[0]
        pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        self.assertEqual(len(pics), 1, "expected exactly one picture (the chart)")
        pic = pics[0]
        grid = self._grid()
        self.assertGreaterEqual(pic.left, grid["margin_x"])
        self.assertGreaterEqual(pic.top, grid["margin_top"])
        self.assertLessEqual(pic.left + pic.width,
                              prs.slide_width - grid["margin_x"])
        self.assertLessEqual(pic.top + pic.height,
                              prs.slide_height - grid["margin_bottom"])

    def test_chart_element_is_not_drawn_as_text(self):
        # D-006 guard: the composed render plans one inline element
        # ({"kind": "chart", "text": "chart: <type>", ...}) purely so lint can
        # see a bbox; it must be removed before primitives.draw and fulfilled
        # afterwards. It must never reach `_add_text` and leave a stray
        # textbox literally reading "chart: column".
        self.assertEqual(self.proc.returncode, 0,
                         f"render failed: {self.proc.stderr}\n{self.proc.stdout}")
        prs = Presentation(self.out)
        slide = prs.slides[0]
        for shape in slide.shapes:
            if shape.has_text_frame:
                self.assertNotIn(
                    "chart:", shape.text_frame.text.lower(),
                    "the chart plan element leaked into draw() as a textbox")

    def test_chart_block_native_renders_graphicframe(self):
        # `native: true` on the same block: exit 0, and the shape drawn for
        # the chart is a real GraphicFrame with an editable chart part
        # (REQ-001), not a picture.
        spec = (
            "---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
            "Title: Native chart slide\n"
            "Block: chart\n" + self.CHART_ITEMS + "native: true\n"
        )
        proc, out = _run(spec, out_name="composed-chart-native.pptx")
        self.assertEqual(proc.returncode, 0,
                         f"render failed: {proc.stderr}\n{proc.stdout}")
        prs = Presentation(out)
        slide = prs.slides[0]
        chart_frames = [s for s in slide.shapes if s.has_chart]
        self.assertEqual(len(chart_frames), 1,
                         "expected exactly one native GraphicFrame chart")
        self.assertTrue(chart_frames[0].has_chart)

    def test_chart_block_placed_region(self):
        # `Block: chart at left` beside `Block: stat-row at right`: the
        # dashboard case (D-006's "who uses this" — a stat-row up top, chart
        # beside it) — both regions tile the grid, lint stays clean, and each
        # block's shapes stay on its own side of the slide midline.
        spec = (
            "---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
            "Title: Split\n"
            "Block: chart at left\n" + self.CHART_ITEMS +
            "Block: stat-row at right\n56 | Days\n4% | Rate\n"
        )
        proc, out = _run(spec, out_name="composed-chart-split.pptx")
        self.assertEqual(proc.returncode, 0,
                         f"render failed: {proc.stderr}\n{proc.stdout}")
        prs = Presentation(out)
        slide = prs.slides[0]
        mid = prs.slide_width // 2
        pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        self.assertEqual(len(pics), 1, "expected exactly one picture (the chart)")
        pic = pics[0]
        self.assertLessEqual(pic.left + pic.width, mid + 10000,
                             "chart picture should sit in the left half")
        stats = [s for s in slide.shapes
                 if not s.is_placeholder
                 and s.shape_type != MSO_SHAPE_TYPE.PICTURE]
        self.assertTrue(stats, "no stat-row shapes drawn")
        self.assertTrue(all(s.left >= mid - 10000 for s in stats),
                        "stat-row should sit in the right half")

    def test_chart_block_missing_type_fails(self):
        # Failure grammar: `Block: chart` reuses `_parse_chart_block`'s own
        # "no 'type'" error rather than inventing a new one. Asserting only
        # the substring "type" would pass vacuously today too (the current
        # "unknown composed block type 'chart'" message also contains the
        # word "type"), so the negative assertion below pins that this is
        # genuinely `_parse_chart_block`'s error, not the generic
        # unknown-block-type one — genuinely red until 'chart' is a
        # recognised composed block type.
        spec = (
            "---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
            "Block: chart\ncategories: Q1, Q2, Q3\nseries Rev: 10, 20, 30\n"
        )
        proc, out = _run(spec, out_name="composed-chart-no-type.pptx")
        self.assertNotEqual(proc.returncode, 0)
        combined = proc.stderr + proc.stdout
        self.assertIn("error:", combined)
        self.assertIn("type", combined)
        self.assertNotIn("unknown composed block type", combined)
        self.assertFalse(os.path.isfile(out), "a half-built .pptx was written")


class ComposedMatplotlibAbsentNoteTest(unittest.TestCase):
    """Architect review FINDING 1 (feat/native-charts): when matplotlib is
    absent, composed `Block: chart` fulfilment collected `chart_notes` as a
    Python list and handed the list itself to `_apply_meta`'s `extra_visual`,
    which f-strung it straight into the note -> `VISUAL TO ADD:
    ['Column chart. ...']`, leaking the list's repr (brackets and quotes)
    instead of clean prose. `_apply_meta` now emits one clean
    'VISUAL TO ADD:' line per chart note.

    matplotlib absence is simulated exactly as test_render.py's
    ChartRenderTest.test_matplotlib_absent_falls_back_to_note does: a
    PYTHONPATH shim directory holding a `matplotlib.py` that raises
    ImportError on import, prepended to the subprocess's PYTHONPATH — no
    existing absence-simulation helper is shared between the two test files,
    so this mirrors that one rather than inventing a new mechanism.
    """

    @classmethod
    def setUpClass(cls):
        cls._tmp = tempfile.mkdtemp(prefix="slides-composed-nomatplotlib-")
        shim = os.path.join(cls._tmp, "shim")
        os.makedirs(shim, exist_ok=True)
        with open(os.path.join(shim, "matplotlib.py"), "w") as fh:
            fh.write('raise ImportError("matplotlib hidden for test")\n')
        cls.env = dict(os.environ)
        cls.env["PYTHONPATH"] = shim + os.pathsep + cls.env.get("PYTHONPATH", "")

    @classmethod
    def tearDownClass(cls):
        shutil.rmtree(cls._tmp, ignore_errors=True)

    def test_single_chart_note_has_no_repr_leak(self):
        spec = (
            "---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
            "Title: Chart slide\n"
            "Block: chart\ntype: column\ncategories: Q1, Q2, Q3\n"
            "series Rev: 10, 20, 30\n"
        )
        proc, out = _run(spec, out_name="nomatplotlib-single.pptx", env=self.env)
        self.assertEqual(proc.returncode, 0,
                         f"render failed: {proc.stderr}\n{proc.stdout}")
        prs = Presentation(out)
        notes = prs.slides[0].notes_slide.notes_text_frame.text
        self.assertIn("VISUAL TO ADD:", notes)
        self.assertNotIn("[", notes, f"list repr leaked into notes: {notes!r}")
        self.assertNotIn("]", notes, f"list repr leaked into notes: {notes!r}")
        self.assertNotIn("'", notes, f"list repr leaked into notes: {notes!r}")
        self.assertEqual(notes.count("VISUAL TO ADD:"), 1)
        after = notes.split("VISUAL TO ADD:", 1)[1].strip()
        self.assertTrue(
            after.startswith("Column chart"),
            f"expected clean prose right after the marker, got: {after!r}")

    def test_two_charts_yield_two_clean_notes(self):
        # A composed slide may carry several `Block: chart`; each must
        # contribute its own clean 'VISUAL TO ADD:' line, not one line
        # holding a Python-list rendering of all of them.
        spec = (
            "---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
            "Title: Two charts\n"
            "Block: chart at left\ntype: column\ncategories: Q1, Q2\n"
            "series Rev: 10, 20\n"
            "Block: chart at right\ntype: bar\ncategories: A, B\n"
            "series Cost: 5, 7\n"
        )
        proc, out = _run(spec, out_name="nomatplotlib-two.pptx", env=self.env)
        self.assertEqual(proc.returncode, 0,
                         f"render failed: {proc.stderr}\n{proc.stdout}")
        prs = Presentation(out)
        notes = prs.slides[0].notes_slide.notes_text_frame.text
        self.assertEqual(
            notes.count("VISUAL TO ADD:"), 2,
            f"expected one clean VISUAL TO ADD line per chart block: {notes!r}")
        self.assertNotIn("[", notes, f"list repr leaked into notes: {notes!r}")
        self.assertNotIn("]", notes, f"list repr leaked into notes: {notes!r}")
        self.assertNotIn("'", notes, f"list repr leaked into notes: {notes!r}")
        self.assertIn("Column chart", notes)
        self.assertIn("Bar chart", notes)


class ComposedChartFontTest(unittest.TestCase):
    """Architect review FINDING 2 (feat/native-charts): the composed
    `Block: chart` image-chart fulfilment called
    `charts_mod.render_png(chart, brand["colours"], None, png)` — font always
    hard-coded to None — while the title-content chart branch registers the
    brand font via `register_brand_font(brand, brand_path)` and passes the
    resolved family name. render.py now threads that same lazily-resolved
    `font_family`/`font_warning` state into `_render_composed_slide` (and
    back out, so it is only ever registered once per run) so composed image
    charts use the identical registered brand font.

    This is pinned at the seam rather than end-to-end: `_render_composed_slide`
    is called directly (in-process, scripts dir on sys.path — the same
    white-box style ComposedLintGateTest and TableRenderTest above already
    use for `lint.check`/`tokens.resolve_tokens`) against a real parsed
    composed spec (via `render.parse_spec`, so the chart dict shape is
    exactly what the real pipeline produces) and a real brand font file (the
    matplotlib-bundled DejaVuSans.ttf, the same fixture
    test_render.py's test_font_file_registered_no_warning reuses). Only the
    `charts` module is stubbed, to record the literal `font` argument
    `render_png` receives — nothing in this suite verifies that matplotlib
    then renders that family's glyphs correctly (not even for the
    title-content path: the existing font coverage only checks that the
    "fallback font" warning disappears), so recording the argument value at
    the exact call site is the most direct, least-mocked way to pin "composed
    passes the same registered brand font title-content does" without
    depending on font-rendering internals this suite doesn't otherwise test.
    """

    def setUp(self):
        sys.path.insert(0, SCRIPTS)
        self._saved_charts = sys.modules.get("charts")

    def tearDown(self):
        # Restore whatever (if anything) really occupied sys.modules['charts']
        # so this stub never leaks into other test modules sharing the process
        # (test_charts.py imports the real module at collection time).
        if self._saved_charts is not None:
            sys.modules["charts"] = self._saved_charts
        else:
            sys.modules.pop("charts", None)

    def test_composed_chart_receives_registered_brand_font(self):
        import types

        import matplotlib
        import render
        import tokens as tokens_mod
        from PIL import Image

        ttf = os.path.join(os.path.dirname(matplotlib.__file__),
                           "mpl-data", "fonts", "ttf", "DejaVuSans.ttf")
        self.assertTrue(os.path.isfile(ttf), "matplotlib DejaVuSans missing")

        recorded_fonts = []

        class _StubChartError(Exception):
            pass

        def _fake_render_png(chart, colours, font, png_path):
            recorded_fonts.append(font)
            # A real, minimal PNG so _place_picture's IHDR read succeeds —
            # its content is irrelevant, only the `font` arg is under test.
            Image.new("RGB", (10, 10), (255, 255, 255)).save(png_path, "PNG")

        stub = types.ModuleType("charts")
        stub.render_png = _fake_render_png
        stub.ChartError = _StubChartError
        sys.modules["charts"] = stub

        tmp = tempfile.mkdtemp()
        brand = dict(BRAND)
        brand["font_files"] = {"DejaVu Sans": ttf}
        brand_path = os.path.join(tmp, "brand.json")
        with open(brand_path, "w", encoding="utf-8") as fh:
            json.dump(brand, fh)

        spec_text = (
            "---\ndeck: d\naudience: a\n---\n\n## Slide 1\nlayout: composed\n"
            "Title: Chart slide\n"
            "Block: chart\ntype: column\ncategories: Q1, Q2\n"
            "series Rev: 10, 20\n"
        )
        spec_path = os.path.join(tmp, "deck.md")
        with open(spec_path, "w", encoding="utf-8") as fh:
            fh.write(spec_text)

        slides = render.parse_spec(spec_path)
        prs = Presentation(brand["template"])
        toks = tokens_mod.resolve_tokens(brand, prs)
        charts_dir = os.path.join(tmp, "charts")

        result = render._render_composed_slide(
            prs, brand, slides[0], toks, charts_dir,
            "unset", None, brand_path,
        )
        _advisories, _dropped_icons, _fallback_notes, font_family, _warning = (
            result
        )

        self.assertEqual(len(recorded_fonts), 1,
                         "expected exactly one render_png call")
        self.assertEqual(
            recorded_fonts[0], "DejaVu Sans",
            "composed image chart must receive the registered brand font, "
            "not None")
        self.assertEqual(font_family, "DejaVu Sans")
