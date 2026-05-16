#!/usr/bin/env python3
"""Generate tests/fixtures/sample-template.pptx — the renderer test fixture.

A "template" and a "deck" are the same thing to python-pptx: both carry a
slide master and a set of named slide layouts. This script produces a minimal
multi-layout template by starting from python-pptx's bundled default
presentation (which ships a full set of named layouts) and stripping any
slides, so the fixture is layouts-only — exactly what render.py expects.

The layouts are renamed to the semantic roles the deck-spec grammar uses, so
inspect_template.py output reads cleanly and the role-to-layout mapping in the
test brand profile is obvious. Run from the repo root:

    python3 tests/generate-fixture-template.py

Re-running always produces the same layout structure (deterministic), though
the .pptx is a zip and is not guaranteed byte-identical across runs.
"""
import os
from pptx import Presentation

# Maps a default-template layout index to the role name render.py resolves.
# The default python-pptx template ships these layouts at stable indices:
#   0 Title Slide | 1 Title and Content | 2 Section Header
#   3 Two Content | 5 Title Only
ROLE_NAMES = {
    0: "title",
    1: "title-content",
    2: "section",
    3: "two-column",
    5: "statement",
}

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(REPO_ROOT, "tests", "fixtures", "sample-template.pptx")


def main():
    prs = Presentation()  # bundled default template: master + named layouts

    # A fresh default Presentation() has no slides; guard anyway so the
    # fixture is always layouts-only even if python-pptx changes.
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        xml_slides.remove(sld)

    for idx, role in ROLE_NAMES.items():
        prs.slide_layouts[idx].name = role

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)

    layouts = Presentation(OUT).slide_layouts
    print(f"wrote {OUT}")
    print(f"{len(layouts)} layouts:")
    for i, layout in enumerate(layouts):
        phs = ", ".join(
            f"idx{ph.placeholder_format.idx}:{ph.placeholder_format.type}"
            for ph in layout.placeholders
        )
        print(f"  [{i}] {layout.name!r} — {phs or 'no placeholders'}")


if __name__ == "__main__":
    main()
