"""Tests for `deck_to_spec.py` — the render.py inverse (T-002 of
`.build/plans/round-trip-revise-plan.md`, REQ-002/003/007).

Pins the extraction contract before the script exists: CLI shape, exit codes,
round-trip fidelity for the fixed roles, `--against` drift detection, the
foreign-tier import with a mandatory report, and composed-slide flattening.

Stdlib `unittest` only — no pytest, no third-party test deps beyond
python-pptx (already a runtime dependency of the pipeline, imported here to
mutate a rendered .pptx for the drift test). Subprocess-driven, mirroring
tests/test_render.py and tests/test_composed.py: deck_to_spec.py is invoked
as a script, never imported. `render.py` is imported in-process only to call
`render.parse_spec`/`render.ROLE_FIELDS` for comparison, following the
sys.path header pattern tests/test_charts.py uses for the same purpose.

Run from the repo root:

    python3 -m unittest tests.test_deck_to_spec -v
    python3 -m unittest discover tests

Until deck_to_spec.py exists, every test in this file FAILS (not skips) at
setUpClass with a named "deck_to_spec.py not yet implemented" AssertionError
— see `_require_deck_to_spec` below. A self-skipping pattern is forbidden
here on purpose: a green run must mean the contract is met, never "there was
nothing to check yet".
"""
import json
import os
import shutil
import subprocess
import sys
import tempfile
import unittest

from pptx import Presentation

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIXTURES = os.path.join(REPO_ROOT, "tests", "fixtures")
TEMPLATE = os.path.join(FIXTURES, "sample-template.pptx")
SCRIPTS = os.path.join(REPO_ROOT, "source", "skills", "build-deck", "scripts")
RENDER_PY = os.path.join(SCRIPTS, "render.py")
DECK_TO_SPEC_PY = os.path.join(SCRIPTS, "deck_to_spec.py")

# tests/test_charts.py's sys.path header pattern: put the scripts dir on the
# path so render.parse_spec/render.ROLE_FIELDS can be used in-process to
# validate the emitted spec and to parse the original fixture spec for
# comparison. render.py already exists (this file does not pin its
# behaviour); only deck_to_spec.py, invoked below via subprocess, is new.
sys.path.insert(0, SCRIPTS)
import render  # noqa: E402

# Same fixture template + layout_map as test_render.py/test_composed.py:
# tests/fixtures/sample-template.pptx has its layouts renamed to role names by
# generate-fixture-template.py. `quote` has no dedicated layout and shares the
# `section` layout (TITLE + BODY) at index 2 — D-006's collision case (not
# exercised directly here: this fixture spec never declares a `section`
# slide, so populated-placeholder-count disambiguation always resolves quote).
LAYOUT_MAP = {
    "title": 0,
    "title-content": 1,
    "section": 2,
    "two-column": 3,
    "statement": 5,
    "quote": 2,
    "composed": 5,
}
BRAND = {
    "template": TEMPLATE,
    "fonts": {"heading": "Calibri", "body": "Calibri"},
    "colours": {
        "primary": "#1F3A5F",
        "accent": "#E07A3F",
        "ink": "#1A1A1A",
        "paper": "#FFFFFF",
    },
    "layout_map": LAYOUT_MAP,
}

# D-007: the extracted frontmatter's audience placeholder, pinned verbatim.
AUDIENCE_PLACEHOLDER = "(unknown — set before re-rendering)"

# One of each fixed role the round trip must prove: title (Title+Subtitle),
# title-content (Title + 3-bullet Body + Notes + Visual), statement, and
# quote (Quote+Attribution).
FIXTURE_SPEC_TEXT = (
    "---\n"
    "deck: The Quarterly Review\n"
    "audience: The leadership team, deciding whether to renew the pilot.\n"
    "register: presented\n"
    "---\n"
    "\n"
    "## Slide 1\n"
    "layout: title\n"
    "Title: The Quarter in One Page\n"
    "Subtitle: A short subtitle for the round trip\n"
    "\n"
    "## Slide 2\n"
    "layout: title-content\n"
    "Title: What We Learned\n"
    "Body:\n"
    "- First bullet about the learning journey\n"
    "- Second bullet naming a blocker\n"
    "- Third bullet naming next steps\n"
    "Notes: A note for the speaker to remember before this slide.\n"
    "Visual: A two-panel diagram showing before and after states.\n"
    "\n"
    "## Slide 3\n"
    "layout: statement\n"
    "Statement: This is the single big claim of the deck.\n"
    "\n"
    "## Slide 4\n"
    "layout: quote\n"
    "Quote: This tool changed how we plan every week.\n"
    "Attribution: A satisfied customer, in a review\n"
)

# A single composed slide (D-005): drawn shapes, not placeholders, so the
# extractor cannot reconstruct `Block:` grammar and must flatten it.
COMPOSED_SPEC_TEXT = (
    "---\n"
    "deck: Quarterly Numbers\n"
    "audience: The leadership team.\n"
    "register: presented\n"
    "---\n"
    "\n"
    "## Slide 1\n"
    "layout: composed\n"
    "Title: What moved this quarter\n"
    "Block: stat-row\n"
    "56 | Days to close\n"
    "4% | Churn rate\n"
)


def _require_deck_to_spec():
    """Fail loudly (never skip) when deck_to_spec.py does not exist yet.

    Called from every TestCase's setUpClass in this file. Raising here (a
    classmethod, so no `self`) fails every test in the class with a clear,
    named reason — 'deck_to_spec.py not yet implemented' — rather than
    letting each test hit a confusing subprocess FileNotFoundError/exit-code
    mismatch on its own.
    """
    if not os.path.isfile(DECK_TO_SPEC_PY):
        raise AssertionError(
            f"deck_to_spec.py not yet implemented at {DECK_TO_SPEC_PY} "
            "(T-005 of .build/plans/round-trip-revise-plan.md); every test "
            "in this file fails until it exists."
        )


def _write(path, text):
    with open(path, "w", encoding="utf-8") as fh:
        fh.write(text)


def _write_brand(path, brand=None):
    with open(path, "w", encoding="utf-8") as fh:
        json.dump(brand or BRAND, fh)


def _render(tmp, spec_text, name="deck"):
    """Write spec_text + brand.json into tmp, render via render.py subprocess.

    Returns (spec_path, brand_path, out_path). Raises RuntimeError (not an
    assertion) if render.py itself fails — that would be a fixture bug in
    this test file, a fault distinct from anything deck_to_spec.py is
    pinned to do, so it must never be mistaken for a deck_to_spec.py
    failure.
    """
    spec_path = os.path.join(tmp, f"{name}.deck.md")
    brand_path = os.path.join(tmp, f"{name}.brand.json")
    out_path = os.path.join(tmp, f"{name}.pptx")
    _write(spec_path, spec_text)
    _write_brand(brand_path)
    proc = subprocess.run(
        [sys.executable, RENDER_PY, "--spec", spec_path,
         "--brand", brand_path, "--out", out_path],
        capture_output=True, text=True,
    )
    if proc.returncode != 0:
        raise RuntimeError(
            "fixture render failed (render.py itself, not deck_to_spec.py "
            f"under test):\nstdout: {proc.stdout}\nstderr: {proc.stderr}"
        )
    return spec_path, brand_path, out_path


def _extract(args):
    """Run deck_to_spec.py <args>. Returns the CompletedProcess."""
    return subprocess.run(
        [sys.executable, DECK_TO_SPEC_PY] + list(args),
        capture_output=True, text=True,
    )


def _norm(value):
    """Whitespace-normalised text for a field value (string or list).

    Pinned comparison rule for round-trip/drift equality: `" ".join(str(x)
    .split())`, applied after flattening a list field (Body/Left/Right) to a
    single string by joining its items with a space.
    """
    if isinstance(value, (list, tuple)):
        value = " ".join(str(v) for v in value)
    return " ".join(str(value).split())


def _read_frontmatter(spec_path):
    """Minimal `key: value` frontmatter reader (stdlib only — no PyYAML
    dependency; mirrors render.py's own hand-rolled frontmatter handling,
    which only locates the `---` block rather than parsing its keys)."""
    with open(spec_path, encoding="utf-8") as fh:
        lines = fh.read().splitlines()
    i = 0
    while i < len(lines) and not lines[i].strip():
        i += 1
    if i >= len(lines) or lines[i].strip() != "---":
        raise ValueError(f"{spec_path} has no opening '---' frontmatter line")
    fm = {}
    i += 1
    while i < len(lines) and lines[i].strip() != "---":
        line = lines[i]
        if ":" in line:
            key, value = line.split(":", 1)
            fm[key.strip()] = value.strip()
        i += 1
    return fm


class RoundTripTest(unittest.TestCase):
    """REQ-002/007 — the load-bearing property: render -> extract -> parse
    round-trips every fixed role, its fields, Notes, and the Visual field
    restored from the rendered 'VISUAL TO ADD:' notes line."""

    @classmethod
    def setUpClass(cls):
        _require_deck_to_spec()
        cls._tmp = tempfile.mkdtemp(prefix="slides-deck-to-spec-roundtrip-")
        cls.spec_path, cls.brand_path, cls.out_path = _render(
            cls._tmp, FIXTURE_SPEC_TEXT)
        cls.extracted_path = os.path.join(cls._tmp, "extracted.deck.md")
        cls.proc = _extract([
            cls.out_path, "--brand", cls.brand_path,
            "--out", cls.extracted_path,
        ])

    @classmethod
    def tearDownClass(cls):
        shutil.rmtree(cls._tmp, ignore_errors=True)

    def test_extraction_exits_zero(self):
        self.assertEqual(
            self.proc.returncode, 0,
            f"deck_to_spec.py exited {self.proc.returncode}.\n"
            f"stdout: {self.proc.stdout}\nstderr: {self.proc.stderr}",
        )

    def test_emitted_spec_parses(self):
        # The extractor's own self-validation contract (Approach, layer 2):
        # never emit a spec render.parse_spec rejects.
        render.parse_spec(self.extracted_path)

    def test_roles_match_in_order(self):
        orig = render.parse_spec(self.spec_path)
        extracted = render.parse_spec(self.extracted_path)
        self.assertEqual(len(orig), len(extracted), "slide count drifted")
        self.assertEqual(
            [s["role"] for s in orig], [s["role"] for s in extracted],
            "roles-in-order drifted",
        )

    def test_fields_match_normalised(self):
        orig = render.parse_spec(self.spec_path)
        extracted = render.parse_spec(self.extracted_path)
        for o, e in zip(orig, extracted):
            role = o["role"]
            for field in render.ROLE_FIELDS[role]:
                with self.subTest(slide=o["number"], field=field):
                    o_val = o["fields"].get(field, "")
                    e_val = e["fields"].get(field, "")
                    self.assertEqual(
                        _norm(o_val), _norm(e_val),
                        f"slide {o['number']} ({role}) field {field!r} "
                        f"drifted: orig={o_val!r} extracted={e_val!r}",
                    )

    def test_notes_preserved(self):
        orig = render.parse_spec(self.spec_path)
        extracted = render.parse_spec(self.extracted_path)
        # Slide 2 (title-content) is the only slide carrying Notes.
        o_notes = orig[1]["meta"].get("Notes", "")
        e_notes = extracted[1]["meta"].get("Notes", "")
        self.assertTrue(o_notes, "fixture sanity: slide 2 must declare Notes")
        self.assertEqual(_norm(o_notes), _norm(e_notes))

    def test_visual_restored_from_visual_to_add_note(self):
        orig = render.parse_spec(self.spec_path)
        extracted = render.parse_spec(self.extracted_path)
        o_visual = orig[1]["meta"].get("Visual", "")
        e_visual = extracted[1]["meta"].get("Visual", "")
        self.assertTrue(e_visual, "Visual field was not restored on slide 2")
        self.assertEqual(_norm(o_visual), _norm(e_visual))

    def test_restored_notes_do_not_carry_the_visual_to_add_prefix(self):
        # render.py joins Notes and 'VISUAL TO ADD: ...' into one speaker-notes
        # text frame; the extractor must split them back apart, not leave the
        # marker sitting inside the restored Notes field.
        extracted = render.parse_spec(self.extracted_path)
        e_notes = extracted[1]["meta"].get("Notes", "")
        self.assertNotIn("VISUAL TO ADD", _norm(e_notes))


class FrontmatterTest(unittest.TestCase):
    """D-007: the extracted frontmatter's `deck` is non-empty, `audience` is
    the named placeholder sentence, and `register` is present (default
    'presented')."""

    @classmethod
    def setUpClass(cls):
        _require_deck_to_spec()
        cls._tmp = tempfile.mkdtemp(prefix="slides-deck-to-spec-frontmatter-")
        cls.spec_path, cls.brand_path, cls.out_path = _render(
            cls._tmp, FIXTURE_SPEC_TEXT)
        cls.extracted_path = os.path.join(cls._tmp, "extracted.deck.md")
        cls.proc = _extract([
            cls.out_path, "--brand", cls.brand_path,
            "--out", cls.extracted_path,
        ])

    @classmethod
    def tearDownClass(cls):
        shutil.rmtree(cls._tmp, ignore_errors=True)

    def test_deck_is_non_empty(self):
        self.assertEqual(
            self.proc.returncode, 0,
            f"stdout: {self.proc.stdout}\nstderr: {self.proc.stderr}",
        )
        fm = _read_frontmatter(self.extracted_path)
        self.assertIn("deck", fm)
        self.assertTrue(fm["deck"].strip(), "frontmatter 'deck' is empty")

    def test_audience_is_the_d007_placeholder(self):
        fm = _read_frontmatter(self.extracted_path)
        self.assertEqual(fm.get("audience"), AUDIENCE_PLACEHOLDER)

    def test_register_present_and_defaults_presented(self):
        fm = _read_frontmatter(self.extracted_path)
        self.assertIn("register", fm)
        self.assertEqual(fm["register"], "presented")


class AgainstTest(unittest.TestCase):
    """REQ-003/D-009: `--against` diffs the deck's extracted text against a
    spec. Identical content exits 0; a hand-edit is named on stdout (slide
    number, field, both texts) and exits 2 — the mechanical hand-edit
    detector."""

    @classmethod
    def setUpClass(cls):
        _require_deck_to_spec()

    def setUp(self):
        self._tmp = tempfile.mkdtemp(prefix="slides-deck-to-spec-against-")
        self.addCleanup(shutil.rmtree, self._tmp, ignore_errors=True)
        self.spec_path, self.brand_path, self.out_path = _render(
            self._tmp, FIXTURE_SPEC_TEXT)

    def test_identical_content_exits_zero(self):
        proc = _extract([
            self.out_path, "--brand", self.brand_path,
            "--against", self.spec_path,
        ])
        self.assertEqual(
            proc.returncode, 0,
            f"expected no drift for an unedited render.\n"
            f"stdout: {proc.stdout}\nstderr: {proc.stderr}",
        )

    def test_hand_edited_title_is_named_and_exits_two(self):
        prs = Presentation(self.out_path)
        original_title = prs.slides[0].shapes.title.text
        mutated_title = "A Totally Different Opening Line"
        prs.slides[0].shapes.title.text = mutated_title
        prs.save(self.out_path)

        proc = _extract([
            self.out_path, "--brand", self.brand_path,
            "--against", self.spec_path,
        ])
        self.assertEqual(
            proc.returncode, 2,
            f"expected drift exit 2.\nstdout: {proc.stdout}\n"
            f"stderr: {proc.stderr}",
        )
        stdout_low = proc.stdout.lower()
        self.assertIn("slide 1", stdout_low)
        self.assertIn("title", stdout_low)
        self.assertIn(original_title.lower(), stdout_low)
        self.assertIn(mutated_title.lower(), stdout_low)


class ForeignTest(unittest.TestCase):
    """REQ-002 foreign tier: a deck with no lineage stamp and no --brand to
    invert (D-002's 'no usable layout_map match') still extracts a spec
    render.parse_spec accepts, via heuristic role inference, plus a
    mandatory import report at an explicit --report path."""

    @classmethod
    def setUpClass(cls):
        _require_deck_to_spec()
        cls._tmp = tempfile.mkdtemp(prefix="slides-deck-to-spec-foreign-")
        cls.spec_path, cls.brand_path, cls.out_path = _render(
            cls._tmp, FIXTURE_SPEC_TEXT)
        # Simulate a deck with no lineage stamp: blank the comments field
        # (T-004's stamp, or any pre-existing template comments) so the
        # foreign path is genuinely forced, present or future.
        prs = Presentation(cls.out_path)
        prs.core_properties.comments = ""
        prs.save(cls.out_path)

        cls.extracted_path = os.path.join(cls._tmp, "foreign.deck.md")
        cls.report_path = os.path.join(cls._tmp, "foreign.import-report.md")
        # Deliberately no --brand: nothing to invert layout_map against,
        # forcing the heuristic foreign-tier path end to end.
        cls.proc = _extract([
            cls.out_path, "--out", cls.extracted_path,
            "--report", cls.report_path,
        ])

    @classmethod
    def tearDownClass(cls):
        shutil.rmtree(cls._tmp, ignore_errors=True)

    def test_extraction_exits_zero(self):
        self.assertEqual(
            self.proc.returncode, 0,
            f"stdout: {self.proc.stdout}\nstderr: {self.proc.stderr}",
        )

    def test_emitted_spec_parses(self):
        render.parse_spec(self.extracted_path)

    def test_report_written_at_the_given_report_path(self):
        self.assertTrue(
            os.path.isfile(self.report_path),
            f"no import report written at --report path {self.report_path}",
        )

    def test_report_mentions_role_inference(self):
        with open(self.report_path, encoding="utf-8") as fh:
            text = fh.read().lower()
        self.assertIn("role", text)
        self.assertIn("infer", text)


class ComposedFlattenTest(unittest.TestCase):
    """D-005: a composed slide's drawn shapes cannot be reconstructed into
    `Block:` grammar. Extraction flattens it to a title-content slide, and
    the import report names the slide and the flattening."""

    @classmethod
    def setUpClass(cls):
        _require_deck_to_spec()
        cls._tmp = tempfile.mkdtemp(prefix="slides-deck-to-spec-composed-")
        cls.spec_path, cls.brand_path, cls.out_path = _render(
            cls._tmp, COMPOSED_SPEC_TEXT, name="composed")
        cls.extracted_path = os.path.join(
            cls._tmp, "composed-extracted.deck.md")
        cls.report_path = os.path.join(
            cls._tmp, "composed.import-report.md")
        cls.proc = _extract([
            cls.out_path, "--brand", cls.brand_path,
            "--out", cls.extracted_path, "--report", cls.report_path,
        ])

    @classmethod
    def tearDownClass(cls):
        shutil.rmtree(cls._tmp, ignore_errors=True)

    def test_extraction_exits_zero(self):
        self.assertEqual(
            self.proc.returncode, 0,
            f"stdout: {self.proc.stdout}\nstderr: {self.proc.stderr}",
        )

    def test_emitted_spec_parses(self):
        render.parse_spec(self.extracted_path)

    def test_composed_slide_flattens_to_title_content(self):
        extracted = render.parse_spec(self.extracted_path)
        self.assertEqual(len(extracted), 1)
        self.assertEqual(extracted[0]["role"], "title-content")

    def test_report_names_the_slide_and_the_flattening(self):
        with open(self.report_path, encoding="utf-8") as fh:
            text = fh.read().lower()
        self.assertIn("slide 1", text)
        self.assertIn("flat", text)  # flatten / flattened / flattening


class ErrorTest(unittest.TestCase):
    """Negative-path coverage matching render.py's `error: <reason>` style
    (CLI contract: exit 1 + 'error:' on unreadable input) plus the plan's
    .potx abuse case (Abuse and edge cases: rejected naming template and
    teach-slides, since a template is teach-slides' job, not revise's)."""

    @classmethod
    def setUpClass(cls):
        _require_deck_to_spec()

    def setUp(self):
        self._tmp = tempfile.mkdtemp(prefix="slides-deck-to-spec-error-")
        self.addCleanup(shutil.rmtree, self._tmp, ignore_errors=True)

    def test_nonexistent_path_exits_one_named_error(self):
        missing = os.path.join(self._tmp, "does-not-exist.pptx")
        proc = _extract([missing])
        self.assertEqual(
            proc.returncode, 1,
            f"stdout: {proc.stdout}\nstderr: {proc.stderr}",
        )
        msg = (proc.stdout + proc.stderr).lower()
        self.assertIn("error:", msg)

    def test_potx_input_is_rejected_naming_template_and_teach_slides(self):
        potx_path = os.path.join(self._tmp, "brand-template.potx")
        # Content is irrelevant: the .potx rejection is a file-kind check on
        # the path (plan Abuse and edge cases), not a content check.
        with open(potx_path, "wb") as fh:
            fh.write(b"")
        proc = _extract([potx_path])
        self.assertEqual(
            proc.returncode, 1,
            f"stdout: {proc.stdout}\nstderr: {proc.stderr}",
        )
        msg = (proc.stdout + proc.stderr).lower()
        self.assertIn("error:", msg)
        self.assertIn("template", msg)
        self.assertIn("teach-slides", msg)


if __name__ == "__main__":
    unittest.main()
