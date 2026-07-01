"""Tests for native_charts.py — the pptx.chart backend for `native: true` charts
(T-003, REQ-001/002/003; native-charts-plan.md D-002/D-003/D-004/D-005).

Pins the contract T-007 must satisfy:

    native_charts.supported(chart) -> (bool, reason_or_None)
    native_charts.insert(slide, chart, colours, region) -> graphic_frame

`chart` is the dict render.py's `_parse_chart_block` produces: for bar/column/
pie {type, categories, series: [{name, values}], emphasis, callout, fmt} plus
the new `native`/`stacked`/`target` keys (T-001); for line/scatter {type,
points, markers, callout, fmt} plus `native`/`target` (no `stacked`, no
`emphasis` — REQ-004/REQ-005 restrict those to bar/column). `colours` is
brand.json's colour dict (charts.py's `_resolve_colours` rule: only the
caller's brand colours are ever used). `region` is (left, top, width, height)
in EMU ints — the same 4-tuple shape render.py's `_chart_region` already
computes for the matplotlib picture path (D-002: a GraphicFrame takes the same
box a picture did).

Reasons pinned verbatim from D-005 (the render-summary fallback-note wording):
  - waterfall:            "waterfall has no native PowerPoint form"
  - target:/callout:/marker: "target:/callout:/marker: are drawn annotations"
Both are asserted for exact equality below (stronger than a substring check;
each also satisfies "contains 'native'" / "contains 'target'" respectively).

All tests below MUST fail red right now: `native_charts.py` does not exist
yet (T-007 creates it). `import native_charts` sits at module top (after the
sys.path insert, mirroring tests/test_charts.py's header), so the whole
module fails to load with one clear ModuleNotFoundError — a single,
unambiguous red, not N different per-test errors.

Run:
    python3 -m unittest tests.test_native_charts -v
"""
import os
import sys
import unittest

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SCRIPTS = os.path.join(REPO_ROOT, "source", "skills", "build-deck", "scripts")
FIXTURES = os.path.join(REPO_ROOT, "tests", "fixtures")
TEMPLATE = os.path.join(FIXTURES, "sample-template.pptx")
sys.path.insert(0, SCRIPTS)

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION  # noqa: E402

import native_charts  # noqa: E402 — module under test; does not exist yet (red)

# The BRAND colours dict shape (tests/test_charts.py's BRAND["colours"]), with
# explicit 'muted'/'spend' added so point-fill assertions below pin an exact
# expected RGB rather than depending on charts.py's private fallback constants
# (_MUTED_FALLBACK etc.) that native_charts.py may or may not reuse verbatim.
COLOURS = {
    "accent": "#4F81BD",
    "accent2": "#C0504D",
    "muted": "#BFBFBF",
    "spend": "#C0504D",
    "ink": "#000000",
    "paper": "#FFFFFF",
}

ACCENT_RGB = RGBColor.from_string(COLOURS["accent"].lstrip("#"))
MUTED_RGB = RGBColor.from_string(COLOURS["muted"].lstrip("#"))

# A generous EMU region well inside the fixture template's slide bounds
# (9144000 x 6858000) — the same 4-tuple shape render.py's _chart_region
# returns for the matplotlib picture path.
REGION = (914400, 914400, 6096000, 3429000)  # left, top, width, height

# D-005's exact fallback-note wording, pinned verbatim.
REASON_WATERFALL = "waterfall has no native PowerPoint form"
REASON_ANNOTATION = "target:/callout:/marker: are drawn annotations"


def _blank_slide():
    """A fresh slide on the fixture template's blank layout — mirrors
    tests/test_composed.py's fixture-template usage, but adds the slide
    directly (no render.py subprocess) since these tests exercise
    native_charts.insert in isolation, like a unit test of a new module."""
    prs = Presentation(TEMPLATE)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # "Blank"
    return prs, slide


def _bar_chart(ctype="column", **overrides):
    """A minimal bar/column chart dict (single series, no emphasis/stacked/
    target/callout) with sane defaults, overridable per test."""
    chart = {
        "type": ctype,
        "categories": ["Q1", "Q2", "Q3"],
        "series": [{"name": "Revenue", "values": [10.0, 20.0, 30.0]}],
        "emphasis": None,
        "callout": None,
        "fmt": {},
        "native": True,
        "stacked": False,
        "target": None,
    }
    chart.update(overrides)
    return chart


def _pie_chart(**overrides):
    chart = {
        "type": "pie",
        "categories": ["A", "B", "C"],
        "series": [{"name": "Share", "values": [50.0, 30.0, 20.0]}],
        "emphasis": None,
        "callout": None,
        "fmt": {},
        "native": True,
        "target": None,
    }
    chart.update(overrides)
    return chart


def _waterfall_chart(**overrides):
    chart = {
        "type": "waterfall",
        "categories": ["Start", "Q1", "Q2"],
        "series": [{"name": "Cash", "values": [40.0, -15.0, 25.0]}],
        "callout": None,
        "fmt": {},
        "total_label": "Total",
        "native": True,
    }
    chart.update(overrides)
    return chart


def _point_chart(ctype="line", **overrides):
    """A minimal line/scatter chart dict (XyChartData grammar: numeric x-y
    points, no categories/series/emphasis/stacked — REQ-004/D-003)."""
    chart = {
        "type": ctype,
        "points": [(0.0, 10.0), (1.0, 20.0), (2.0, 15.0)],
        "markers": [],
        "callout": None,
        "fmt": {},
        "native": True,
        "target": None,
    }
    chart.update(overrides)
    return chart


class TestSupported(unittest.TestCase):
    """native_charts.supported(chart) -> (bool, reason_or_None) truth table
    (D-005, REQ-003). Supported: plain bar/column/pie/line/scatter, and a
    stacked column. Unsupported: waterfall (no chartEx in python-pptx) and
    anything carrying target:/callout:/marker: (native charts have no
    annotation vocabulary) — reason is None only when supported is True.
    """

    def _assert_supported(self, chart):
        ok, reason = native_charts.supported(chart)
        self.assertTrue(ok, f"expected supported, got reason={reason!r}")
        self.assertIsNone(reason)

    def _assert_unsupported(self, chart, expected_reason):
        ok, reason = native_charts.supported(chart)
        self.assertFalse(ok)
        self.assertEqual(reason, expected_reason)

    # -- supported: five plain chart types + stacked column ------------------

    def test_bar_plain_supported(self):
        self._assert_supported(_bar_chart(ctype="bar"))

    def test_column_plain_supported(self):
        self._assert_supported(_bar_chart(ctype="column"))

    def test_pie_plain_supported(self):
        self._assert_supported(_pie_chart())

    def test_line_plain_supported(self):
        self._assert_supported(_point_chart(ctype="line"))

    def test_scatter_plain_supported(self):
        self._assert_supported(_point_chart(ctype="scatter"))

    def test_column_stacked_supported(self):
        chart = _bar_chart(
            stacked=True,
            categories=["Q1", "Q2"],
            series=[{"name": "Revenue", "values": [10.0, 20.0]},
                    {"name": "Cost", "values": [5.0, 8.0]}],
        )
        self._assert_supported(chart)

    # -- unsupported: waterfall ------------------------------------------------

    def test_waterfall_unsupported(self):
        self._assert_unsupported(_waterfall_chart(), REASON_WATERFALL)
        # The pinned wording contains 'native' (satisfies the plain-language
        # requirement independent of the exact-string pin above).
        _ok, reason = native_charts.supported(_waterfall_chart())
        self.assertIn("native", reason)

    # -- unsupported: target:/callout:/marker: (shared annotation wording) ---

    def test_target_unsupported(self):
        chart = _bar_chart(target={"value": 50.0, "label": "goal"})
        self._assert_unsupported(chart, REASON_ANNOTATION)
        _ok, reason = native_charts.supported(chart)
        self.assertIn("target", reason)
        self.assertIn("drawn", reason)

    def test_callout_unsupported(self):
        chart = _bar_chart(callout="Q2 dipped")
        self._assert_unsupported(chart, REASON_ANNOTATION)
        _ok, reason = native_charts.supported(chart)
        self.assertIn("callout", reason)
        self.assertIn("drawn", reason)

    def test_line_with_markers_unsupported(self):
        chart = _point_chart(ctype="line", markers=[{"x": 1.0, "label": "peak"}])
        self._assert_unsupported(chart, REASON_ANNOTATION)
        _ok, reason = native_charts.supported(chart)
        self.assertIn("marker", reason)
        self.assertIn("drawn", reason)

    def test_scatter_with_markers_unsupported(self):
        chart = _point_chart(ctype="scatter",
                              markers=[{"x": 1.0, "label": "peak"}])
        self._assert_unsupported(chart, REASON_ANNOTATION)


class TestInsertColumn(unittest.TestCase):
    """native_charts.insert(slide, chart, colours, region) -> GraphicFrame,
    read back on a fixture-template slide (A-001: the read-back tests prove
    each REQ-002 styling claim mechanically)."""

    def test_returns_graphicframe_with_chart_on_the_slide(self):
        prs, slide = _blank_slide()
        gframe = native_charts.insert(slide, _bar_chart(), COLOURS, REGION)
        self.assertTrue(gframe.has_chart)
        on_slide = [s.shape_id for s in slide.shapes
                    if getattr(s, "has_chart", False)]
        self.assertIn(gframe.shape_id, on_slide)

    def test_chart_type_is_column_clustered(self):
        prs, slide = _blank_slide()
        gframe = native_charts.insert(slide, _bar_chart(), COLOURS, REGION)
        self.assertEqual(gframe.chart.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    def test_one_series_with_given_values(self):
        prs, slide = _blank_slide()
        gframe = native_charts.insert(slide, _bar_chart(), COLOURS, REGION)
        plot = gframe.chart.plots[0]
        self.assertEqual(len(plot.series), 1)
        self.assertEqual(list(plot.series[0].values), [10.0, 20.0, 30.0])

    def test_emphasis_category_point_is_accent_others_muted(self):
        prs, slide = _blank_slide()
        chart = _bar_chart(emphasis="Q2")  # categories[1]
        gframe = native_charts.insert(slide, chart, COLOURS, REGION)
        series = gframe.chart.plots[0].series[0]
        fills = [pt.format.fill.fore_color.rgb for pt in series.points]
        self.assertEqual(fills[1], ACCENT_RGB)
        self.assertEqual(fills[0], MUTED_RGB)
        self.assertEqual(fills[2], MUTED_RGB)

    def test_dollar_shorthand_turns_on_data_labels_with_number_format(self):
        # D-004: fmt {"prefix": "$"} (the `$` shorthand) -> Excel '"$"#,##0'.
        prs, slide = _blank_slide()
        chart = _bar_chart(fmt={"prefix": "$"})
        gframe = native_charts.insert(slide, chart, COLOURS, REGION)
        plot = gframe.chart.plots[0]
        self.assertTrue(plot.has_data_labels)
        self.assertEqual(plot.data_labels.number_format, '"$"#,##0')

    def test_single_series_has_no_legend(self):
        prs, slide = _blank_slide()
        gframe = native_charts.insert(slide, _bar_chart(), COLOURS, REGION)
        self.assertFalse(gframe.chart.has_legend)

    def test_two_series_has_legend_below(self):
        chart = _bar_chart(
            series=[{"name": "Revenue", "values": [10.0, 20.0, 30.0]},
                    {"name": "Cost", "values": [5.0, 8.0, 12.0]}],
        )
        prs, slide = _blank_slide()
        gframe = native_charts.insert(slide, chart, COLOURS, REGION)
        self.assertTrue(gframe.chart.has_legend)
        self.assertEqual(gframe.chart.legend.position, XL_LEGEND_POSITION.BOTTOM)

    def test_value_axis_hidden_no_gridlines(self):
        prs, slide = _blank_slide()
        gframe = native_charts.insert(slide, _bar_chart(), COLOURS, REGION)
        value_axis = gframe.chart.value_axis
        self.assertFalse(value_axis.has_major_gridlines)
        self.assertFalse(value_axis.visible)


class TestNumberFormatMap(unittest.TestCase):
    """D-004's full fmt -> Excel number_format map, pinned via insert() +
    read-back for every row (not a private helper name, so the contract
    survives whatever internal shape T-007 chooses). Each fmt dict below is
    exactly what render.py's CHART_FORMAT_SHORTHANDS produces for the named
    `format:` shorthand."""

    ROWS = [
        ({}, "General"),                                    # no format / 'compact'
        ({"abbreviate": False}, "General"),                 # 'plain'
        ({"prefix": "$"}, '"$"#,##0'),                       # '$' / 'currency'
        ({"suffix": "%"}, '0"%"'),                          # '%' / 'percent'
        ({"prefix": "$", "suffix": "k"}, '"$"#,##0,"k"'),    # '$k'
        ({"prefix": "$", "suffix": "M"}, '"$"#,##0,,"M"'),   # '$m'
        ({"suffix": "k"}, '#,##0,"k"'),                      # 'k'
        ({"suffix": "M"}, '#,##0,,"M"'),                     # 'm'
    ]

    def test_number_format_map(self):
        for fmt, expected in self.ROWS:
            with self.subTest(fmt=fmt):
                prs, slide = _blank_slide()
                chart = _bar_chart(fmt=fmt)
                gframe = native_charts.insert(slide, chart, COLOURS, REGION)
                plot = gframe.chart.plots[0]
                self.assertTrue(plot.has_data_labels)
                self.assertEqual(plot.data_labels.number_format, expected)


class TestInsertVariants(unittest.TestCase):
    """stacked column, pie with emphasis, and the XY (scatter/line) types."""

    def test_stacked_two_series_column_is_column_stacked(self):
        chart = _bar_chart(
            stacked=True,
            categories=["Q1", "Q2"],
            series=[{"name": "Revenue", "values": [10.0, 20.0]},
                    {"name": "Cost", "values": [5.0, 8.0]}],
        )
        prs, slide = _blank_slide()
        gframe = native_charts.insert(slide, chart, COLOURS, REGION)
        self.assertEqual(gframe.chart.chart_type, XL_CHART_TYPE.COLUMN_STACKED)

    def test_pie_with_emphasis_slice_accent_others_muted(self):
        chart = _pie_chart(emphasis="B")  # categories[1]
        prs, slide = _blank_slide()
        gframe = native_charts.insert(slide, chart, COLOURS, REGION)
        self.assertEqual(gframe.chart.chart_type, XL_CHART_TYPE.PIE)
        series = gframe.chart.plots[0].series[0]
        fills = [pt.format.fill.fore_color.rgb for pt in series.points]
        self.assertEqual(fills[1], ACCENT_RGB)
        self.assertEqual(fills[0], MUTED_RGB)
        self.assertEqual(fills[2], MUTED_RGB)

    def test_scatter_chart_dict_is_xy_scatter(self):
        chart = {
            "type": "scatter", "points": [(1, 2), (3, 4)],
            "markers": [], "callout": None, "fmt": {},
            "native": True, "target": None,
        }
        prs, slide = _blank_slide()
        gframe = native_charts.insert(slide, chart, COLOURS, REGION)
        self.assertEqual(gframe.chart.chart_type, XL_CHART_TYPE.XY_SCATTER)

    def test_line_chart_dict_is_xy_scatter_lines_no_markers(self):
        chart = {
            "type": "line", "points": [(1, 2), (3, 4)],
            "markers": [], "callout": None, "fmt": {},
            "native": True, "target": None,
        }
        prs, slide = _blank_slide()
        gframe = native_charts.insert(slide, chart, COLOURS, REGION)
        self.assertEqual(gframe.chart.chart_type,
                         XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS)


if __name__ == "__main__":
    unittest.main()
