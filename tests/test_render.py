"""Renderer tests for source/skills/build-deck/scripts/render.py.

Stdlib `unittest` only — no pytest, no third-party test deps (python-pptx is
the renderer's own runtime dependency and is imported here to inspect output).

Run from the repo root:

    python3 -m unittest tests.test_render
    python3 -m unittest discover tests

Until render.py exists these tests FAIL by design (REQ-003 Wave 0 gate): the
subprocess call returns non-zero and no .pptx is produced. They pass once
render.py turns the sample deck spec into a valid .pptx.
"""
import json
import os
import shutil
import subprocess
import sys
import tempfile
import unittest

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIXTURES = os.path.join(REPO_ROOT, "tests", "fixtures")
TEMPLATE = os.path.join(FIXTURES, "sample-template.pptx")
DECK_SPEC = os.path.join(FIXTURES, "sample-deck.md")
RENDER_PY = os.path.join(
    REPO_ROOT, "source", "skills", "build-deck", "scripts", "render.py"
)

# The parse-contract tests at the bottom of this file (TableParseTest,
# WaterfallParseTest) call render.py's parse functions in-process, so import it
# as a module. Its own directory carries pptxlib; put that on the path first.
# The subprocess-based tests above do not use this import; it is purely
# additive.
sys.path.insert(0, os.path.dirname(RENDER_PY))
import render  # noqa: E402

# tests/fixtures/sample-template.pptx has its layouts renamed to role names by
# generate-fixture-template.py. `quote` has no dedicated layout and shares the
# `section` layout (TITLE + BODY) — index 2.
LAYOUT_MAP = {
    "title": 0,
    "title-content": 1,
    "section": 2,
    "two-column": 3,
    "statement": 5,
    "quote": 2,
}
LAYOUT_NAME_FOR_ROLE = {
    "title": "title",
    "section": "section",
    "statement": "statement",
    "title-content": "title-content",
    "two-column": "two-column",
    "quote": "section",  # quote shares the section layout in this fixture
}

# What sample-deck.md declares, slide by slide: (role, primary-field text).
# The primary field of each role fills the slide's title placeholder.
EXPECTED_SLIDES = [
    ("title", "From status meeting to written update"),
    ("section", "Where the hour goes"),
    ("statement", "The weekly status meeting costs the team six hours every week."),
    ("title-content", "What the meeting does, and does not, do well"),
    ("two-column", "Two ways to spend the same hour"),
    ("quote", "I skip half of what I say in standup because it does not apply to most of the room."),
]


class RenderTest(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls._tmp = tempfile.mkdtemp(prefix="slides-render-test-")
        cls.out_path = os.path.join(cls._tmp, "out.pptx")
        brand_path = os.path.join(cls._tmp, "brand.json")
        with open(brand_path, "w", encoding="utf-8") as fh:
            json.dump(
                {
                    "template": TEMPLATE,
                    "fonts": {"heading": "Calibri", "body": "Calibri"},
                    "colours": {
                        "primary": "#1F3A5F",
                        "accent": "#E07A3F",
                        "ink": "#1A1A1A",
                        "paper": "#FFFFFF",
                    },
                    "layout_map": LAYOUT_MAP,
                },
                fh,
            )
        cls.result = subprocess.run(
            [
                sys.executable,
                RENDER_PY,
                "--spec", DECK_SPEC,
                "--brand", brand_path,
                "--out", cls.out_path,
            ],
            capture_output=True,
            text=True,
        )

    def _presentation(self):
        self.assertTrue(
            os.path.exists(self.out_path),
            f"render.py produced no .pptx.\nstdout: {self.result.stdout}\n"
            f"stderr: {self.result.stderr}",
        )
        return Presentation(self.out_path)

    def test_render_exits_zero(self):
        self.assertEqual(
            self.result.returncode,
            0,
            f"render.py exited {self.result.returncode}.\n"
            f"stdout: {self.result.stdout}\nstderr: {self.result.stderr}",
        )

    def test_output_reopens(self):
        # A valid .pptx that python-pptx can parse again.
        self._presentation()

    def test_slide_count_matches_spec(self):
        prs = self._presentation()
        self.assertEqual(
            len(prs.slides),
            len(EXPECTED_SLIDES),
            "slide count must equal the number of `## Slide` sections in the spec",
        )

    def test_each_slide_uses_the_mapped_layout(self):
        prs = self._presentation()
        for i, (role, _) in enumerate(EXPECTED_SLIDES):
            with self.subTest(slide=i + 1, role=role):
                self.assertEqual(
                    prs.slides[i].slide_layout.name,
                    LAYOUT_NAME_FOR_ROLE[role],
                    f"slide {i + 1} ({role}) must use the layout brand.json maps it to",
                )

    def test_primary_field_fills_the_title_placeholder(self):
        prs = self._presentation()
        for i, (role, expected) in enumerate(EXPECTED_SLIDES):
            with self.subTest(slide=i + 1, role=role):
                title = prs.slides[i].shapes.title
                self.assertIsNotNone(title, f"slide {i + 1} has no title placeholder")
                self.assertEqual(title.text.strip(), expected)

    def test_body_field_fills_a_content_placeholder(self):
        prs = self._presentation()
        # Slide 4 is the title-content slide; its Body has four bullets.
        body_text = " ".join(
            ph.text
            for ph in prs.slides[3].placeholders
            if ph.placeholder_format.idx != 0
        )
        for bullet in (
            "surfaces blockers fast",
            "no record anyone can search",
        ):
            self.assertIn(bullet, body_text)

    def test_visual_field_is_recorded_in_speaker_notes(self):
        prs = self._presentation()
        # Slide 5 carries a Visual: field; render.py records it in the notes.
        notes = prs.slides[4].notes_slide.notes_text_frame.text
        self.assertIn("VISUAL TO ADD", notes)
        self.assertIn("two-panel diagram", notes)

    def test_no_text_box_outside_template_placeholders(self):
        # The structural guarantee against an injected strapline: render.py
        # fills only the template's own placeholders and never adds a shape.
        prs = self._presentation()
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    with self.subTest(slide=i + 1, shape=shape.shape_id):
                        self.assertTrue(
                            shape.is_placeholder,
                            f"slide {i + 1} has text in a non-placeholder shape",
                        )

    @classmethod
    def tearDownClass(cls):
        shutil.rmtree(cls._tmp, ignore_errors=True)


class RenderErrorTest(unittest.TestCase):
    """Negative-path coverage: malformed input must exit non-zero with a
    message naming the fault, and never emit a .pptx."""

    @classmethod
    def setUpClass(cls):
        cls._tmp = tempfile.mkdtemp(prefix="slides-render-error-test-")

    @classmethod
    def tearDownClass(cls):
        # rmtree (not os.remove): a chart error may leave a sidecar charts dir.
        shutil.rmtree(cls._tmp, ignore_errors=True)

    def _run(self, name, spec_text, brand=None):
        """Write a spec (+ brand) and run render.py. Returns CompletedProcess."""
        if brand is None:
            brand = {
                "template": TEMPLATE,
                "fonts": {"heading": "Calibri", "body": "Calibri"},
                "colours": {"ink": "#1A1A1A"},
                "layout_map": dict(LAYOUT_MAP),
            }
        spec_path = os.path.join(self._tmp, f"{name}.deck.md")
        brand_path = os.path.join(self._tmp, f"{name}.brand.json")
        out_path = os.path.join(self._tmp, f"{name}.pptx")
        with open(spec_path, "w", encoding="utf-8") as fh:
            fh.write(spec_text)
        with open(brand_path, "w", encoding="utf-8") as fh:
            json.dump(brand, fh)
        result = subprocess.run(
            [sys.executable, RENDER_PY, "--spec", spec_path,
             "--brand", brand_path, "--out", out_path],
            capture_output=True, text=True,
        )
        self.assertFalse(
            os.path.exists(out_path),
            "render.py must not emit a .pptx on malformed input",
        )
        return result

    def _assert_named_error(self, result, *needles):
        self.assertEqual(
            result.returncode, 1,
            f"expected exit 1; got {result.returncode}\nstderr: {result.stderr}",
        )
        msg = (result.stderr + result.stdout).lower()
        self.assertIn("error:", msg)
        for needle in needles:
            self.assertIn(needle.lower(), msg)

    def test_slide_number_gap_is_named(self):
        spec = ("---\ndeck: d\naudience: a\n---\n\n"
                "## Slide 1\nlayout: section\nTitle: One\n\n"
                "## Slide 3\nlayout: section\nTitle: Three\n")
        self._assert_named_error(self._run("gap", spec), "slide", "3")

    def test_unknown_role_is_named(self):
        spec = ("---\ndeck: d\naudience: a\n---\n\n"
                "## Slide 1\nlayout: splash\nTitle: One\n")
        self._assert_named_error(self._run("role", spec), "slide 1", "splash")

    def test_unrecognised_field_is_named(self):
        spec = ("---\ndeck: d\naudience: a\n---\n\n"
                "## Slide 1\nlayout: title\nTitle: Hi\n"
                "Strapline: Innovate. Accelerate. Dominate.\n")
        self._assert_named_error(self._run("field", spec), "slide 1", "strapline")

    def test_field_overflow_is_named(self):
        # two-column carries 3 fields; point it at the 1-placeholder
        # statement layout so the fill overflows.
        brand = {
            "template": TEMPLATE,
            "fonts": {"heading": "Calibri", "body": "Calibri"},
            "colours": {"ink": "#1A1A1A"},
            "layout_map": {**LAYOUT_MAP, "two-column": 5},
        }
        spec = ("---\ndeck: d\naudience: a\n---\n\n"
                "## Slide 1\nlayout: two-column\nTitle: T\nLeft: L\nRight: R\n")
        self._assert_named_error(
            self._run("overflow", spec, brand), "slide 1", "placeholder")

    def test_missing_brand_key_is_named(self):
        brand = {
            "template": TEMPLATE,
            "fonts": {"heading": "Calibri", "body": "Calibri"},
            "colours": {"ink": "#1A1A1A"},
        }  # no layout_map
        spec = ("---\ndeck: d\naudience: a\n---\n\n"
                "## Slide 1\nlayout: section\nTitle: One\n")
        self._assert_named_error(
            self._run("brandkey", spec, brand), "layout_map")

    # --- chart faults (REQ-005): every malformed Chart block fails loudly ---

    def _chart_brand(self, colours=None):
        return {
            "template": TEMPLATE,
            "fonts": {"heading": "Helvetica Neue", "body": "Helvetica Neue"},
            "colours": CHART_COLOURS if colours is None else colours,
            "layout_map": dict(LAYOUT_MAP),
        }

    def _chart_spec(self, chart_block, role="title-content", title="T",
                    body="A line."):
        head = f"## Slide 1\nlayout: {role}\nTitle: {title}\n"
        if body:
            head += f"Body: {body}\n"
        return "---\ndeck: d\naudience: a\n---\n\n" + head + chart_block

    def test_chart_unknown_type_is_named(self):
        spec = self._chart_spec(
            "Chart:\n  type: donut\n  categories: A, B\n  series X: 1, 2\n")
        self._assert_named_error(
            self._run("ctype", spec, self._chart_brand()), "slide 1", "donut")

    def test_chart_length_mismatch_is_named(self):
        spec = self._chart_spec(
            "Chart:\n  type: column\n  categories: A, B, C\n  series X: 1, 2\n")
        self._assert_named_error(
            self._run("clen", spec, self._chart_brand()), "slide 1", "length")

    def test_chart_non_numeric_is_named(self):
        spec = self._chart_spec(
            "Chart:\n  type: column\n  categories: A, B\n  series X: 1, two\n")
        self._assert_named_error(
            self._run("cnum", spec, self._chart_brand()), "slide 1", "two")

    def test_chart_unknown_emphasis_is_named(self):
        spec = self._chart_spec(
            "Chart:\n  type: column\n  emphasis: Z\n"
            "  categories: A, B\n  series X: 1, 2\n")
        self._assert_named_error(
            self._run("cemph", spec, self._chart_brand()), "slide 1", "z")

    def test_chart_on_wrong_role_is_named(self):
        spec = self._chart_spec(
            "Chart:\n  type: column\n  categories: A, B\n  series X: 1, 2\n",
            role="statement", body=None)
        # statement carries Statement, not Title; give it one plus the Chart.
        spec = spec.replace("Title: T\n", "Statement: S\n")
        self._assert_named_error(
            self._run("crole", spec, self._chart_brand()),
            "slide 1", "title-content")

    def test_chart_empty_colours_is_named(self):
        spec = self._chart_spec(
            "Chart:\n  type: column\n  categories: A, B\n  series X: 1, 2\n")
        self._assert_named_error(
            self._run("ccol", spec, self._chart_brand(colours={})),
            "slide 1", "colour")

    def test_chart_inline_value_is_named(self):
        spec = self._chart_spec("Chart: column\n")
        self._assert_named_error(
            self._run("cinline", spec, self._chart_brand()),
            "slide 1", "block")

    def test_chart_empty_block_is_named(self):
        # A Chart: label with no body lines before end of slide.
        spec = self._chart_spec("Chart:\n")
        self._assert_named_error(
            self._run("cempty", spec, self._chart_brand()), "slide 1", "chart")

    def test_chart_marker_missing_point_is_named(self):
        spec = self._chart_spec(
            "Chart:\n  type: line\n  points: 0 10, 1 20\n  marker: 5 Late\n")
        self._assert_named_error(
            self._run("cmark", spec, self._chart_brand()), "slide 1", "marker")

    def test_chart_line_emphasis_is_named(self):
        # emphasis is meaningless on a line chart; fail loudly, not silently.
        spec = self._chart_spec(
            "Chart:\n  type: line\n  emphasis: X\n  points: 0 10, 1 20\n")
        self._assert_named_error(
            self._run("clineemph", spec, self._chart_brand()),
            "slide 1", "emphasis")

    def test_chart_pie_multi_series_is_named(self):
        # A pie shows one whole; more than one series is ambiguous.
        spec = self._chart_spec(
            "Chart:\n  type: pie\n  categories: A, B\n"
            "  series X: 1, 2\n  series Y: 3, 4\n")
        self._assert_named_error(
            self._run("cpie", spec, self._chart_brand()), "slide 1", "pie")

    def test_chart_scatter_emphasis_is_named(self):
        # scatter has no categories to emphasise; fail loudly.
        spec = self._chart_spec(
            "Chart:\n  type: scatter\n  emphasis: Z\n  points: 0 1, 1 2\n")
        self._assert_named_error(
            self._run("cscatteremph", spec, self._chart_brand()),
            "slide 1", "emphasis")


# tests/fixtures/sample-template.pptx maps title-content to layout index 1,
# whose content placeholder (idx 1) hosts the resized Body line; the chart is a
# separate picture placed below it. Distinct accent/muted colours let the pixel
# histogram prove emphasis colouring without sampling exact coordinates.
CHART_COLOURS = {
    "accent": "#2E8B6F",
    "muted": "#8AA3A0",
    "spend": "#D98E5A",
    "ink": "#1F3A34",
}
COLUMN_BLOCK = (
    "Chart:\n  type: column\n  emphasis: 2031\n"
    "  categories: 2026, 2027, 2028, 2029, 2030, 2031\n"
    "  series Balance: 76900, 34300, 37400, 21900, 24600, 27300\n"
)
BAR_BLOCK = (
    "Chart:\n  type: bar\n  emphasis: Conservatory\n"
    "  categories: Conservatory, Car, Kitchen\n  series Cost: 49, 18, 14\n"
)
LINE_BLOCK = (
    "Chart:\n  type: line\n"
    "  points: 0 76900, 12 34300, 26 19600, 60 27300\n  marker: 26 Car\n"
)
PIE_BLOCK = (
    "Chart:\n  type: pie\n  emphasis: Rent\n"
    "  categories: Rent, Food, Travel, Savings\n"
    "  series Spend: 1200, 400, 200, 300\n"
)
SCATTER_BLOCK = (
    "Chart:\n  type: scatter\n"
    "  points: 1 2, 2 4, 3 5, 4 8\n  marker: 4 Peak\n"
)


class ChartRenderTest(unittest.TestCase):
    """REQ-001..004, 009: charts render to PNG and place as a picture."""

    @classmethod
    def setUpClass(cls):
        cls._tmp = tempfile.mkdtemp(prefix="slides-chart-render-")

    @classmethod
    def tearDownClass(cls):
        shutil.rmtree(cls._tmp, ignore_errors=True)

    def _render(self, name, chart_block, brand=None, body="A line.", env=None,
                font_files=None):
        if brand is None:
            brand = {
                "template": TEMPLATE,
                "fonts": {"heading": "Helvetica Neue", "body": "Helvetica Neue"},
                "colours": CHART_COLOURS,
                "layout_map": dict(LAYOUT_MAP),
            }
            if font_files is not None:
                brand["font_files"] = font_files
        head = "## Slide 1\nlayout: title-content\nTitle: T\n"
        if body:
            head += f"Body: {body}\n"
        spec = "---\ndeck: d\naudience: a\n---\n\n" + head + chart_block
        spec_path = os.path.join(self._tmp, f"{name}.deck.md")
        brand_path = os.path.join(self._tmp, f"{name}.brand.json")
        out_path = os.path.join(self._tmp, f"{name}.pptx")
        charts_dir = os.path.join(self._tmp, f"{name}.charts")
        with open(spec_path, "w", encoding="utf-8") as fh:
            fh.write(spec)
        with open(brand_path, "w", encoding="utf-8") as fh:
            json.dump(brand, fh)
        result = subprocess.run(
            [sys.executable, RENDER_PY, "--spec", spec_path,
             "--brand", brand_path, "--out", out_path,
             "--charts-dir", charts_dir],
            capture_output=True, text=True, env=env,
        )
        return result, out_path, charts_dir

    def _ok(self, result, out_path):
        self.assertEqual(
            result.returncode, 0,
            f"render failed.\nstdout: {result.stdout}\nstderr: {result.stderr}")
        self.assertTrue(os.path.exists(out_path), "no .pptx produced")
        return Presentation(out_path)

    def _pictures(self, slide):
        return [s for s in slide.shapes
                if s.shape_type == MSO_SHAPE_TYPE.PICTURE]

    def test_chart_slide_has_one_picture(self):
        result, out, _ = self._render("col", COLUMN_BLOCK)
        prs = self._ok(result, out)
        self.assertEqual(len(self._pictures(prs.slides[0])), 1)

    def test_sidecar_png_written(self):
        result, out, charts_dir = self._render("png", COLUMN_BLOCK)
        self._ok(result, out)
        pngs = [f for f in os.listdir(charts_dir) if f.endswith(".png")]
        self.assertTrue(pngs, f"no PNG in {charts_dir}")

    def test_picture_sits_below_body_within_slide(self):
        result, out, _ = self._render("region", COLUMN_BLOCK)
        prs = self._ok(result, out)
        slide = prs.slides[0]
        pic = self._pictures(slide)[0]
        body = next(p for p in slide.placeholders
                    if p.placeholder_format.idx != 0)
        self.assertIsNotNone(pic.top)
        self.assertGreaterEqual(pic.top, body.top + body.height)
        self.assertGreaterEqual(pic.left, 0)
        self.assertLessEqual(pic.left + pic.width, prs.slide_width + 1)
        self.assertLessEqual(pic.top + pic.height, prs.slide_height + 1)

    def test_each_type_renders_a_picture(self):
        for name, block in (("t-bar", BAR_BLOCK), ("t-col", COLUMN_BLOCK),
                            ("t-line", LINE_BLOCK), ("t-pie", PIE_BLOCK),
                            ("t-scatter", SCATTER_BLOCK)):
            with self.subTest(chart=name):
                result, out, _ = self._render(name, block)
                prs = self._ok(result, out)
                self.assertEqual(len(self._pictures(prs.slides[0])), 1)

    def test_emphasis_colour_present(self):
        # Whole-image histogram: both the brand accent and the muted colour
        # must appear, proving emphasis colouring sourced from brand.json.
        from PIL import Image

        result, out, charts_dir = self._render("hist", COLUMN_BLOCK)
        self._ok(result, out)
        png = os.path.join(
            charts_dir,
            next(f for f in os.listdir(charts_dir) if f.endswith(".png")))
        img = Image.open(png).convert("RGB")
        histogram = img.getcolors(maxcolors=2 ** 24) or []

        def near(target, tol=10):
            tr, tg, tb = target
            return sum(n for n, (r, g, b) in histogram
                       if abs(r - tr) <= tol and abs(g - tg) <= tol
                       and abs(b - tb) <= tol)

        accent = (0x2E, 0x8B, 0x6F)
        muted = (0x8A, 0xA3, 0xA0)
        self.assertGreater(near(accent), 200, "brand accent not found in chart")
        self.assertGreater(near(muted), 200, "muted colour not found in chart")

    def test_missing_font_warns(self):
        result, out, _ = self._render("nofont", COLUMN_BLOCK)
        self._ok(result, out)
        self.assertIn("font", (result.stdout + result.stderr).lower())
        self.assertIn("fallback", (result.stdout + result.stderr).lower())

    def test_font_file_registered_no_warning(self):
        import matplotlib
        ttf = os.path.join(os.path.dirname(matplotlib.__file__),
                           "mpl-data", "fonts", "ttf", "DejaVuSans.ttf")
        self.assertTrue(os.path.isfile(ttf), "matplotlib DejaVuSans missing")
        result, out, _ = self._render(
            "withfont", COLUMN_BLOCK, font_files={"DejaVu Sans": ttf})
        self._ok(result, out)
        self.assertNotIn("fallback", (result.stdout + result.stderr).lower())

    def test_matplotlib_absent_falls_back_to_note(self):
        # REQ-009/D-011: hide matplotlib via a PYTHONPATH shim; the chart slide
        # falls back to a VISUAL TO ADD note and the deck still builds.
        shim = os.path.join(self._tmp, "shim")
        os.makedirs(shim, exist_ok=True)
        with open(os.path.join(shim, "matplotlib.py"), "w") as fh:
            fh.write('raise ImportError("matplotlib hidden for test")\n')
        env = dict(os.environ)
        env["PYTHONPATH"] = shim + os.pathsep + env.get("PYTHONPATH", "")
        result, out, _ = self._render("absent", COLUMN_BLOCK, env=env)
        prs = self._ok(result, out)
        slide = prs.slides[0]
        self.assertEqual(len(self._pictures(slide)), 0, "should not draw a chart")
        notes = slide.notes_slide.notes_text_frame.text
        self.assertIn("VISUAL TO ADD", notes)
        summary = (result.stdout + result.stderr).lower()
        self.assertIn("matplotlib", summary)
        self.assertIn("pip install matplotlib", summary)


# --- T-001 (Wave 0, red-first): parse contract for the `table` composed block
# and the `waterfall` chart type. Both features are unimplemented, so every case
# below fails right now for a named pre-implementation reason:
#   * a `table` block  -> SpecError "unknown composed block type 'table'"
#     (raised at the top of _parse_composed_block, before any table logic — so
#      the shape/CSV cases raise it instead of returning, and the malformed
#      cases raise it instead of their eventual table-specific message).
#   * a `waterfall` chart -> SpecError "unknown chart type 'waterfall'"
#     (raised after the _parse_chart_block key loop). The two `total:`-bearing
#     cases raise earlier, on the not-yet-recognised `total:` key -> SpecError
#     "unknown chart key 'total'"; the chart_to_note case falls through the
#     point-chart branch -> KeyError 'points'. Each is the genuine current fault.
# These assert the desired end-state behaviour; they go green when T-006 lands.


class TableParseTest(unittest.TestCase):
    """REQ-001/002/003: the `table` composed block parse contract.

    The parse branch must return, nested under a `"table"` key exactly as
    `matrix` nests under `"spec"`:

        {"type": "table",
         "table": {"header": [...],
                   "rows": [{"cells": [...], "emphasis": bool}, ...]}}
    """

    def _parse_table(self, body, files=None):
        """Write a one-slide composed spec (plus optional sidecar files) and
        parse it through the public `parse_spec` entry point, so `spec_dir` is
        threaded exactly as at runtime. Returns the parsed first block dict, or
        propagates SpecError for a malformed spec."""
        tmp = tempfile.mkdtemp(prefix="slides-table-parse-")
        self.addCleanup(shutil.rmtree, tmp, ignore_errors=True)
        for name, content in (files or {}).items():
            with open(os.path.join(tmp, name), "w", encoding="utf-8") as fh:
                fh.write(content)
        spec = ("---\ndeck: d\naudience: a\n---\n\n"
                "## Slide 1\nlayout: composed\n" + body)
        spec_path = os.path.join(tmp, "deck.md")
        with open(spec_path, "w", encoding="utf-8") as fh:
            fh.write(spec)
        return render.parse_spec(spec_path)[0]["blocks"][0]

    def test_header_and_rows_shape(self):
        block = self._parse_table(
            "Block: table\n"
            "Name | Q1 | Q2\n"
            "Alpha | 1 | 2\n"
            "Beta | 3 | 4\n"
        )
        self.assertEqual(block["type"], "table")
        self.assertEqual(block["table"], {
            "header": ["Name", "Q1", "Q2"],
            "rows": [
                {"cells": ["Alpha", "1", "2"], "emphasis": False},
                {"cells": ["Beta", "3", "4"], "emphasis": False},
            ],
        })

    def test_emphasis_marker_sets_row_flag(self):
        block = self._parse_table(
            "Block: table\n"
            "Name | Q1\n"
            "Alpha | 1\n"
            "! Beta | 2\n"
        )
        rows = block["table"]["rows"]
        self.assertEqual(rows[0], {"cells": ["Alpha", "1"], "emphasis": False})
        self.assertEqual(rows[1], {"cells": ["Beta", "2"], "emphasis": True})

    def test_ragged_row_raises(self):
        # Row cell count != header cell count.
        with self.assertRaisesRegex(render.SpecError, r"(?i)cell|ragged"):
            self._parse_table(
                "Block: table\n"
                "Name | Q1 | Q2\n"
                "Alpha | 1\n"
            )

    def test_one_column_header_raises(self):
        # A single-column table is a list; the message points elsewhere.
        with self.assertRaisesRegex(render.SpecError, r"(?i)column"):
            self._parse_table(
                "Block: table\n"
                "OnlyOne\n"
                "Alpha\n"
            )

    def test_data_csv_loads_header_and_rows(self):
        block = self._parse_table(
            "Block: table\n"
            "data: costs.csv\n",
            files={"costs.csv": "Item,Q1,Q2\nAlpha,1,2\nBeta,3,4\n"},
        )
        self.assertEqual(block["type"], "table")
        self.assertEqual(block["table"], {
            "header": ["Item", "Q1", "Q2"],
            "rows": [
                {"cells": ["Alpha", "1", "2"], "emphasis": False},
                {"cells": ["Beta", "3", "4"], "emphasis": False},
            ],
        })

    def test_data_plus_inline_rows_raises(self):
        # `data:` and inline rows are mutually exclusive.
        with self.assertRaisesRegex(render.SpecError, r"(?i)inline|combined"):
            self._parse_table(
                "Block: table\n"
                "data: costs.csv\n"
                "Extra | 9 | 9\n",
                files={"costs.csv": "Item,Q1,Q2\nAlpha,1,2\n"},
            )

    def test_emphasis_key_marks_matching_row(self):
        block = self._parse_table(
            "Block: table\n"
            "data: costs.csv\n"
            "emphasis: Beta\n",
            files={"costs.csv": "Item,Q1\nAlpha,1\nBeta,2\n"},
        )
        rows = block["table"]["rows"]
        self.assertEqual(rows[0], {"cells": ["Alpha", "1"], "emphasis": False})
        self.assertEqual(rows[1], {"cells": ["Beta", "2"], "emphasis": True})

    def test_emphasis_key_no_match_raises(self):
        # An emphasis label matching no first cell fails, naming the label.
        with self.assertRaisesRegex(render.SpecError, r"Zeta"):
            self._parse_table(
                "Block: table\n"
                "data: costs.csv\n"
                "emphasis: Zeta\n",
                files={"costs.csv": "Item,Q1\nAlpha,1\nBeta,2\n"},
            )


class WaterfallParseTest(unittest.TestCase):
    """REQ-006/007/010: the `waterfall` chart parse contract.

    A waterfall takes `categories` plus exactly one `series` of signed deltas,
    accepts an optional `total:` (default label "Total", `none` disables it),
    and rejects `emphasis:` (the sign colouring already encodes emphasis).
    """

    def _parse_chart(self, lines, spec_dir=None):
        return render._parse_chart_block(1, lines, spec_dir)

    def test_type_accepted_with_single_series(self):
        chart = self._parse_chart([
            "type: waterfall",
            "categories: Start, Rent, Food, Save",
            "series Cash: 40, -15, -10, 25",
        ])
        self.assertEqual(chart["type"], "waterfall")
        self.assertEqual(chart["categories"], ["Start", "Rent", "Food", "Save"])
        self.assertEqual(len(chart["series"]), 1)
        self.assertEqual(chart["series"][0]["values"], [40, -15, -10, 25])

    def test_two_series_raises(self):
        with self.assertRaisesRegex(render.SpecError, r"(?i)series"):
            self._parse_chart([
                "type: waterfall",
                "categories: A, B",
                "series X: 1, 2",
                "series Y: 3, 4",
            ])

    def test_emphasis_raises(self):
        with self.assertRaisesRegex(render.SpecError, r"(?i)emphasis"):
            self._parse_chart([
                "type: waterfall",
                "emphasis: A",
                "categories: A, B",
                "series X: 1, 2",
            ])

    def test_total_defaults_to_total_label(self):
        chart = self._parse_chart([
            "type: waterfall",
            "categories: A, B, C",
            "series X: 10, -5, 20",
        ])
        self.assertEqual(chart.get("total_label"), "Total")

    def test_total_none_disables_total(self):
        chart = self._parse_chart([
            "type: waterfall",
            "categories: A, B",
            "series X: 10, -5",
            "total: none",
        ])
        self.assertIsNone(chart["total_label"])

    def test_total_on_non_waterfall_raises(self):
        with self.assertRaisesRegex(render.SpecError, r"(?i)waterfall"):
            self._parse_chart([
                "type: column",
                "categories: A, B",
                "series X: 1, 2",
                "total: 5",
            ])

    def test_data_csv_single_series(self):
        tmp = tempfile.mkdtemp(prefix="slides-waterfall-parse-")
        self.addCleanup(shutil.rmtree, tmp, ignore_errors=True)
        with open(os.path.join(tmp, "wf.csv"), "w", encoding="utf-8") as fh:
            fh.write("Stage,Delta\nStart,40\nRent,-15\nFood,-10\n")
        chart = self._parse_chart(
            ["type: waterfall", "data: wf.csv"], spec_dir=tmp)
        self.assertEqual(chart["type"], "waterfall")
        self.assertEqual(chart["categories"], ["Start", "Rent", "Food"])
        self.assertEqual(len(chart["series"]), 1)
        self.assertEqual(chart["series"][0]["values"], [40, -15, -10])

    def test_chart_to_note_includes_signed_deltas_and_total(self):
        chart = {
            "type": "waterfall",
            "categories": ["Start", "Rent", "Save"],
            "series": [{"name": "Cash", "values": [40, -15, 25]}],
            "emphasis": None,
            "callout": None,
            "fmt": {},
            "total_label": "Total",
        }
        note = render.chart_to_note(chart)
        self.assertIn("+40", note)
        self.assertTrue(
            "-15" in note or "−15" in note,
            f"signed negative delta missing from note: {note!r}",
        )
        self.assertIn("+25", note)
        self.assertIn("50", note)  # computed total: 40 - 15 + 25


# --- T-001 (Wave 0, red-first), continued: native-charts-plan.md REQ-001,
# REQ-003, REQ-004, REQ-005, REQ-006; decisions D-001, D-005, D-009.
#
# `native:`, `stacked:`, and `target:` are not yet recognised chart keys, so
# every line carrying one of them raises SpecError "unknown chart key
# '<key>'" the moment `_parse_chart_block`'s key loop reaches it — the
# genuine current fault, regardless of what other keys/lines surround it.
# Positive-path cases below (no key mismatch to check) are left uncaught and
# so simply error out on that same SpecError right now. Cross-validation
# cases assert the *desired* end-state wording (naming bar/column, 'series',
# 'emphasis', or the allowed point-chart types) — text the current "unknown
# chart key" message does not contain, so these fail red as assertion
# mismatches, not as accidental passes. Both go green once T-006 lands.


class NativeChartParseTest(unittest.TestCase):
    """REQ-001/D-001: the `native:` chart key parse contract.

    Pinned shape: `chart["native"]` is always present as a bool (key present
    even when the spec never mentions `native:` — default False), set from
    `true`/`false` (case-insensitive); anything else is a SpecError naming
    'native'.
    """

    def _parse_chart(self, lines, spec_dir=None):
        return render._parse_chart_block(1, lines, spec_dir)

    def test_native_true_parses_true(self):
        chart = self._parse_chart([
            "type: column",
            "categories: A, B",
            "series X: 1, 2",
            "native: true",
        ])
        self.assertIs(chart["native"], True)

    def test_native_false_parses_false(self):
        chart = self._parse_chart([
            "type: column",
            "categories: A, B",
            "series X: 1, 2",
            "native: false",
        ])
        self.assertIs(chart["native"], False)

    def test_native_absent_defaults_false(self):
        # No `native:` line at all — the key must still be present (pinned
        # shape: key always present, boolean), not merely absent-equivalent.
        chart = self._parse_chart([
            "type: column",
            "categories: A, B",
            "series X: 1, 2",
        ])
        self.assertIn("native", chart)
        self.assertIs(chart["native"], False)

    def test_native_invalid_value_raises(self):
        # Strengthened past a bare 'native' substring check (today's "unknown
        # chart key 'native'" message would already satisfy that trivially,
        # coincidentally passing before implementation) by also requiring the
        # bad token 'maybe' appear, so this is genuinely red right now.
        with self.assertRaises(render.SpecError) as cm:
            self._parse_chart([
                "type: column",
                "categories: A, B",
                "series X: 1, 2",
                "native: maybe",
            ])
        msg = str(cm.exception).lower()
        self.assertIn("native", msg)
        self.assertIn("maybe", msg)


class StackedTargetParseTest(unittest.TestCase):
    """REQ-004/005, D-009: the `stacked:` and `target:` chart key parse
    contract.

    Pinned shapes:
      chart["stacked"] -> bool, legal only on multi-series bar/column.
      chart["target"]  -> {"value": float, "label": str|None}, legal only on
                           column/bar/line (D-009's optional '| label').
    """

    def _parse_chart(self, lines, spec_dir=None):
        return render._parse_chart_block(1, lines, spec_dir)

    def _assert_chart_error(self, lines, *needles, spec_dir=None):
        with self.assertRaises(render.SpecError) as cm:
            self._parse_chart(lines, spec_dir=spec_dir)
        msg = str(cm.exception).lower()
        for needle in needles:
            self.assertIn(needle.lower(), msg)

    # --- stacked ---

    def test_stacked_true_two_series_column(self):
        chart = self._parse_chart([
            "type: column",
            "categories: A, B",
            "series X: 1, 2",
            "series Y: 3, 4",
            "stacked: true",
        ])
        self.assertIs(chart["stacked"], True)

    def test_stacked_on_pie_raises_naming_bar_column(self):
        self._assert_chart_error(
            ["type: pie", "categories: A, B", "series X: 1, 2",
             "stacked: true"],
            "bar", "column",
        )

    def test_stacked_single_series_raises(self):
        self._assert_chart_error(
            ["type: column", "categories: A, B", "series X: 1, 2",
             "stacked: true"],
            "series",
        )

    def test_stacked_with_emphasis_raises(self):
        self._assert_chart_error(
            ["type: column", "emphasis: A", "categories: A, B",
             "series X: 1, 2", "series Y: 3, 4", "stacked: true"],
            "emphasis",
        )

    # --- target ---

    def test_target_value_only(self):
        chart = self._parse_chart([
            "type: column",
            "categories: A, B",
            "series X: 1, 2",
            "target: 50",
        ])
        self.assertEqual(chart["target"], {"value": 50.0, "label": None})

    def test_target_with_label(self):
        chart = self._parse_chart([
            "type: column",
            "categories: A, B",
            "series X: 1, 2",
            "target: 50 | goal",
        ])
        self.assertEqual(chart["target"], {"value": 50.0, "label": "goal"})

    def test_target_non_numeric_raises(self):
        # Strengthened past a bare 'target' substring check (same rationale
        # as test_native_invalid_value_raises) by requiring the bad token.
        with self.assertRaises(render.SpecError) as cm:
            self._parse_chart([
                "type: column", "categories: A, B", "series X: 1, 2",
                "target: abc",
            ])
        self.assertIn("abc", str(cm.exception).lower())

    def test_target_on_scatter_raises_naming_allowed_types(self):
        self._assert_chart_error(
            ["type: scatter", "points: 0 1, 1 2", "target: 50"],
            "column", "bar", "line",
        )


class ComposedChartBlockTest(unittest.TestCase):
    """REQ-006: `Block: chart` on a composed slide reuses `_parse_chart_block`'s
    grammar (including `data:` CSV, resolved against spec_dir), nested under a
    `"chart"` key exactly as `table` nests under `"table"`:

        {"type": "chart", "chart": {<same dict _parse_chart_block returns>}}

    `chart` is not yet in COMPOSED_BLOCK_TYPES, so every case below raises
    SpecError "unknown composed block type 'chart'" right now — the genuine
    current fault, raised at the top of _parse_composed_block before any
    chart-specific logic runs (so the CSV case fails the same way as the
    inline case). These assert the desired end-state; they go green once
    T-006 adds 'chart' to COMPOSED_BLOCK_TYPES and delegates to
    _parse_chart_block(items, spec_dir).
    """

    def _parse_composed(self, body, files=None):
        tmp = tempfile.mkdtemp(prefix="slides-composed-chart-parse-")
        self.addCleanup(shutil.rmtree, tmp, ignore_errors=True)
        for name, content in (files or {}).items():
            with open(os.path.join(tmp, name), "w", encoding="utf-8") as fh:
                fh.write(content)
        spec = ("---\ndeck: d\naudience: a\n---\n\n"
                "## Slide 1\nlayout: composed\n" + body)
        spec_path = os.path.join(tmp, "deck.md")
        with open(spec_path, "w", encoding="utf-8") as fh:
            fh.write(spec)
        return render.parse_spec(spec_path)[0]["blocks"][0]

    def test_chart_block_parses_inline(self):
        block = self._parse_composed(
            "Block: chart\n"
            "type: column\n"
            "categories: A, B\n"
            "series X: 1, 2\n"
        )
        self.assertEqual(block["type"], "chart")
        chart = block["chart"]
        self.assertEqual(chart["type"], "column")
        self.assertEqual(chart["categories"], ["A", "B"])
        self.assertEqual(len(chart["series"]), 1)
        self.assertEqual(chart["series"][0]["name"], "X")
        self.assertEqual(chart["series"][0]["values"], [1, 2])

    def test_chart_block_data_csv_resolves_against_spec_dir(self):
        block = self._parse_composed(
            "Block: chart\n"
            "type: column\n"
            "data: sales.csv\n",
            files={"sales.csv": "Category,X\nA,1\nB,2\n"},
        )
        self.assertEqual(block["type"], "chart")
        chart = block["chart"]
        self.assertEqual(chart["type"], "column")
        self.assertEqual(chart["categories"], ["A", "B"])
        self.assertEqual(chart["series"][0]["values"], [1, 2])


if __name__ == "__main__":
    unittest.main()
