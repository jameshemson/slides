"""native_charts.py — render a deck spec's Chart block to a native, editable
PowerPoint chart (`native: true`; native-charts-plan.md T-007, D-002).

`render.py` calls `supported(chart)` to decide whether a chart can be drawn
natively (REQ-001/003, D-005) and, when it can, `insert(slide, chart, colours,
region)`. This module owns ALL `pptx.chart` use, mirroring charts.py's
ownership of matplotlib — the two backends stay symmetric. It is imported
lazily by render.py only when a slide asks for `native: true`, so a chartless
(or image-only) deck never touches it.

Relationship to D-002: render.py never lets this module touch the slide or
the template beyond the single region it is given. Every colour comes from
the caller-supplied `colours` dict (brand.json's colour dict) — resolved with
the same fallback rule as charts._resolve_colours, reimplemented locally so
this module never imports charts.py (which imports matplotlib at module
scope; native charts must work without it — D-010). There are no colour or
font literals here. `region` is the same (left, top, width, height) EMU tuple
render.py's `_chart_region` already computes for the matplotlib picture path:
a GraphicFrame takes the same box a picture did.
"""
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# D-005's exact fallback-note wording — pinned verbatim, read by render.py's
# run summary (`slide N: <reason>; drawn as an image`).
REASON_WATERFALL = "waterfall has no native PowerPoint form"
REASON_ANNOTATION = "target:/callout:/marker: are drawn annotations"

# D-003's type map. Category charts (CategoryChartData) index on `stacked:`;
# pie is never stacked (REQ-004 restricts `stacked:` to bar/column). Point
# charts (XyChartData) carry no stacked/emphasis grammar at all.
_CATEGORY_XL_TYPE = {
    "column": (XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.COLUMN_STACKED),
    "bar": (XL_CHART_TYPE.BAR_CLUSTERED, XL_CHART_TYPE.BAR_STACKED),
    "pie": (XL_CHART_TYPE.PIE, XL_CHART_TYPE.PIE),
}
_POINT_XL_TYPE = {
    "line": XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}


class ChartError(Exception):
    """A native chart that cannot be drawn (mirrors charts.ChartError)."""


# --- supportability (D-005, REQ-003) ------------------------------------------


def supported(chart):
    """(bool, reason_or_None): can `chart` be drawn as a native PowerPoint
    chart? Pure — no pptx use — so render.py can call it before deciding
    which backend to import.

    False for `waterfall` (no chartEx in python-pptx: Out of scope, see the
    plan) and for anything carrying `target:`, `callout:`, or a non-empty
    `markers:` (native charts have no annotation vocabulary — REQ-003). True
    for plain or stacked bar/column/pie/line/scatter, reason None.
    """
    if chart.get("type") == "waterfall":
        return False, REASON_WATERFALL
    if chart.get("target") or chart.get("callout") or chart.get("markers"):
        return False, REASON_ANNOTATION
    return True, None


# --- colour resolution (mirrors charts._resolve_colours) ----------------------


def _normalise_hex(value):
    """Return '#RRGGBB' (uppercase) for a hex string, or None if unparseable.
    Accepts '#abc', 'abc', '#aabbcc', 'aabbcc' — identical to charts.py's
    helper, reimplemented so this module never imports charts.py."""
    if not isinstance(value, str):
        return None
    text = value.strip().lstrip("#")
    if len(text) == 3:
        text = "".join(ch * 2 for ch in text)
    if len(text) != 6:
        return None
    try:
        int(text, 16)
    except ValueError:
        return None
    return "#" + text.upper()


# Same generic fallbacks as charts.py (never the brand): a de-emphasised
# point or second series drawn in paper would be invisible on paper.
_MUTED_FALLBACK = "#BFBFBF"
_INK_FALLBACK = "#333333"


def _resolve_colours(colours):
    """Resolve (emphasis, muted, spend, ink) from a brand colour dict —
    identical fallback rule to charts._resolve_colours:

    emphasis = colours['accent'] else 'growth' else first value.
    muted    = colours['muted']  else a neutral grey (never paper).
    spend    = colours['spend']  else emphasis.
    ink      = colours['ink']    else a dark default.
    """
    if not isinstance(colours, dict) or not colours:
        raise ChartError("chart needs at least one brand colour, but "
                         "brand.json 'colours' is empty")
    ordered = [c for c in (_normalise_hex(v) for v in colours.values()) if c]
    if not ordered:
        raise ChartError("brand.json 'colours' has no valid hex value for "
                         "the chart")

    def pick(key):
        return _normalise_hex(colours.get(key))

    emphasis = pick("accent") or pick("growth") or ordered[0]
    muted = pick("muted") or _MUTED_FALLBACK
    spend = pick("spend") or emphasis
    ink = pick("ink") or _INK_FALLBACK
    return emphasis, muted, spend, ink


def _rgb(hex_colour):
    return RGBColor.from_string(hex_colour.lstrip("#"))


# --- number format (D-004) ----------------------------------------------------


def fmt_to_number_format(fmt):
    """A chart's `fmt` (prefix/suffix/abbreviate, from CHART_FORMAT_SHORTHANDS)
    -> an Excel number format for a native data label (D-004):

    `$` -> '"$"#,##0' · `%` -> '0"%"' (a literal suffix — spec values are
    already percentages; Excel's true `%` format would multiply by 100) ·
    `$k` -> '"$"#,##0,"k"' · `$m` -> '"$"#,##0,,"M"' · `k` -> '#,##0,"k"' ·
    `m` -> '#,##0,,"M"' · no prefix/suffix (`plain`/`compact`/absent) ->
    'General'. The image path's default large-number abbreviation has no
    Excel-format equivalent, hence the `General` fallback here.
    """
    fmt = fmt or {}
    prefix, suffix = fmt.get("prefix", ""), fmt.get("suffix", "")
    if suffix == "%":
        return '0"%"'
    if not prefix and not suffix:
        return "General"
    base = '"$"#,##0' if prefix == "$" else "#,##0"
    if suffix == "k":
        return base + ',"k"'
    if suffix == "M":
        return base + ',,"M"'
    return base


# --- public entry point -------------------------------------------------------


def insert(slide, chart, colours, region):
    """Insert `chart` as a native, editable pptx chart into `slide` within
    `region` (left, top, width, height — EMU ints). Returns the GraphicFrame.

    `chart` is the dict render.py's `_parse_chart_block` produces; callers
    must have already checked `supported(chart)`.
    """
    emphasis, muted, spend, ink = _resolve_colours(colours)
    ctype = chart["type"]
    left, top, width, height = region

    if ctype in _POINT_XL_TYPE:
        xd = XyChartData()
        series = xd.add_series("Series 1")
        for x, y in chart["points"]:
            series.add_data_point(x, y)
        gframe = slide.shapes.add_chart(
            _POINT_XL_TYPE[ctype], left, top, width, height, xd
        )
        _style_point_series(gframe.chart.plots[0].series[0], ctype, emphasis)
        _style_legend(gframe.chart, 1)
        return gframe

    cats, series_list = chart["categories"], chart["series"]
    stacked = bool(chart.get("stacked"))
    xl_type = _CATEGORY_XL_TYPE[ctype][1 if stacked else 0]
    cd = CategoryChartData()
    cd.categories = cats
    for s in series_list:
        cd.add_series(s["name"], s["values"])
    gframe = slide.shapes.add_chart(xl_type, left, top, width, height, cd)

    _fill_category(gframe.chart, chart, ctype, emphasis, muted, spend, ink)
    _style_data_labels(gframe.chart, chart.get("fmt"))
    _style_legend(gframe.chart, len(series_list))
    _hide_value_axis(gframe.chart)
    return gframe


# --- styling (REQ-002) --------------------------------------------------------


def _fill_category(pptx_chart, chart, ctype, emphasis, muted, spend, ink):
    """Per-point/per-series fills, mirroring charts._draw_bars/_draw_pie:

    - single series + `emphasis:` -> that category's point is accent, the
      rest muted (same rule for bar/column/pie — a pie slice is a point too);
    - single series, no emphasis -> every point accent, EXCEPT pie, which
      rotates the 4-colour palette (accent/muted/spend/ink) across slices —
      matching charts._draw_pie's own no-emphasis behaviour exactly;
    - 2+ series -> one palette colour per series (accent/muted/spend,
      repeating), no per-point override — a stacked/grouped chart reads by
      series, not by category.
    """
    plot = pptx_chart.plots[0]
    series = plot.series
    cats = chart["categories"]
    emph = chart.get("emphasis")

    if len(series) == 1:
        if emph is not None:
            point_colours = [emphasis if c == emph else muted for c in cats]
        elif ctype == "pie":
            palette = [emphasis, muted, spend, ink]
            point_colours = [palette[i % len(palette)] for i in range(len(cats))]
        else:
            point_colours = [emphasis] * len(cats)
        for point, colour in zip(series[0].points, point_colours):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _rgb(colour)
    else:
        palette = [emphasis, muted, spend]
        for i, s in enumerate(series):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = _rgb(palette[i % len(palette)])


def _style_point_series(series, ctype, emphasis):
    """A line/scatter chart is always a single, un-emphasised series
    (D-003/REQ-004 restrict `stacked:`/`emphasis:` to bar/column) — one accent
    colour on the line (line) or the markers (scatter)."""
    if ctype == "scatter":
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = _rgb(emphasis)
    else:
        series.format.line.color.rgb = _rgb(emphasis)


def _style_data_labels(pptx_chart, fmt):
    """Data labels on, formatted per D-004 — the direct-value-label craft
    that keeps a category chart readable without a value axis."""
    plot = pptx_chart.plots[0]
    plot.has_data_labels = True
    labels = plot.data_labels
    labels.number_format = fmt_to_number_format(fmt)
    labels.number_format_is_linked = False


def _style_legend(pptx_chart, n_series):
    """No legend for a single series (the chart already reads on its own);
    2+ series get a legend below the plot, out of its layout, so it never
    crowds the slide title above — the same placement charts.py's grouped
    bars use."""
    if n_series >= 2:
        pptx_chart.has_legend = True
        pptx_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        pptx_chart.legend.include_in_layout = False
    else:
        pptx_chart.has_legend = False


def _hide_value_axis(pptx_chart):
    """Category charts (bar/column/pie) hide the value axis and its
    gridlines entirely — the data labels already carry the numbers. Pie has
    no value axis; python-pptx signals that with a ValueError, which is the
    guard here rather than a type check."""
    try:
        value_axis = pptx_chart.value_axis
    except ValueError:
        return
    value_axis.has_major_gridlines = False
    value_axis.has_minor_gridlines = False
    value_axis.visible = False
