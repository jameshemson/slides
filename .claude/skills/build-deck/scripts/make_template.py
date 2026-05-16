#!/usr/bin/env python3
"""make_template.py — generate a themed starter .pptx for a user with no template.

Most users have a real PowerPoint template; `build-deck` fills it. A user with
none needs a starting point. This script produces one: it begins from
python-pptx's bundled default presentation (which ships a full set of named
layouts), strips any slides, renames the layouts `build-deck` uses to the six
semantic role names, and applies the user's brand fonts and colours to the
theme. A reasonable themed starter is the bar — not a design showcase.

The output is a normal .pptx the user can open, refine in PowerPoint, and then
point brand.json's `template` key at.

Usage:

    python3 make_template.py --out starter.pptx \\
        --colours '#1A1A2E,#E94560' \\
        --heading-font Georgia --body-font Verdana

Exit status: 0 on success, non-zero with a message on stderr otherwise.
"""
import argparse
import os
import sys

from pptx import Presentation

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptxlib import apply_theme, _normalise_hex  # noqa: E402

# Default-template layout index -> semantic role name. The default python-pptx
# template ships these layouts at these stable indices. `quote` has no
# dedicated layout; it shares `section` (a title + body layout), matching how
# brand.json's layout_map points two roles at one layout.
ROLE_LAYOUTS = {
    0: "title",
    1: "title-content",
    2: "section",
    3: "two-column",
    5: "statement",
}


def main(argv=None):
    parser = argparse.ArgumentParser(
        description="Generate a themed starter .pptx with role-named layouts."
    )
    parser.add_argument("--out", required=True, help="output .pptx path")
    parser.add_argument(
        "--colours",
        default="",
        help="comma-separated brand hex colours, e.g. '#1A1A2E,#E94560'",
    )
    parser.add_argument(
        "--heading-font", default="", help="heading (major) typeface name"
    )
    parser.add_argument("--body-font", default="", help="body (minor) typeface name")
    args = parser.parse_args(argv)

    colours = _parse_colours(args.colours)
    if args.colours.strip() and not [c for c in colours if _normalise_hex(c)]:
        print(
            f"error: no usable hex colours parsed from {args.colours!r}; "
            f"expected comma-separated 6-digit hex like '#1A1A2E,#E94560'",
            file=sys.stderr,
        )
        return 1

    prs = Presentation()  # bundled default template: master + named layouts

    # A fresh Presentation() has no slides; guard anyway so the starter is
    # always layouts-only, exactly what render.py expects of a template.
    sld_id_lst = prs.slides._sldIdLst
    for sld in list(sld_id_lst):
        sld_id_lst.remove(sld)

    for idx, role in ROLE_LAYOUTS.items():
        prs.slide_layouts[idx].name = role

    # apply_theme reads only the brand dict we hand it — no literals here.
    brand = {
        "fonts": {"heading": args.heading_font, "body": args.body_font},
        "colours": {f"brand{n}": hex_ for n, hex_ in enumerate(colours, start=1)},
    }
    apply_theme(prs, brand)

    out_dir = os.path.dirname(os.path.abspath(args.out))
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    try:
        prs.save(args.out)
    except Exception as exc:  # noqa: BLE001
        print(f"error: could not write {args.out}: {exc}", file=sys.stderr)
        return 1

    layout_count = len(Presentation(args.out).slide_layouts)
    print(
        f"wrote {args.out} — {layout_count} layouts, "
        f"roles renamed: {', '.join(sorted(ROLE_LAYOUTS.values()))}"
    )
    return 0


def _parse_colours(raw):
    """Split a comma-separated string into a clean list of hex colour tokens."""
    return [token.strip() for token in raw.split(",") if token.strip()]


if __name__ == "__main__":
    sys.exit(main())
