#!/usr/bin/env python3
"""deck_to_spec.py — the render.py inverse: extract a `.deck.md` from a `.pptx`.

Given a rendered deck, reconstructs a deck spec `render.parse_spec` accepts,
plus a sidecar import report naming how each slide was read. Round-trip plan:
`.build/plans/round-trip-revise-plan.md`, REQ-002/003, D-004..D-009.

Two tiers, resolved per slide (not per deck — a mostly-pack deck with one
hand-added slide degrades only on that slide):

  Tier 1 (lineage, `--brand <brand.json>` given): a slide's rendered layout is
  identified against the template's `slide_layouts` and looked up in the
  brand's (inverted) `layout_map`. Placeholder text fills the resolved role's
  fields in `render.ROLE_FIELDS` order — the exact inverse of
  `render.fill_placeholders`/`_ordered_role_fields`. A `quote`/`section`
  collision at one layout index (D-006) is broken by populated-placeholder
  count, then by precedence; the report names a precedence resolution.

  Tier 2 (foreign, no `--brand` or an unmapped layout index): a slide's role is
  guessed from its placeholder types and counts.

Structural detection runs BEFORE either tier and wins regardless of it: a
slide whose real content sits in drawn, non-placeholder shapes or tables (a
pack `composed` slide, or a foreign diagram) cannot be reconstructed into
`Block:` grammar (D-005, out of scope — see the plan) and is flattened into a
`title-content` slide instead, with a report line naming the slide. A picture
anywhere on a slide (a chart PNG, a photo) cannot be reversed to data either
(D-008); it becomes a `Visual:` note asking for re-placement, without forcing
a flatten by itself — text is what forces a flatten, not a picture alone.

The frontmatter's `audience` is always the fixed D-007 placeholder sentence
below and `register` always defaults to `presented` — a human names the real
audience before any re-render; extraction cannot recover it. `deck` is the
first `title`-role slide's `Title` text, or the pptx's filename stem.

Before ever writing a byte, the assembled spec text is round-tripped through
`render.parse_spec` (self-validation): a spec this script cannot itself parse
is never shipped — see `main`.

Usage:

    python3 deck_to_spec.py <deck.pptx> [--brand <brand.json>] [--out <spec>]
                             [--against <spec>] [--report <path>]

Exit status: 0 on a clean extraction, or on `--against` with no drift; 1 with
`error: <reason>` on unreadable input or a `.potx` given by mistake; 2 on
`--against` drift (differences printed to stdout, never stderr).
"""
import argparse
import os
import re
import sys
import tempfile

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER  # noqa: E402
from pptxlib import CONTENT_PLACEHOLDER_TYPES, load_template  # noqa: E402
import render  # noqa: E402

# D-007: the placeholder sentence extraction always writes for `audience` — a
# real audience cannot be recovered from a rendered deck, so a human is asked
# for it (the `revise` skill's job) before any re-render.
AUDIENCE_PLACEHOLDER = "(unknown — set before re-rendering)"
DEFAULT_REGISTER = "presented"

# D-006: precedence once populated-placeholder count doesn't settle a
# layout-index collision between two fixed roles (quote/section share one in
# the fixture template; nothing stops a user brand.json from creating others).
ROLE_PRECEDENCE = [
    "title-content", "two-column", "title", "statement", "section", "quote",
]

VISUAL_PREFIX = "VISUAL TO ADD:"
# D-008: a picture cannot be reversed to the data/photo it came from.
PICTURE_NOTE = (
    "a picture on this slide needs re-placing — re-place it or re-declare "
    "its Chart block"
)

STAMP_RE = re.compile(
    r"^slides-spec:\s*(?P<name>\S.*?)\s+sha256:(?P<sha>[0-9a-f]{64})\s*$"
)


class ExtractError(Exception):
    """A deck that cannot be opened, or the wrong kind of file entirely.

    Caller prints 'error: <message>' and exits 1 — render.py's own style.
    """


def main(argv=None):
    parser = argparse.ArgumentParser(
        description="Extract a deck spec (.deck.md) from a rendered .pptx — "
        "the inverse of render.py."
    )
    parser.add_argument("deck", help="path to a rendered .pptx")
    parser.add_argument(
        "--brand", default=None,
        help="path to brand.json; enables lineage-tier role mapping",
    )
    parser.add_argument(
        "--out", default=None,
        help="output .deck.md path (default: <deck-stem>.deck.md)",
    )
    parser.add_argument(
        "--against", default=None,
        help="diff the deck's extracted text against an existing spec",
    )
    parser.add_argument(
        "--report", default=None,
        help="import report path (default: <out-or-deck-stem>.import-report.md)",
    )
    args = parser.parse_args(argv)

    try:
        prs, stamp_note = _open_deck(args.deck)
        layout_map = None
        if args.brand:
            layout_map = render.load_brand(args.brand)["layout_map"]
        slides, slide_reports = _extract_slides(prs, layout_map)
        deck_title = _frontmatter_deck(args.deck, slides)
        spec_text = _format_spec(
            deck_title, AUDIENCE_PLACEHOLDER, DEFAULT_REGISTER, slides
        )
        # Self-validation (Approach, layer 2): never emit a spec render.py
        # itself would reject. A failure here is a bug in the emitter, not a
        # user-facing condition — surfaced as a named error, not a traceback.
        extracted_parsed = _self_validate(spec_text)
    except ExtractError as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 1
    except render.SpecError as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 1

    report_path = _report_path(args.deck, args.out, args.report)
    _write_report(report_path, args.deck, stamp_note, slide_reports)

    if args.against:
        return _run_against(extracted_parsed, args.against)

    out_path = args.out or (os.path.splitext(args.deck)[0] + ".deck.md")
    out_dir = os.path.dirname(os.path.abspath(out_path))
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    with open(out_path, "w", encoding="utf-8") as fh:
        fh.write(spec_text)
    print(
        f"extracted {len(slides)} slide(s) to {out_path}; "
        f"report at {report_path}"
    )
    return 0


# --- opening the deck / lineage stamp ----------------------------------------


def _open_deck(path):
    """Open `path` as a presentation, or raise ExtractError naming why not."""
    if path.lower().endswith(".potx"):
        raise ExtractError(
            "revise takes a deck (.pptx), not a template (.potx) — "
            "templates are teach-slides' job"
        )
    if not os.path.isfile(path):
        raise ExtractError(f"cannot open {path}: no such file")
    try:
        prs = load_template(path)
    except Exception as exc:  # noqa: BLE001 — any open failure, named cleanly
        raise ExtractError(f"cannot open {path}: {exc}")
    return prs, _read_stamp(prs)


def _read_stamp(prs):
    """REQ-001: read render.py's lineage stamp back out, for the report only.

    Never used to locate or re-read the original spec's *content* — that is
    the `revise` skill's job once a stamp is found; here it only tells the
    report whether this deck is pack-rendered or a foreign import.
    """
    comments = (prs.core_properties.comments or "").strip()
    match = STAMP_RE.match(comments)
    if match:
        return (
            f"lineage stamp found: spec {match.group('name')!r}, "
            f"sha256:{match.group('sha')[:12]}…"
        )
    return "no lineage stamp found — treated as a foreign-tier import"


# --- placeholder / shape text helpers ----------------------------------------


def _flatten_text(text):
    """Collapse all whitespace (including embedded newlines) to single spaces.

    Every inline field value and every block-field item is passed through
    this before it is written, so nothing extracted can ever break the
    single-line 'Field: value' grammar `render.parse_spec` expects.
    """
    return " ".join((text or "").split())


def _content_placeholders(slide):
    """Slide placeholders that hold content, in idx order. No furniture.

    Mirrors pptxlib's own (private) helper of the same shape — the inverse
    read needs the identical idx-ordering fill_placeholders wrote with.
    """
    phs = [
        ph for ph in slide.placeholders
        if ph.placeholder_format.type in CONTENT_PLACEHOLDER_TYPES
    ]
    phs.sort(key=lambda ph: ph.placeholder_format.idx)
    return phs


def _placeholder_inline(ph):
    return _flatten_text(ph.text_frame.text)


def _placeholder_block(ph):
    return [
        t for t in (_flatten_text(p.text) for p in ph.text_frame.paragraphs) if t
    ]


def _harvest_shape(shape):
    """Recursively pull text out of one shape. Returns (texts, has_picture).

    Handles the shapes composed slides and foreign decks actually carry:
    textboxes/autoshapes (has_text_frame), tables (has_table, one row per
    line), grouped shapes (walked via .shapes) — and two kinds of D-008
    "cannot reverse to data" content, neither treated as text: pictures
    (chart PNGs, photos) and native chart GraphicFrames (has_chart — a real
    pptx chart object, not text or a picture, but just as unrecoverable as
    one here; both are flagged the same way, as a picture would be).
    """
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or getattr(shape, "has_chart", False):
        return [], True
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        texts, has_picture = [], False
        for sub in shape.shapes:
            sub_texts, sub_pic = _harvest_shape(sub)
            texts.extend(sub_texts)
            has_picture = has_picture or sub_pic
        return texts, has_picture
    if getattr(shape, "has_table", False):
        rows = []
        for row in shape.table.rows:
            cells = [c for c in (_flatten_text(cell.text) for cell in row.cells) if c]
            if cells:
                rows.append(" | ".join(cells))
        return rows, False
    if getattr(shape, "has_text_frame", False):
        text = _flatten_text(shape.text_frame.text)
        return ([text] if text else []), False
    return [], False


def _harvest_drawn_content(slide):
    """Text (and a has_picture flag) from every non-placeholder shape.

    A non-empty result is the structural signal (D-005) that this slide's
    real content lives outside any placeholder — a composed pack slide or a
    foreign diagram — and must be flattened, regardless of what its layout
    index would otherwise map to.
    """
    texts, has_picture = [], False
    for shape in slide.shapes:
        if shape.is_placeholder:
            continue
        shape_texts, shape_pic = _harvest_shape(shape)
        texts.extend(shape_texts)
        has_picture = has_picture or shape_pic
    return texts, has_picture


# --- role resolution: mapped (Tier 1), heuristic (Tier 2), flattened --------


def _invert_layout_map(layout_map):
    """layout index -> [role, ...] for the SIX FIXED roles only.

    'composed' is deliberately dropped even when brand.json maps it (as the
    fixture does, sharing an index with 'statement'): a real composed slide
    is always caught first by _harvest_drawn_content (every Block: type
    draws at least one shape), so it never reaches this table — dropping it
    here means a plain fixed-role slide sharing that index resolves cleanly,
    with no D-006 collision to disambiguate.
    """
    inv = {}
    for role, idx in (layout_map or {}).items():
        if role not in render.ROLE_FIELDS:
            continue
        inv.setdefault(idx, []).append(role)
    return inv


def _layout_index(prs, slide):
    """The template layout index a slide was built from (A-003): identity of
    slide.slide_layout against prs.slide_layouts[i], by part name (stable —
    prs.slide_layouts[i] is the same object/part on every access, unlike
    placeholders, which are rewrapped)."""
    target = slide.slide_layout.part.partname
    for i, layout in enumerate(prs.slide_layouts):
        if layout.part.partname == target:
            return i
    return None


def _disambiguate(slide, candidates):
    """D-006: break a role collision at one layout index.

    2 populated content placeholders -> quote (has Attribution); 1 -> section.
    Anything else falls to a fixed precedence — the caller's report line
    names that resolution as by-precedence, not exact.
    """
    if len(candidates) == 1:
        return candidates[0], "mapped"
    populated = sum(
        1 for ph in _content_placeholders(slide)
        if _placeholder_inline(ph) or _placeholder_block(ph)
    )
    if populated == 2 and "quote" in candidates:
        return "quote", "mapped (disambiguated: 2 populated placeholders)"
    if populated == 1 and "section" in candidates:
        return "section", "mapped (disambiguated: 1 populated placeholder)"
    for role in ROLE_PRECEDENCE:
        if role in candidates:
            return role, "mapped (disambiguated by precedence, not exactly)"
    return candidates[0], "mapped (disambiguated: arbitrary)"


def _extract_mapped_fields(slide, role):
    """The inverse of render._ordered_role_fields/fill_placeholders: zip the
    role's fields against the slide's content placeholders in idx order.

    Symmetric by construction for a pack deck — render.py dropped any
    placeholder an optional, absent field left unfilled (_drop_unused), so
    what's left, in idx order, is exactly what ROLE_FIELDS[role] filled.
    """
    fields = {}
    for field_name, ph in zip(render.ROLE_FIELDS[role], _content_placeholders(slide)):
        if field_name in render.BLOCK_FIELDS:
            items = _placeholder_block(ph)
            if items:
                fields[field_name] = items
        else:
            value = _placeholder_inline(ph)
            if value:
                fields[field_name] = value
    return fields


def _infer_heuristic_role(slide, number):
    """Tier 2: guess a role from placeholder types/counts alone (no brand).

    CENTER_TITLE -> title; title + one other -> title-content; title + two
    others -> two-column; exactly one populated placeholder -> statement;
    nothing usable -> title-content with whatever text exists as Body (never
    an empty required field — always parses).
    """
    content_phs = _content_placeholders(slide)
    title_ph = next(
        (ph for ph in content_phs
         if ph.placeholder_format.type == PP_PLACEHOLDER.CENTER_TITLE), None
    )
    is_center = title_ph is not None
    if title_ph is None:
        title_ph = next(
            (ph for ph in content_phs
             if ph.placeholder_format.type == PP_PLACEHOLDER.TITLE), None
        )
    others = [ph for ph in content_phs if ph is not title_ph]

    if is_center:
        fields = {"Title": _placeholder_inline(title_ph)}
        if others:
            sub = _placeholder_inline(others[0])
            if sub:
                fields["Subtitle"] = sub
        return "title", fields

    if title_ph is not None and len(others) == 1:
        body = _placeholder_block(others[0]) or [
            "(this slide's content placeholder was empty)"
        ]
        return "title-content", {"Title": _placeholder_inline(title_ph), "Body": body}

    if title_ph is not None and len(others) == 2:
        left, right = _placeholder_block(others[0]), _placeholder_block(others[1])
        if left and right:
            return "two-column", {
                "Title": _placeholder_inline(title_ph), "Left": left, "Right": right,
            }
        # A required two-column side is empty: fall back rather than emit an
        # invalid slide (both a-column sides are required, unlike Body).
        body = (left or []) + (right or []) or [
            "(this slide's content placeholders were empty)"
        ]
        return "title-content", {"Title": _placeholder_inline(title_ph), "Body": body}

    populated = [t for t in (_placeholder_inline(ph) for ph in content_phs) if t]
    if len(populated) == 1:
        return "statement", {"Statement": populated[0]}

    title_text = populated[0] if populated else f"Slide {number}"
    body = populated[1:] or ["(no extractable text content on this slide)"]
    return "title-content", {"Title": title_text, "Body": body}


def _flatten(slide, number, drawn_texts):
    """D-005: flatten a composed/foreign-diagram slide's drawn text into
    title-content. Block: grammar is never reconstructed from geometry."""
    title = None
    title_ph = slide.shapes.title
    if title_ph is not None:
        text = _flatten_text(title_ph.text_frame.text)
        if text:
            title = text
    items = list(drawn_texts)
    if title is None and items:
        title = items.pop(0)
    if title is None:
        title = f"Slide {number}"
    if not items:
        items = ["(content on this slide could not be extracted automatically)"]
    note = "drawn shapes flattened into Body (Block: grammar not reconstructed)"
    return "title-content", {"Title": title, "Body": items}, note


def _resolve_slide(prs, slide, number, layout_roles, has_brand):
    """Resolve one slide's role and fields. Returns (role, fields, meta, how,
    note) — `how`/`note` feed the per-slide report line."""
    notes_text, visual_from_notes = _read_notes(slide)
    drawn_texts, has_picture = _harvest_drawn_content(slide)

    if drawn_texts:
        role, fields, note = _flatten(slide, number, drawn_texts)
        meta = _build_meta(notes_text, visual_from_notes, has_picture)
        return role, fields, meta, "flattened: drawn shapes/tables, not placeholders", note

    role, fields, how = None, None, None
    idx = _layout_index(prs, slide)
    if idx in layout_roles:
        role, how = _disambiguate(slide, layout_roles[idx])
        fields = _extract_mapped_fields(slide, role)
        try:
            render._validate_slide_fields(number, role, fields)
        except render.SpecError:
            role = None  # mapped extraction produced an invalid slide; fall back

    if role is None:
        role, fields = _infer_heuristic_role(slide, number)
        how = "heuristic role inference" + (
            " (layout not in brand's layout_map)" if has_brand
            else " (no --brand given)"
        )

    meta = _build_meta(notes_text, visual_from_notes, has_picture)
    note = "picture on this slide flagged for re-placement" if has_picture else None
    return role, fields, meta, how, note


# --- meta: Notes / Visual, restored from the rendered speaker notes ---------


def _read_notes(slide):
    """Split a rendered slide's speaker notes back into (Notes, Visual).

    render._apply_meta wrote "\\n\\n".join([Notes?, 'VISUAL TO ADD: ...'?,
    ...]); this is the exact inverse — any VISUAL TO ADD line is pulled back
    into Visual and never left sitting inside the restored Notes text.
    """
    if not getattr(slide, "has_notes_slide", False):
        return "", None
    raw = slide.notes_slide.notes_text_frame.text
    if not raw.strip():
        return "", None
    note_bits, visual_bits = [], []
    for part in (p.strip() for p in raw.split("\n\n") if p.strip()):
        if part.startswith(VISUAL_PREFIX):
            visual_bits.append(_flatten_text(part[len(VISUAL_PREFIX):]))
        else:
            note_bits.append(part)
    notes = _flatten_text(" ".join(note_bits))
    visual = "; ".join(v for v in visual_bits if v) or None
    return notes, visual


def _build_meta(notes_text, visual_from_notes, has_picture):
    meta = {}
    if notes_text:
        meta["Notes"] = notes_text
    bits = [visual_from_notes] if visual_from_notes else []
    if has_picture:
        bits.append(PICTURE_NOTE)
    if bits:
        meta["Visual"] = "; ".join(bits)
    return meta


# --- walking the deck / frontmatter / spec text formatting -----------------


def _extract_slides(prs, layout_map):
    has_brand = bool(layout_map)
    layout_roles = _invert_layout_map(layout_map) if has_brand else {}
    slides, reports = [], []
    for i, slide in enumerate(prs.slides, start=1):
        role, fields, meta, how, note = _resolve_slide(
            prs, slide, i, layout_roles, has_brand
        )
        slides.append({"number": i, "role": role, "fields": fields, "meta": meta})
        line = f"- Slide {i}: role `{role}` ({how})"
        if note:
            line += f" — {note}"
        reports.append(line)
    return slides, reports


def _frontmatter_deck(pptx_path, slides):
    """deck = the first `title`-role slide's Title, or the pptx filename stem."""
    for s in slides:
        if s["role"] == "title":
            title = s["fields"].get("Title")
            if title:
                return title
    return os.path.splitext(os.path.basename(pptx_path))[0]


def _format_spec(deck_title, audience, register, slides):
    lines = [
        "---",
        f"deck: {_flatten_text(deck_title)}",
        f"audience: {audience}",
        f"register: {register}",
        "---",
        "",
    ]
    for s in slides:
        lines.append(f"## Slide {s['number']}")
        lines.append(f"layout: {s['role']}")
        for field in render.ROLE_FIELDS[s["role"]]:
            value = s["fields"].get(field)
            if not value:
                continue
            if field in render.BLOCK_FIELDS:
                lines.append(f"{field}:")
                lines.extend(f"- {item}" for item in value)
            else:
                lines.append(f"{field}: {value}")
        if s["meta"].get("Notes"):
            lines.append(f"Notes: {s['meta']['Notes']}")
        if s["meta"].get("Visual"):
            lines.append(f"Visual: {s['meta']['Visual']}")
        lines.append("")
    return "\n".join(lines) + "\n"


def _self_validate(spec_text):
    """Round-trip the assembled spec text through render.parse_spec before
    anything is written. Returns the parsed slide list (also the canonical
    representation --against compares, so the diff sees exactly what would be
    written). Raises render.SpecError, naming the fault, on failure."""
    fd, tmp_path = tempfile.mkstemp(suffix=".deck.md")
    os.close(fd)
    try:
        with open(tmp_path, "w", encoding="utf-8") as fh:
            fh.write(spec_text)
        return render.parse_spec(tmp_path)
    finally:
        os.remove(tmp_path)


# --- import report ------------------------------------------------------------


def _report_path(deck_path, out_path, report_arg):
    if report_arg:
        return report_arg
    base = out_path or deck_path
    return os.path.splitext(base)[0] + ".import-report.md"


def _write_report(path, deck_path, stamp_note, slide_reports):
    lines = [
        f"# Import report — {os.path.basename(deck_path)}",
        "",
        stamp_note,
        "",
        "## Slides",
        "",
        *slide_reports,
        "",
    ]
    out_dir = os.path.dirname(os.path.abspath(path))
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    with open(path, "w", encoding="utf-8") as fh:
        fh.write("\n".join(lines))


# --- --against: the mechanical hand-edit detector ---------------------------


def _norm(value):
    """Whitespace-normalised comparison text — list fields join with a space
    first, matching the round-trip/drift equality rule pinned by the tests."""
    if isinstance(value, (list, tuple)):
        value = " ".join(str(v) for v in value)
    return " ".join(str(value).split())


def _collect_strings(obj):
    """Every string leaf in a nested dict/list — used only to give a composed
    slide's parsed Block: items a flat text form to diff against (D-009: no
    geometry, content only; D-005: never reconstructed, only compared)."""
    if isinstance(obj, str):
        text = obj.strip()
        return [text] if text else []
    if isinstance(obj, dict):
        out = []
        for v in obj.values():
            out.extend(_collect_strings(v))
        return out
    if isinstance(obj, (list, tuple)):
        out = []
        for v in obj:
            out.extend(_collect_strings(v))
        return out
    return []


def _diff_composed(number, pptx_slide, spec_slide):
    """A composed spec slide vs. its (always-flattened) pptx extraction:
    compare flattened text, not Block: grammar (D-005/D-009, best-effort —
    an item's non-text keys, e.g. an icon name, are not represented on the
    rendered slide and are not diffed)."""
    spec_text = _norm(
        _collect_strings(spec_slide.get("fields", {}))
        + _collect_strings(spec_slide.get("blocks", []))
    )
    pptx_text = _norm(_collect_strings(pptx_slide.get("fields", {})))
    diffs = []
    if pptx_text != spec_text:
        diffs.append(f"slide {number} Body: pptx={pptx_text!r} spec={spec_text!r}")
    pn = _norm(pptx_slide["meta"].get("Notes", ""))
    sn = _norm(spec_slide["meta"].get("Notes", ""))
    if pn != sn:
        diffs.append(f"slide {number} Notes: pptx={pn!r} spec={sn!r}")
    return diffs


def _diff_slides(pptx_slides, spec_slides):
    by_pptx = {s["number"]: s for s in pptx_slides}
    by_spec = {s["number"]: s for s in spec_slides}
    diffs = []
    for n in sorted(set(by_pptx) | set(by_spec)):
        p, s = by_pptx.get(n), by_spec.get(n)
        if p is None:
            diffs.append(f"slide {n}: exists in the spec but not the pptx")
            continue
        if s is None:
            diffs.append(f"slide {n}: exists in the pptx but not the spec")
            continue
        if s["role"] == "composed":
            diffs.extend(_diff_composed(n, p, s))
            continue
        if p["role"] != s["role"]:
            diffs.append(f"slide {n} Role: pptx={p['role']!r} spec={s['role']!r}")
            continue
        for field in render.ROLE_FIELDS.get(p["role"], []):
            pv = _norm(p["fields"].get(field, ""))
            sv = _norm(s["fields"].get(field, ""))
            if pv != sv:
                diffs.append(f"slide {n} {field}: pptx={pv!r} spec={sv!r}")
        pn = _norm(p["meta"].get("Notes", ""))
        sn = _norm(s["meta"].get("Notes", ""))
        if pn != sn:
            diffs.append(f"slide {n} Notes: pptx={pn!r} spec={sn!r}")
    return diffs


def _run_against(extracted_parsed, against_path):
    if not os.path.isfile(against_path):
        print(f"error: spec not found: {against_path}", file=sys.stderr)
        return 1
    try:
        spec_parsed = render.parse_spec(against_path)
    except render.SpecError as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 1

    diffs = _diff_slides(extracted_parsed, spec_parsed)
    for line in diffs:
        print(line)
    return 2 if diffs else 0


if __name__ == "__main__":
    sys.exit(main())
